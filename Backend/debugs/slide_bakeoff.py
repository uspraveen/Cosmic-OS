"""Slide-model bake-off: render identical prompts through two Fireworks models.

Same prompt, same planner-style JSON contract, same deterministic PPTX builder —
so any difference the reviewer sees comes from the model, not the harness.
"""

from __future__ import annotations

import json
import pathlib
import re
import sys
import time

import httpx

sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[1]))

from agents.slide_agent import llm_client
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ENV_FILE = pathlib.Path(__file__).resolve().parents[1] / "visual_enhancement.env"
OUT_DIR = pathlib.Path(r"C:\Users\Praveen Raj U S\Downloads\slide-bakeoff")

THEMES = [
    "The Future of Renewable Energy",
    "How Coffee Conquered the World",
    "Building Habits That Stick",
]

PLANNER_SYSTEM = "You are a presentation content planner for a slide-generation pipeline. Return strict JSON only."
PLANNER_USER = (
    'Plan a 3-slide presentation on the theme: "{theme}".\n'
    "Slide 1 is the title slide; slides 2 and 3 are content slides.\n"
    "Return exactly this JSON shape and nothing else:\n"
    "{{\n"
    '  "deck_title": string,\n'
    '  "subtitle": string,\n'
    '  "slides": [\n'
    '    {{"title": string, "bullets": [string, string, string], "speaker_note": string}},\n'
    '    {{"title": string, "bullets": [string, string, string, string], "speaker_note": string}},\n'
    '    {{"title": string, "bullets": [string, string, string, string], "speaker_note": string}}\n'
    "  ]\n"
    "}}\n"
    "Rules: catchy but professional wording; concrete specifics over vague claims; no markdown."
)

NAVY = RGBColor(0x1F, 0x38, 0x64)
BODY = RGBColor(0x2E, 0x2E, 0x2E)
ACCENT = RGBColor(0x2A, 0x9D, 0x8F)
GRAY = RGBColor(0x8A, 0x8A, 0x8A)


def load_key() -> str:
    for line in ENV_FILE.read_text(encoding="utf-8").splitlines():
        if line.startswith("VISUAL_ENHANCEMENT_FIREWORKS_API_KEY="):
            return line.split("=", 1)[1].strip()
    raise SystemExit("no fireworks key found")


def call_model(model: str, messages: list[dict], key: str) -> tuple[dict, dict]:
    """Call with the slide agent's payload shape; retry once without reasoning_effort."""
    payload = llm_client._build_payload(
        messages,
        temperature=0.2,
        model=model,
        max_tokens=4000,
        response_schema=None,
        force_json_object=True,
        reasoning_effort="low",
    )
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    started = time.perf_counter()
    with httpx.Client(timeout=300) as client:
        resp = client.post(
            "https://api.fireworks.ai/inference/v1/chat/completions",
            headers=headers,
            json=payload,
        )
        if resp.status_code == 400 and "reasoning_effort" in resp.text:
            payload.pop("reasoning_effort", None)
            resp = client.post(
                "https://api.fireworks.ai/inference/v1/chat/completions",
                headers=headers,
                json=payload,
            )
    latency = time.perf_counter() - started
    resp.raise_for_status()
    data = resp.json()
    msg = (data.get("choices") or [{}])[0].get("message", {})
    content = msg.get("content") or msg.get("reasoning_content") or ""
    usage = data.get("usage") or {}
    plan = llm_client.parse_json_response(content)
    return plan, {"latency_s": round(latency, 1), "usage": usage}


def slug(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")[:40]


def add_title_bar(slide, prs) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def build_pptx(plan: dict, out_path: pathlib.Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, prs)
    title = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.8))
    para = title.text_frame.paragraphs[0]
    para.text = str(plan.get("deck_title") or plan.get("subtitle") or "Untitled Deck")
    para.font.size = Pt(44)
    para.font.bold = True
    para.font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.8))
    sub_para = sub.text_frame.paragraphs[0]
    sub_para.text = str(plan.get("subtitle") or "")
    sub_para.font.size = Pt(20)
    sub_para.font.color.rgb = GRAY

    # Content slides
    for index, spec in enumerate(plan.get("slides", []) or [], start=1):
        if not isinstance(spec, dict):
            continue
        slide = prs.slides.add_slide(blank)
        add_title_bar(slide, prs)
        head = slide.shapes.add_textbox(Inches(0.9), Inches(0.5), Inches(11.5), Inches(1.0))
        head_para = head.text_frame.paragraphs[0]
        head_para.text = str(spec.get("title") or f"Slide {index}")
        head_para.font.size = Pt(32)
        head_para.font.bold = True
        head_para.font.color.rgb = NAVY
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.8))
        frame = body.text_frame
        frame.word_wrap = True
        bullets = spec.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        for bullet_index, bullet in enumerate(bullets):
            bullet_para = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
            bullet_para.text = f"•  {bullet}"
            bullet_para.font.size = Pt(20)
            bullet_para.font.color.rgb = BODY
            bullet_para.space_after = Pt(12)
        note = slide.shapes.add_textbox(Inches(0.9), Inches(6.9), Inches(11.5), Inches(0.4))
        note_para = note.text_frame.paragraphs[0]
        note_para.text = f"{index + 1} / {len(plan.get('slides', []) or []) + 1}"
        note_para.font.size = Pt(12)
        note_para.font.color.rgb = GRAY
        if spec.get("speaker_note"):
            slide.notes_slide.notes_text_frame.text = str(spec["speaker_note"])

    prs.save(out_path)


def main() -> int:
    model = sys.argv[1]
    key = load_key()
    short = model.rsplit("/", 1)[-1]
    model_dir = OUT_DIR / short
    model_dir.mkdir(parents=True, exist_ok=True)

    summary = {"model": model, "runs": []}
    for theme in THEMES:
        print(f"[{short}] theme: {theme}", flush=True)
        messages = [
            {"role": "system", "content": PLANNER_SYSTEM},
            {"role": "user", "content": PLANNER_USER.format(theme=theme)},
        ]
        try:
            plan, meta = call_model(model, messages, key)
        except Exception as exc:
            print(f"  FAILED: {exc}", flush=True)
            summary["runs"].append({"theme": theme, "error": str(exc)[:300]})
            continue
        plan_path = model_dir / f"plan_{slug(theme)}.json"
        plan_path.write_text(json.dumps(plan, indent=2, ensure_ascii=False), encoding="utf-8")
        pptx_path = model_dir / f"{slug(theme)}.pptx"
        build_pptx(plan, pptx_path)
        usage = meta.get("usage") or {}
        print(
            f"  ok in {meta['latency_s']}s -> {pptx_path.name} "
            f"(tokens: {usage.get('prompt_tokens')}in/{usage.get('completion_tokens')}out)",
            flush=True,
        )
        summary["runs"].append(
            {
                "theme": theme,
                "plan": plan_path.name,
                "pptx": pptx_path.name,
                **meta,
            }
        )

    (model_dir / "summary.json").write_text(
        json.dumps(summary, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    print(f"[{short}] done", flush=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
