"""HTML → native PPTX converter for the advanced slide workflow.

Design each slide in HTML/CSS (where the LLM is strongest), render it in
Chromium, then instead of screenshotting — like the html workflow does —
extract the resolved geometry of every visual element from the DOM and emit
REAL PowerPoint objects: text boxes, filled/outline shapes, and pictures.
The output stays fully editable in PowerPoint.

Extraction runs in the browser (page.evaluate), so any CSS layout — flex,
grid, absolute positioning — works: we only consume final computed boxes.
The converter deliberately consumes a restricted visual vocabulary (flat
fills, borders, rounded corners, linear gradients, plain text, images) and
drops the rest (shadows, blurs, transforms); the native workflow's design
prompts steer the model into that vocabulary.
"""

from __future__ import annotations

import colorsys
import logging
import re
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

logger = logging.getLogger(__name__)

# The HTML design canvas. Must match html_workflow.BASE_HTML_CSS (1280×720).
CANVAS_WIDTH_PX = 1280
CANVAS_HEIGHT_PX = 720

PX_TO_PT = 0.75  # CSS px → font points at 96 dpi

_FULL_BLEED_COVERAGE = 0.97  # a rect covering this much of the slide is "full-bleed"

_EXTRACT_ATOMS_JS = r"""
() => {
  const slide = document.querySelector('.slide');
  const out = { slideBg: null, atoms: [] };
  if (!slide) return out;
  const sb = slide.getBoundingClientRect();
  out.slideBg = getComputedStyle(slide).backgroundColor;
  const rel = (r) => ({
    x: Math.round((r.left - sb.left) * 100) / 100,
    y: Math.round((r.top - sb.top) * 100) / 100,
    w: Math.round(r.width * 100) / 100,
    h: Math.round(r.height * 100) / 100,
  });
  const TRANSPARENT = new Set(['transparent', 'rgba(0, 0, 0, 0)', 'none', '']);
  const colorOf = (raw) => {
    const v = (raw || '').trim().toLowerCase();
    return TRANSPARENT.has(v) ? null : v;
  };
  const visible = (cs) =>
    cs.display !== 'none' && cs.visibility !== 'hidden' &&
    parseFloat(cs.opacity || '1') > 0.01;
  const BLOCK_TAGS = new Set(['DIV','P','H1','H2','H3','H4','H5','H6','UL','OL','LI','SECTION','FIGURE','BLOCKQUOTE','TABLE','TR','HEADER','FOOTER','MAIN','ASIDE','NAV']);
  const isBlock = (el) => BLOCK_TAGS.has(el.tagName) || getComputedStyle(el).display.startsWith('block') || getComputedStyle(el).display === 'list-item';

  // Elements whose inline text was already consumed by an enclosing text atom.
  const consumed = new Set();

  const styleOfEl = (el) => {
    const cs = getComputedStyle(el);
    return {
      color: colorOf(cs.color) || '#111111',
      family: (cs.fontFamily || 'Arial').split(',')[0].replace(/["']/g, '').trim(),
      sizePx: parseFloat(cs.fontSize) || 18,
      weight: parseInt(cs.fontWeight, 10) || 400,
      style: cs.fontStyle || 'normal',
      align: cs.textAlign || 'left',
    };
  };

  // Collect styled runs under `root`, stopping at block boundaries (block
  // children are left for the walker to visit on their own, each with its
  // own box). Returns paragraphs: array of arrays of {text, style}.
  const collectRuns = (root, baseStyle) => {
    const paragraphs = [];
    let para = [];
    const pushPara = () => { if (para.length) { paragraphs.push(para); para = []; } };
    const styleOf = (el, fallback) => {
      if (!el) return { ...fallback };
      const cs = getComputedStyle(el);
      return {
        color: colorOf(cs.color) || fallback.color,
        family: (cs.fontFamily || fallback.family).split(',')[0].replace(/["']/g, '').trim(),
        sizePx: parseFloat(cs.fontSize) || fallback.sizePx,
        weight: parseInt(cs.fontWeight, 10) || 400,
        style: cs.fontStyle || 'normal',
        align: cs.textAlign || fallback.align,
      };
    };
    const visit = (node, style) => {
      if (node.nodeType === 3) {
        const text = node.textContent;
        if (text && text.trim()) para.push({ text: text.replace(/\s+/g, ' '), style });
        return;
      }
      if (node.nodeType !== 1) return;
      const tag = node.tagName;
      if (tag === 'SCRIPT' || tag === 'STYLE') return;
      if (tag === 'BR') { pushPara(); return; }
      const next = styleOf(node, style);
      if (node !== root && isBlock(node)) { pushPara(); return; }
      for (const child of node.childNodes) visit(child, next);
    };
    for (const node of root.childNodes) visit(node, baseStyle);
    pushPara();
    return paragraphs;
  };

  const markConsumed = (root) => {
    consumed.add(root);
    for (const el of root.querySelectorAll('*')) {
      if (!isBlock(el)) consumed.add(el);
    }
  };

  const emitText = (el) => {
    const cs = getComputedStyle(el);
    const box = rel(el.getBoundingClientRect());
    const runs = collectRuns(el, styleOfEl(el));
    if (!runs.length) return;
    const listItem = cs.display === 'list-item';
    if (listItem && cs.listStyleType !== 'none' && runs[0] && runs[0].length) {
      runs[0][0] = { ...runs[0][0], text: '•  ' + runs[0][0].text };
    }
    out.atoms.push({
      kind: 'text',
      box,
      runs,
      lineSpacing: cs.lineHeight && cs.lineHeight !== 'normal'
        ? parseFloat(cs.lineHeight) / (parseFloat(cs.fontSize) || 18) : null,
      transform: cs.textTransform || 'none',
    });
    markConsumed(el);
  };

  const emitRect = (el, cs, box) => {
    const bg = colorOf(cs.backgroundColor);
    const bgImage = cs.backgroundImage && cs.backgroundImage !== 'none' ? cs.backgroundImage : null;
    const borderWidth = parseFloat(cs.borderTopWidth) || 0;
    const borderStyle = cs.borderTopStyle || 'none';
    const borderColor = colorOf(cs.borderTopColor);
    const radiusRaw = cs.borderRadius || '0';
    let radius = 0;
    try {
      radius = Math.min(...radiusRaw.split('/').map(part =>
        Math.min(...part.trim().split(/\s+/).map(v => parseFloat(v) || 0))));
    } catch (e) { radius = 0; }
    const hasFace = bg || bgImage;
    const hasBorder = borderWidth > 0 && borderStyle !== 'none' && borderColor;
    if (!hasFace && !hasBorder) return;
    out.atoms.push({
      kind: 'rect',
      box,
      bg,
      bgGradient: bgImage && bgImage.startsWith('linear-gradient') ? bgImage : null,
      radius,
      border: hasBorder ? { widthPx: borderWidth, color: borderColor } : null,
    });
  };

  const walk = (el) => {
    for (const child of el.children) {
      if (consumed.has(child)) continue;
      const cs = getComputedStyle(child);
      if (!visible(cs)) continue;
      const r = child.getBoundingClientRect();
      const box = rel(r);
      if (box.w < 0.5 || box.h < 0.5) continue;
      const tag = child.tagName.toLowerCase();
      if (tag === 'img') {
        out.atoms.push({
          kind: 'image',
          box,
          src: child.currentSrc || child.src || '',
          fit: cs.objectFit || 'fill',
          radius: parseFloat(cs.borderRadius) || 0,
        });
        continue;
      }
      if (tag === 'svg') continue;
      emitRect(child, cs, box);
      const hasDirectText = Array.from(child.childNodes)
        .some(n => n.nodeType === 3 && n.textContent.trim());
      if (hasDirectText || child.children.length === 0) {
        if (child.textContent.trim()) emitText(child);
      }
      if (child.children.length > 0) walk(child);
    }
  };

  // The slide root itself: paint its background as the bottom-most rect.
  const slideCs = getComputedStyle(slide);
  const slideBg = colorOf(slideCs.backgroundColor);
  if (slideBg) {
    out.atoms.push({ kind: 'rect', box: { x: 0, y: 0, w: sb.width, h: sb.height }, bg: slideBg, bgGradient: null, radius: 0, border: null, root: true });
  }
  walk(slide);
  return out;
}
"""


# ── color helpers ──────────────────────────────────────────────────────────────

def parse_css_color(value: str | None) -> tuple[int, int, int, float] | None:
    """Parse a CSS color into (r, g, b, alpha 0-1)."""
    if not value:
        return None
    raw = value.strip().lower()
    if raw in {"transparent", "none", ""}:
        return None
    m = re.match(r"rgba?\(\s*([\d.]+)\s*,\s*([\d.]+)\s*,\s*([\d.]+)\s*(?:,\s*([\d.]+)\s*)?\)", raw)
    if m:
        r, g, b = (int(float(m.group(i))) for i in (1, 2, 3))
        alpha = float(m.group(4)) if m.group(4) is not None else 1.0
        return (max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)), max(0.0, min(1.0, alpha)))
    hex_match = re.match(r"#([0-9a-f]{3}|[0-9a-f]{6})$", raw)
    if hex_match:
        hex_part = hex_match.group(1)
        if len(hex_part) == 3:
            hex_part = "".join(ch * 2 for ch in hex_part)
        return (
            int(hex_part[0:2], 16),
            int(hex_part[2:4], 16),
            int(hex_part[4:6], 16),
            1.0,
        )
    return None


def blend_rgb(
    base: tuple[int, int, int],
    top: tuple[int, int, int],
    alpha: float,
) -> tuple[int, int, int]:
    alpha = max(0.0, min(1.0, alpha))
    return tuple(
        max(0, min(255, round(base[i] * (1 - alpha) + top[i] * alpha)))
        for i in range(3)
    )


def resolve_color(value: str | None, backdrop: tuple[int, int, int]) -> tuple[int, int, int] | None:
    """Solid RGB for a CSS color; translucent colors are flattened against backdrop."""
    parsed = parse_css_color(value)
    if parsed is None:
        return None
    r, g, b, alpha = parsed
    if alpha >= 0.999:
        return (r, g, b)
    if alpha <= 0.01:
        return None
    return blend_rgb(backdrop, (r, g, b), alpha)


_GRADIENT_STOP_RE = re.compile(
    r"(rgba?\([^)]+\)|#[0-9a-fA-F]{3,8}|(?:aqua|black|blue|fuchsia|gray|green|lime|maroon|navy|olive|purple|red|silver|teal|white|yellow))"
)


def parse_linear_gradient(value: str) -> tuple[tuple[int, int, int], tuple[int, int, int], float] | None:
    """Two-stop approximation of a CSS linear-gradient → (start, end, angle)."""
    if not value:
        return None
    stops = _GRADIENT_STOP_RE.findall(value)
    if len(stops) < 2:
        return None
    start = parse_css_color(stops[0])
    end = parse_css_color(stops[-1])
    if not start or not end:
        return None
    angle = 180.0
    header = value.split("(", 1)[1].split(")")[0] if "(" in value else ""
    deg_match = re.match(r"\s*([\d.]+)deg", header)
    if deg_match:
        angle = float(deg_match.group(1))
    elif "to right" in value:
        angle = 90.0
    elif "to left" in value:
        angle = 270.0
    elif "to top" in value:
        angle = 0.0
    return (start[:3], end[:3], angle)


# ── geometry / text helpers ────────────────────────────────────────────────────

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "start": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "end": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

_TEXT_TRANSFORM_MAP = {"uppercase": str.upper, "lowercase": str.lower, "capitalize": str.title}


def _apply_text_transform(text: str, transform: str) -> str:
    fn = _TEXT_TRANSFORM_MAP.get((transform or "").strip().lower())
    return fn(text) if fn and text else text


def _font_name(family: str) -> str:
    name = (family or "Arial").split(",")[0].replace('"', "").replace("'", "").strip()
    return name or "Arial"


def _is_bold(weight: int) -> bool:
    return weight >= 600


def _bounding_area(box: dict[str, float], canvas_w: float, canvas_h: float) -> float:
    return (box["w"] * box["h"]) / max(1.0, canvas_w * canvas_h)


# ── extraction ─────────────────────────────────────────────────────────────────

def extract_slide_atoms(page: Any, html_uri: str) -> dict[str, Any]:
    """Load one slide HTML in Chromium and return resolved visual atoms."""
    page.goto(html_uri, wait_until="load", timeout=45000)
    page.wait_for_timeout(150)
    return page.evaluate(_EXTRACT_ATOMS_JS)


# ── conversion ─────────────────────────────────────────────────────────────────

class _SlideCanvas:
    """Per-slide px → EMU mapping and drawing helpers."""

    def __init__(self, prs: Presentation, slide: Any, *, template_backed: bool) -> None:
        self.prs = prs
        self.slide = slide
        self.template_backed = template_backed
        self.canvas_w = float(CANVAS_WIDTH_PX)
        self.canvas_h = float(CANVAS_HEIGHT_PX)
        self.scale_x = int(prs.slide_width) / self.canvas_w
        self.scale_y = int(prs.slide_height) / self.canvas_h

    def emu_x(self, px: float) -> Emu:
        return Emu(int(px * self.scale_x))

    def emu_y(self, px: float) -> Emu:
        return Emu(int(px * self.scale_y))

    def add_shape(self, shape_type: Any, box: dict[str, float]) -> Any:
        return self.slide.shapes.add_shape(
            shape_type,
            self.emu_x(box["x"]),
            self.emu_y(box["y"]),
            self.emu_x(box["w"]),
            self.emu_y(box["h"]),
        )

    def add_picture(self, image_path: Path, box: dict[str, float], fit: str) -> Any:
        try:
            with Image.open(image_path) as img:
                src_w, src_h = img.size
        except Exception:
            src_w, src_h = 0, 0
        picture = self.slide.shapes.add_picture(
            str(image_path),
            self.emu_x(box["x"]),
            self.emu_y(box["y"]),
            self.emu_x(box["w"]),
            self.emu_y(box["h"]),
        )
        if fit == "cover" and src_w and src_h:
            # Fill the box by center-cropping the source, like object-fit: cover.
            box_ratio = box["w"] / max(box["h"], 0.01)
            src_ratio = src_w / max(src_h, 1)
            if src_ratio > box_ratio:  # source wider → crop left/right
                keep = box_ratio / src_ratio
                crop = max(0.0, (1.0 - keep) / 2)
                picture.crop_left = crop
                picture.crop_right = crop
            elif src_ratio < box_ratio:  # source taller → crop top/bottom
                keep = src_ratio / box_ratio
                crop = max(0.0, (1.0 - keep) / 2)
                picture.crop_top = crop
                picture.crop_bottom = crop
        return picture


def _apply_fill(shape: Any, rgb: tuple[int, int, int]) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _apply_gradient(shape: Any, start: tuple[int, int, int], end: tuple[int, int, int], angle: float) -> bool:
    try:
        shape.fill.gradient()
        stops = shape.fill.gradient_stops
        stops[0].color.rgb = RGBColor(*start)
        stops[0].position = 0.0
        stops[1].color.rgb = RGBColor(*end)
        stops[1].position = 1.0
        try:
            shape.fill.gradient_angle = angle
        except Exception:
            pass
        shape.line.fill.background()
        shape.shadow.inherit = False
        return True
    except Exception:
        logger.debug("html2pptx: gradient fill failed; falling back to solid", exc_info=True)
        return False


def _apply_border(shape: Any, border: dict[str, Any] | None) -> None:
    if not border:
        shape.line.fill.background()
    else:
        width_pt = max(0.75, border.get("widthPx", 1.0) * PX_TO_PT)
        rgb = resolve_color(border.get("color"), (17, 17, 17))
        if rgb is None:
            shape.line.fill.background()
            return
        shape.line.color.rgb = RGBColor(*rgb)
        shape.line.width = Pt(width_pt)
    shape.shadow.inherit = False


def _rounded_shape_kind(radius_px: float, box: dict[str, float]) -> Any:
    if radius_px >= 2:
        return MSO_SHAPE.ROUNDED_RECTANGLE
    return MSO_SHAPE.RECTANGLE


def _set_rounding(shape: Any, radius_px: float, box: dict[str, float]) -> None:
    if radius_px < 2:
        return
    try:
        min_dim_px = min(box["w"], box["h"]) or 1.0
        # roundRect adjustment = corner radius as a fraction of min(width, height)
        shape.adjustments[0] = max(0.0, min(0.5, radius_px / min_dim_px))
    except Exception:
        logger.debug("html2pptx: rounding adjustment failed", exc_info=True)


def _emit_text_atom(canvas: _SlideCanvas, atom: dict[str, Any], backdrop: tuple[int, int, int]) -> None:
    box = atom["box"]
    textbox = canvas.slide.shapes.add_textbox(
        canvas.emu_x(box["x"]),
        canvas.emu_y(box["y"]),
        canvas.emu_x(box["w"]),
        canvas.emu_y(box["h"]),
    )
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    paragraphs = atom.get("runs") or []
    line_spacing = atom.get("lineSpacing")
    first = True
    for para_runs in paragraphs:
        if not para_runs:
            continue
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        style = para_runs[0].get("style") or {}
        paragraph.alignment = _ALIGN_MAP.get(str(style.get("align") or "left"), PP_ALIGN.LEFT)
        if line_spacing and 0.5 < line_spacing < 4.0:
            paragraph.line_spacing = line_spacing
        for run_spec in para_runs:
            text = _apply_text_transform(str(run_spec.get("text") or ""), atom.get("transform") or style.get("transform") or "")
            if not text:
                continue
            run = paragraph.add_run()
            run.text = text
            run_style = run_spec.get("style") or style
            rgb = resolve_color(run_style.get("color"), backdrop)
            if rgb is None:
                rgb = backdrop
            font = run.font
            font.name = _font_name(run_style.get("family") or "Arial")
            font.size = Pt(max(1.0, float(run_style.get("sizePx") or 18) * PX_TO_PT))
            font.bold = _is_bold(int(run_style.get("weight") or 400))
            font.italic = (run_style.get("style") or "normal") == "italic"
            font.color.rgb = RGBColor(*rgb)


def _backdrop_for(atoms: list[dict[str, Any]], fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    """Best-guess page background used to flatten translucent colors."""
    for atom in atoms:
        if atom.get("kind") == "rect" and atom.get("root"):
            rgb = parse_css_color(atom.get("bg"))
            if rgb:
                return (rgb[0], rgb[1], rgb[2])
    return fallback


def convert_slide_atoms(
    canvas: _SlideCanvas,
    atoms: dict[str, Any],
    *,
    template_backed: bool,
    asset_root: Path | None,
    backdrop_fallback: tuple[int, int, int] = (17, 17, 17),
) -> None:
    """Draw one slide's atoms as native PowerPoint objects."""
    all_atoms = atoms.get("atoms") or []
    backdrop = _backdrop_for(all_atoms, backdrop_fallback)

    for atom in all_atoms:
        kind = atom.get("kind")
        box = atom.get("box") or {}
        if kind == "rect":
            coverage = _bounding_area(box, canvas.canvas_w, canvas.canvas_h)
            # In template mode the template's own slide design provides the
            # background; skip our full-bleed fills so it stays visible.
            if template_backed and atom.get("root"):
                continue
            if template_backed and coverage >= _FULL_BLEED_COVERAGE and not atom.get("bgGradient"):
                continue
            shape = canvas.add_shape(
                _rounded_shape_kind(float(atom.get("radius") or 0.0), box),
                box,
            )
            _set_rounding(shape, float(atom.get("radius") or 0.0), box)
            rgb = resolve_color(atom.get("bg"), backdrop)
            gradient = parse_linear_gradient(atom.get("bgGradient"))
            if gradient is not None:
                start, end, angle = gradient
                if not _apply_gradient(shape, start, end, angle) and rgb is not None:
                    _apply_fill(shape, rgb)
            elif rgb is not None:
                _apply_fill(shape, rgb)
            else:
                shape.fill.background()
                _apply_border(shape, atom.get("border"))
                continue
            _apply_border(shape, atom.get("border"))
        elif kind == "image":
            src = str(atom.get("src") or "")
            image_path = _local_image_path(src, asset_root)
            if image_path is None:
                logger.warning("html2pptx: dropping unresolvable image %s", src[:120])
                continue
            canvas.add_picture(image_path, box, str(atom.get("fit") or "fill"))
        elif kind == "text":
            _emit_text_atom(canvas, atom, backdrop)


def _local_image_path(src: str, asset_root: Path | None) -> Path | None:
    if not src:
        return None
    if src.startswith("file://"):
        path = Path(unquote(urlparse(src).path))
        if path.exists():
            return path
    if src.startswith("data:image"):
        return None
    candidate = Path(src)
    if candidate.is_absolute() and candidate.exists():
        return candidate
    if asset_root is not None:
        relative = unquote(urlparse(src).path).lstrip("/")
        under_root = (asset_root / relative).resolve()
        if under_root.exists():
            return under_root
        fallback = (asset_root / Path(src).name).resolve()
        if fallback.exists():
            return fallback
    if candidate.exists():
        return candidate
    return None


def convert_html_deck_to_native_pptx(
    manifest: list[dict[str, Any]],
    output_path: Path,
    *,
    deck_title: str = "",
    template_path: Path | None = None,
    template_layout_numbers: dict[int, int] | None = None,
    backdrop_rgb: tuple[int, int, int] = (17, 17, 17),
) -> Path:
    """Convert per-slide HTML designs into one native, editable PPTX.

    Args:
        manifest: html_workflow-style entries with html_path, slide_number,
            speaker_notes, and the slide dir (for resolving relative assets).
        template_path: when given, the deck is built ON this template —
            its masters, theme, and layout backgrounds are preserved and
            each new slide is placed on a template layout.
        template_layout_numbers: deck slide number → template slide number
            (from the layout selector) used to pick the layout each new
            slide starts from.
        backdrop_rgb: fallback backdrop used to flatten translucent CSS
            colors (template-backed slides have no background rect to
            sample from).
    """
    from html_workflow import _launch_browser

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if template_path is not None:
        prs, layout_by_number = _open_template_with_layouts(Path(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Emu(12192000)  # 13.333in
        prs.slide_height = Emu(6858000)  # 7.5in
        layout_by_number = {}

    template_backed = template_path is not None

    pw, browser = _launch_browser()
    converted = 0
    try:
        page = browser.new_page(viewport={"width": CANVAS_WIDTH_PX, "height": max(CANVAS_HEIGHT_PX, 1000)})
        try:
            for item in manifest:
                slide_number = int(item.get("slide_number") or (converted + 1))
                html_path = Path(item["html_path"])
                layout = _layout_for_slide(layout_by_number, template_layout_numbers, slide_number)
                if layout is not None:
                    slide = prs.slides.add_slide(layout)
                    _strip_empty_placeholders(slide)
                else:
                    slide = prs.slides.add_slide(_default_layout(prs))
                canvas = _SlideCanvas(prs, slide, template_backed=template_backed)
                atoms = extract_slide_atoms(page, html_path.as_uri())
                convert_slide_atoms(
                    canvas,
                    atoms,
                    template_backed=template_backed,
                    asset_root=html_path.parent,
                    backdrop_fallback=backdrop_rgb,
                )
                notes = str(item.get("speaker_notes") or "")
                if notes:
                    try:
                        slide.notes_slide.notes_text_frame.text = notes
                    except Exception:
                        logger.debug("html2pptx: failed to set notes on slide %d", slide_number)
                converted += 1
        finally:
            page.close()
    finally:
        browser.close()
        pw.stop()

    if deck_title:
        try:
            prs.core_properties.title = deck_title
        except Exception:
            logger.debug("html2pptx: failed to set deck title metadata")

    logger.info(
        "html2pptx: converted %d slides → %s%s",
        converted,
        output_path.name,
        " (template-backed)" if template_backed else "",
    )
    prs.save(str(output_path))
    return output_path


def _default_layout(prs: Presentation) -> Any:
    """A blank-ish layout: index 6 when available, otherwise the last one."""
    layouts = prs.slide_layouts
    if len(layouts) > 6:
        return layouts[6]
    return layouts[len(layouts) - 1]


def _open_template_with_layouts(template_path: Path) -> tuple[Presentation, dict[int, Any]]:
    """Open a template PPTX, keep its masters/theme, drop its slides, and
    remember which layout each template slide used so new slides can inherit
    the template's designed chrome."""
    prs = Presentation(str(template_path))
    layout_by_number: dict[int, Any] = {}
    for index, slide in enumerate(list(prs.slides)):
        layout_by_number[index + 1] = slide.slide_layout
    slide_id_list = prs.slides._sldIdLst
    for _ in range(len(slide_id_list)):
        r_id = slide_id_list[0].rId
        prs.part.drop_rel(r_id)
        del slide_id_list[0]
    return prs, layout_by_number


def _layout_for_slide(
    layout_by_number: dict[int, Any],
    template_layout_numbers: dict[int, int] | None,
    slide_number: int,
) -> Any | None:
    if not layout_by_number:
        return None
    template_slide_number = int((template_layout_numbers or {}).get(slide_number) or 1)
    return layout_by_number.get(template_slide_number) or next(iter(layout_by_number.values()))


def _strip_empty_placeholders(slide: Any) -> None:
    """Remove layout-inherited placeholders — none of them carry content in
    this workflow, and empty ones would show prompt text in PowerPoint."""
    for shape in list(slide.placeholders):
        has_text = shape.has_text_frame and bool(shape.text_frame.text.strip())
        if not has_text:
            shape._element.getparent().remove(shape._element)
