"""PPTX assembly helpers for the HTML workflow."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5


def build_pptx_from_images(
    slide_pngs: Iterable[Path],
    output_path: Path,
    *,
    deck_title: str = "",
    speaker_notes: list[str] | None = None,
) -> Path:
    """Create a widescreen PPTX with one rendered image per slide."""
    slide_pngs = [Path(p) for p in slide_pngs]
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    while len(prs.slides) > 0:
        slide_id_list = prs.slides._sldIdLst
        r_id = slide_id_list[0].rId
        prs.part.drop_rel(r_id)
        del slide_id_list[0]

    for idx, png in enumerate(slide_pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if speaker_notes and idx < len(speaker_notes) and speaker_notes[idx]:
            try:
                slide.notes_slide.notes_text_frame.text = speaker_notes[idx]
            except Exception:
                logger.debug("pptx_writer: failed to set notes on slide %d", idx + 1)

    if deck_title:
        try:
            prs.core_properties.title = deck_title
        except Exception:
            logger.debug("pptx_writer: failed to set deck title metadata")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
