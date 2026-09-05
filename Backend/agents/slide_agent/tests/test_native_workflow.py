"""The advanced (native) slide workflow: HTML-designed, editable PPTX.

The converter is the risky half — HTML geometry must land as real PowerPoint
objects — so these tests exercise it directly: pure color/gradient parsing,
atom→shape conversion without a browser, template theme extraction, a full
extract→convert pass through a real Chromium render, and the agent wiring
('advanced' as the default workflow, template-backed dispatch, edit routing).
"""

from __future__ import annotations

import asyncio
import json
import shutil
import sys
from pathlib import Path
from typing import Any
from uuid import uuid4

import pytest
from pptx import Presentation
from pptx.util import Emu

SLIDE_AGENT_DIR = Path(__file__).resolve().parent.parent
if str(SLIDE_AGENT_DIR) not in sys.path:
    sys.path.insert(0, str(SLIDE_AGENT_DIR))

import html2pptx  # noqa: E402
from agents.slide_agent.agent import SlideAgent  # noqa: E402
from agents.slide_agent.config import SlideAgentConfig  # noqa: E402
from shared.contracts import TaskEnvelope  # noqa: E402

TEST_RUNTIME_ROOT = Path(__file__).resolve().parent / "_runtime"


def _runtime_dir() -> Path:
    path = TEST_RUNTIME_ROOT / uuid4().hex
    path.mkdir(parents=True, exist_ok=False)
    return path


def _task(
    *,
    task_id: str,
    input_payload: dict[str, Any],
    intent: str = "slide.create",
) -> TaskEnvelope:
    return TaskEnvelope(
        task_id=task_id,
        task_list_id="task_list_slide_tests",
        session_id="sess_slide_tests",
        sender="cosmic/orchestrator:1.0.0",
        recipient="cosmic/slide-agent:1.0.0",
        intent=intent,
        input=input_payload,
        input_artifacts=[],
        idempotency_key=f"idem_{task_id}",
        signature="test-signature",
    )


def _agent(runtime_dir: Path) -> SlideAgent:
    root = runtime_dir / "runs" / "artifacts"
    return SlideAgent(
        None,
        config=SlideAgentConfig(
            artifacts_root=root,
            templates_dir=runtime_dir / "templates",
            catalogs_dir=runtime_dir / "catalogs",
            assets_cache_dir=runtime_dir / "assets" / "cache",
            validate_outputs=False,
        ),
    )


# ── color + gradient parsing ──────────────────────────────────────────────────


def test_parse_css_color_handles_hex_rgb_and_alpha() -> None:
    assert html2pptx.parse_css_color("#FF8800") == (255, 136, 0, 1.0)
    assert html2pptx.parse_css_color("#08c") == (0, 136, 204, 1.0)
    assert html2pptx.parse_css_color("rgb(10, 20, 30)") == (10, 20, 30, 1.0)
    assert html2pptx.parse_css_color("rgba(10, 20, 30, 0.5)") == (10, 20, 30, 0.5)
    assert html2pptx.parse_css_color("transparent") is None
    assert html2pptx.parse_css_color("") is None
    assert html2pptx.parse_css_color(None) is None


def test_resolve_color_flattens_translucent_colors_against_backdrop() -> None:
    white = (255, 255, 255)
    assert html2pptx.resolve_color("rgba(0, 0, 0, 0.5)", white) == (128, 128, 128)
    # Fully transparent resolves to None (nothing to draw).
    assert html2pptx.resolve_color("rgba(0, 0, 0, 0)", white) is None
    assert html2pptx.resolve_color("#FF0000", white) == (255, 0, 0)


def test_parse_linear_gradient_extracts_two_stops_and_angle() -> None:
    parsed = html2pptx.parse_linear_gradient(
        "linear-gradient(90deg, #ff0000 0%, #0000ff 100%)"
    )
    assert parsed is not None
    start, end, angle = parsed
    assert start == (255, 0, 0)
    assert end == (0, 0, 255)
    assert angle == 90.0

    directional = html2pptx.parse_linear_gradient(
        "linear-gradient(to right, rgb(255,0,0), rgb(0,0,255))"
    )
    assert directional is not None and directional[2] == 90.0

    # A single stop cannot become a gradient.
    assert html2pptx.parse_linear_gradient("linear-gradient(#fff)") is None


def test_apply_text_transform_maps_case() -> None:
    assert html2pptx._apply_text_transform("hello world", "uppercase") == "HELLO WORLD"
    assert html2pptx._apply_text_transform("Hello", "none") == "Hello"


def test_font_name_takes_first_family_and_strips_quotes() -> None:
    assert html2pptx._font_name("'Bahnschrift', Arial, sans-serif") == "Bahnschrift"
    assert html2pptx._font_name("") == "Arial"


# ── atom → native conversion (no browser) ─────────────────────────────────────


def _fresh_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs


def _text_of(slide) -> list[str]:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return texts


def test_convert_rect_and_text_atoms_to_native_shapes() -> None:
    prs = _fresh_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas = html2pptx._SlideCanvas(prs, slide, template_backed=False)

    atoms = {
        "slideBg": "#0d0d0f",
        "atoms": [
            {"kind": "rect", "root": True, "box": {"x": 0, "y": 0, "w": 1280, "h": 720},
             "bg": "#0d0d0f", "bgGradient": None, "radius": 0, "border": None},
            {"kind": "rect", "box": {"x": 80, "y": 120, "w": 520, "h": 300},
             "bg": "rgba(255, 255, 255, 0.16)", "bgGradient": None, "radius": 26, "border": None},
            {"kind": "text", "box": {"x": 80, "y": 60, "w": 900, "h": 120},
             "runs": [[{"text": "Quarterly Review", "style": {
                 "color": "#f5f1e8", "family": "Bahnschrift", "sizePx": 88,
                 "weight": 800, "style": "normal", "align": "left"}}]],
             "lineSpacing": 0.92, "transform": "none"},
        ],
    }

    html2pptx.convert_slide_atoms(canvas, atoms, template_backed=False, asset_root=None)

    texts = _text_of(slide)
    assert any("Quarterly Review" in text for text in texts)
    # Root backdrop + one panel must exist as autoshapes (shape_type 1).
    assert len([s for s in slide.shapes if s.shape_type == 1]) == 2

    run = slide.shapes[-1].text_frame.paragraphs[0].runs[0]
    assert run.text == "Quarterly Review"
    assert run.font.bold is True
    assert run.font.size.pt == pytest.approx(66.0)  # 88px * 0.75


def test_template_backed_conversion_skips_full_bleed_backgrounds() -> None:
    prs = _fresh_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas = html2pptx._SlideCanvas(prs, slide, template_backed=True)

    atoms = {
        "slideBg": "#0d0d0f",
        "atoms": [
            {"kind": "rect", "root": True, "box": {"x": 0, "y": 0, "w": 1280, "h": 720},
             "bg": "#0d0d0f", "bgGradient": None, "radius": 0, "border": None},
            {"kind": "text", "box": {"x": 80, "y": 60, "w": 900, "h": 120},
             "runs": [[{"text": "On template", "style": {
                 "color": "#111111", "family": "Georgia", "sizePx": 40,
                 "weight": 700, "style": "normal", "align": "left"}}]],
             "lineSpacing": None, "transform": "none"},
        ],
    }
    html2pptx.convert_slide_atoms(canvas, atoms, template_backed=True, asset_root=None)

    # The full-bleed backdrop was dropped so the template design shows through.
    assert not [s for s in slide.shapes if s.shape_type == 1]
    assert any("On template" in text for text in _text_of(slide))


def test_image_atom_places_picture_with_cover_crop(tmp_path: Path) -> None:
    from PIL import Image

    image_path = tmp_path / "photo.png"
    Image.new("RGB", (400, 200), (200, 30, 30)).save(image_path)

    prs = _fresh_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas = html2pptx._SlideCanvas(prs, slide, template_backed=False)
    atoms = {
        "atoms": [
            {"kind": "image", "box": {"x": 100, "y": 100, "w": 600, "h": 300},
             "src": str(image_path), "fit": "cover", "radius": 0},
        ],
    }
    html2pptx.convert_slide_atoms(canvas, atoms, template_backed=False, asset_root=tmp_path)

    pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    assert len(pictures) == 1
    picture = pictures[0]
    # A 2:1 source in a 2:1 box needs no crop.
    assert picture.crop_left == 0 and picture.crop_top == 0

    wide_box = {
        "atoms": [
            {"kind": "image", "box": {"x": 100, "y": 100, "w": 800, "h": 200},
             "src": str(image_path), "fit": "cover", "radius": 0},
        ],
    }
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    canvas2 = html2pptx._SlideCanvas(prs, slide2, template_backed=False)
    html2pptx.convert_slide_atoms(canvas2, wide_box, template_backed=False, asset_root=tmp_path)
    picture2 = [s for s in slide2.shapes if s.shape_type == 13][0]
    # 4:1 box from a 2:1 source crops top/bottom.
    assert picture2.crop_top > 0 and picture2.crop_bottom > 0


def test_unresolvable_image_is_dropped_not_fatal() -> None:
    prs = _fresh_prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    canvas = html2pptx._SlideCanvas(prs, slide, template_backed=False)
    atoms = {
        "atoms": [
            {"kind": "image", "box": {"x": 10, "y": 10, "w": 100, "h": 100},
             "src": "https://example.com/missing.png", "fit": "cover", "radius": 0},
        ],
    }
    html2pptx.convert_slide_atoms(canvas, atoms, template_backed=False, asset_root=None)
    assert not [s for s in slide.shapes if s.shape_type == 13]


# ── template theme extraction ─────────────────────────────────────────────────


def test_extract_template_theme_from_a_real_pptx(tmp_path: Path) -> None:
    from native_workflow import extract_template_theme

    template_path = tmp_path / "brand_template.pptx"
    prs = Presentation()
    prs.save(str(template_path))  # default python-pptx package carries a full theme1.xml

    theme = extract_template_theme(template_path)
    assert theme is not None
    assert set(theme["css_variables"]) == {
        "bg", "fg", "muted", "accent", "accent_2", "panel", "panel_strong", "line",
    }
    assert theme["css_variables"]["bg"].startswith("#")
    assert theme["css_variables"]["accent"].startswith("#")
    for stack in theme["font_stacks"].values():
        assert stack.split(",")[0].strip()
    assert theme["deck_guidelines"]


# ── full extract → convert pass through a real browser render ─────────────────


def test_browser_extract_and_convert_end_to_end(tmp_path: Path) -> None:
    pytest.importorskip("playwright.sync_api")
    slide_html = tmp_path / "slide.html"
    slide_html.write_text(
        """<!doctype html>
<html><head><style>
html, body { margin: 0; width: 1280px; height: 720px; overflow: hidden; }
.slide { position: relative; width: 1280px; height: 720px; background: #101418; color: #f5f1e8;
         font-family: Georgia, serif; }
.card { position: absolute; left: 80px; top: 200px; width: 480px; height: 240px;
        background: rgba(255,255,255,0.12); border: 1px solid rgba(255,255,255,0.3);
        border-radius: 24px; }
.hero { position: absolute; left: 80px; top: 60px; font-size: 72px; font-weight: 800;
        font-family: Bahnschrift, Arial, sans-serif; }
.copy { position: absolute; left: 104px; top: 224px; width: 430px; font-size: 20px; }
.shot { position: absolute; left: 640px; top: 200px; width: 560px; height: 320px;
        object-fit: cover; }
</style></head>
<body><section class="slide">
  <h1 class="hero">Native Deck</h1>
  <div class="card"><p class="copy">Real editable text inside a real card.</p></div>
  <img class="shot" src="photo.png" />
</section></body></html>
""",
        encoding="utf-8",
    )
    from PIL import Image

    Image.new("RGB", (320, 180), (30, 90, 160)).save(tmp_path / "photo.png")

    manifest = [{
        "slide_number": 1,
        "title": "Native Deck",
        "speaker_notes": "Speaker notes survive conversion.",
        "html_path": str(slide_html),
    }]
    output_path = tmp_path / "deck.pptx"
    html2pptx.convert_html_deck_to_native_pptx(
        manifest,
        output_path,
        deck_title="Native Deck",
    )

    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    joined = " ".join(_text_of(slide))
    assert "Native Deck" in joined
    assert "Real editable text inside a real card." in joined
    # Backdrop + card + picture + two text boxes exist as native shapes.
    assert len([s for s in slide.shapes if s.shape_type == 13]) == 1
    filled = [s for s in slide.shapes if s.shape_type == 1]  # AUTO_SHAPE
    assert len(filled) >= 2
    # Speaker notes were carried over.
    assert "Speaker notes survive conversion." in slide.notes_slide.notes_text_frame.text


def test_template_backed_conversion_preserves_template_dimensions(tmp_path: Path) -> None:
    pytest.importorskip("playwright.sync_api")
    template_path = tmp_path / "brand.pptx"
    template = Presentation()
    template.slide_width = Emu(9144000)   # 10in (4:3)
    template.slide_height = Emu(6858000)  # 7.5in
    template.slides.add_slide(template.slide_layouts[6])
    template.save(str(template_path))

    canvas_w, canvas_h = 1280, round(1280 * 6858000 / 9144000)
    slide_html = tmp_path / "slide.html"
    slide_html.write_text(
        f"""<!doctype html>
<html><head><style>
html, body {{ margin: 0; width: {canvas_w}px; height: {canvas_h}px; overflow: hidden; }}
.slide {{ position: relative; width: {canvas_w}px; height: {canvas_h}px; background: transparent;
         font-family: Georgia, serif; }}
.copy {{ position: absolute; left: 100px; top: 100px; width: 800px; font-size: 40px; }}
</style></head>
<body><section class="slide">
  <p class="copy">Template-backed slide</p>
</section></body></html>
""",
        encoding="utf-8",
    )
    output_path = tmp_path / "out.pptx"
    html2pptx.convert_html_deck_to_native_pptx(
        [{"slide_number": 1, "title": "T", "speaker_notes": "", "html_path": str(slide_html)}],
        output_path,
        template_path=template_path,
    )
    prs = Presentation(str(output_path))
    assert int(prs.slide_width) == 9144000
    assert len(prs.slides) == 1
    assert any("Template-backed slide" in text for text in _text_of(prs.slides[0]))
    # Transparent slide background must not have painted a full-bleed backdrop.
    assert not [s for s in prs.slides[0].shapes if s.shape_type == 1]


# ── agent wiring ───────────────────────────────────────────────────────────────


def test_normalize_workflow_accepts_advanced() -> None:
    assert SlideAgent._normalize_workflow("advanced") == "advanced"
    assert SlideAgent._normalize_workflow("ADVANCED") == "advanced"
    assert SlideAgent._normalize_workflow("natively") == ""


def test_auto_without_template_resolves_to_advanced(monkeypatch: Any) -> None:
    runtime_dir = _runtime_dir()
    try:
        agent = _agent(runtime_dir)
        captured: dict[str, Any] = {}

        async def fake_emit(task_id: str, event_type: str, payload: dict[str, Any]) -> str:
            return f"evt_{len(captured)}"

        async def fake_source_context(*_args: Any, **_kwargs: Any) -> dict[str, Any]:
            return {"source_artifacts": [], "documents": [], "visual_assets": []}

        def fake_run_native(
            description: str,
            template_path: Path | None,
            output_dir: Path,
            max_slides: int | None,
            validate: bool,
            force_catalog: bool,
            content_plan: dict[str, Any] | None = None,
        ) -> dict[str, Any]:
            captured["template_path"] = template_path
            captured["validate"] = validate
            output_dir.mkdir(parents=True, exist_ok=True)
            pptx_path = output_dir / "deck.pptx"
            plan_path = output_dir / "plan.json"
            theme_path = output_dir / "theme.json"
            report_path = output_dir / "build_report.json"
            pptx_path.write_bytes(b"native deck")
            plan_path.write_text(json.dumps({"deck_title": "Plan", "slides": [1]}), encoding="utf-8")
            theme_path.write_text(json.dumps({"theme_name": "planned"}), encoding="utf-8")
            report_path.write_text(json.dumps({"workflow": "advanced"}), encoding="utf-8")
            return {
                "workflow": "advanced",
                "pptx_path": str(pptx_path),
                "plan_path": str(plan_path),
                "theme_path": str(theme_path),
                "contact_sheet": "",
                "pdf_path": "",
                "validation_results": [],
            }

        monkeypatch.setattr(agent, "emit_event", fake_emit)
        monkeypatch.setattr(agent, "_prepare_source_materials_for_generation", fake_source_context)
        monkeypatch.setattr(agent, "_run_native", fake_run_native)

        result = asyncio.run(
            agent.handle_slide_create(
                _task(
                    task_id="tsk_native_auto",
                    input_payload={"description": "A native editable deck about COSMIC.", "workflow": "auto"},
                )
            )
        )

        assert result.status == "completed"
        assert result.output["workflow"] == "advanced"
        assert result.output["editable"] is True
        assert result.output["template_path"] is None
        assert captured["template_path"] is None
    finally:
        shutil.rmtree(runtime_dir, ignore_errors=True)


def test_advanced_with_template_backs_the_deck_on_it(monkeypatch: Any) -> None:
    runtime_dir = _runtime_dir()
    try:
        agent = _agent(runtime_dir)
        template_path = runtime_dir / "templates" / "brand.pptx"
        template_path.parent.mkdir(parents=True, exist_ok=True)
        Presentation().save(str(template_path))

        captured: dict[str, Any] = {}

        async def fake_emit(task_id: str, event_type: str, payload: dict[str, Any]) -> str:
            return "evt"

        async def fake_source_context(*_args: Any, **_kwargs: Any) -> dict[str, Any]:
            return {"source_artifacts": [], "documents": [], "visual_assets": []}

        def fake_run_native(
            description: str,
            template: Path | None,
            output_dir: Path,
            max_slides: int | None,
            validate: bool,
            force_catalog: bool,
            content_plan: dict[str, Any] | None = None,
        ) -> dict[str, Any]:
            captured["template"] = template
            output_dir.mkdir(parents=True, exist_ok=True)
            pptx_path = output_dir / "deck.pptx"
            plan_path = output_dir / "plan.json"
            theme_path = output_dir / "theme.json"
            pptx_path.write_bytes(b"native template deck")
            plan_path.write_text("{}", encoding="utf-8")
            theme_path.write_text("{}", encoding="utf-8")
            return {
                "workflow": "advanced",
                "pptx_path": str(pptx_path),
                "plan_path": str(plan_path),
                "theme_path": str(theme_path),
                "contact_sheet": "",
                "pdf_path": "",
                "validation_results": [],
            }

        monkeypatch.setattr(agent, "emit_event", fake_emit)
        monkeypatch.setattr(agent, "_prepare_source_materials_for_generation", fake_source_context)
        monkeypatch.setattr(agent, "_run_native", fake_run_native)

        result = asyncio.run(
            agent.handle_slide_create(
                _task(
                    task_id="tsk_native_template",
                    input_payload={
                        "description": "A deck on my template.",
                        "workflow": "advanced",
                        "template_path": str(template_path),
                    },
                )
            )
        )

        assert result.status == "completed"
        assert result.output["workflow"] == "advanced"
        assert result.output["editable"] is True
        assert captured["template"] == template_path
    finally:
        shutil.rmtree(runtime_dir, ignore_errors=True)


def test_edit_of_advanced_deck_reroutes_to_native_regeneration(monkeypatch: Any) -> None:
    runtime_dir = _runtime_dir()
    try:
        agent = _agent(runtime_dir)
        source_dir = runtime_dir / "native_source"
        source_dir.mkdir()
        source_deck = source_dir / "deck.pptx"
        source_deck.write_bytes(b"native deck")
        (source_dir / "build_report.json").write_text(
            json.dumps({"workflow": "advanced", "description": "A native deck.", "template_backed": False}),
            encoding="utf-8",
        )

        captured: dict[str, Any] = {}

        async def fake_emit(task_id: str, event_type: str, payload: dict[str, Any]) -> str:
            return "evt"

        async def fake_source_context(*_args: Any, **_kwargs: Any) -> dict[str, Any]:
            return {"source_artifacts": [], "documents": [], "visual_assets": []}

        def fail_run_template(*_args: Any, **_kwargs: Any) -> dict[str, Any]:
            raise AssertionError("advanced-deck edits must not fall through to template rebuilding")

        def fake_run_native(
            description: str,
            template_path: Path | None,
            output_dir: Path,
            max_slides: int | None,
            validate: bool,
            force_catalog: bool,
            content_plan: dict[str, Any] | None = None,
        ) -> dict[str, Any]:
            captured["description"] = description
            captured["template_path"] = template_path
            output_dir.mkdir(parents=True, exist_ok=True)
            pptx_path = output_dir / "deck.pptx"
            plan_path = output_dir / "plan.json"
            theme_path = output_dir / "theme.json"
            pptx_path.write_bytes(b"regenerated native deck")
            plan_path.write_text("{}", encoding="utf-8")
            theme_path.write_text("{}", encoding="utf-8")
            return {
                "workflow": "advanced",
                "pptx_path": str(pptx_path),
                "plan_path": str(plan_path),
                "theme_path": str(theme_path),
                "contact_sheet": "",
                "pdf_path": "",
                "validation_results": [],
            }

        monkeypatch.setattr(agent, "emit_event", fake_emit)
        monkeypatch.setattr(agent, "_prepare_source_materials_for_generation", fake_source_context)
        monkeypatch.setattr(agent, "_count_pptx_slides", lambda _path: 3)
        monkeypatch.setattr(agent, "_run_native", fake_run_native)
        monkeypatch.setattr(agent, "_run_template", fail_run_template)

        result = asyncio.run(
            agent.handle_slide_edit(
                _task(
                    task_id="tsk_edit_native",
                    intent="slide.edit",
                    input_payload={
                        "edit_request": "Tighten the copy on slide 2.",
                        "source_pptx_path": str(source_deck),
                        "validate": False,
                    },
                )
            )
        )

        assert result.status == "completed"
        assert result.output["workflow"] == "advanced"
        assert result.output["editable"] is True
        assert result.output["edit_mode"] == "advanced_regeneration"
        assert "Tighten the copy on slide 2." in captured["description"]
    finally:
        shutil.rmtree(runtime_dir, ignore_errors=True)
