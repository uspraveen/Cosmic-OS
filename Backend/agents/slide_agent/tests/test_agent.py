"""Tests for the slide agent."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

BACKEND_ROOT = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(BACKEND_ROOT))


class TestConfig:
    def test_default_config(self):
        from agents.slide_agent.config import SlideAgentConfig

        cfg = SlideAgentConfig()
        assert cfg.mimo_base_url == "https://openrouter.ai/api/v1"
        assert cfg.mimo_model == "qwen/qwen3.6-plus:free"
        assert cfg.mimo_temperature == 1.0
        assert cfg.mimo_reasoning_enabled is True
        assert cfg.mimo_reasoning_max_tokens == 256
        assert cfg.slide_use_langgraph is True
        assert cfg.default_template == "corporate-dark"
        assert cfg.export_pdf is True

    def test_config_from_env(self):
        from agents.slide_agent.config import SlideAgentConfig

        cfg = SlideAgentConfig.from_env()
        assert cfg.redis_url


class TestAgentImport:
    """Verify the agent is importable and constructable."""

    def test_import_slide_agent(self):
        from agents.slide_agent.agent import SlideAgent

        assert SlideAgent is not None

    def test_import_slide_graph(self):
        from agents.slide_agent.slide_graph import run_slide_langgraph

        assert run_slide_langgraph is not None

    def test_construct_with_mock_redis(self):
        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from unittest.mock import MagicMock

        cfg = SlideAgentConfig()
        mock_redis = MagicMock()
        agent = SlideAgent(redis_client=mock_redis, config=cfg)
        assert agent.agent_id == "cosmic/slide-agent:1.0.0"
        assert agent._cfg.default_template == "corporate-dark"

    def test_correct_super_init_kwarg(self):
        """Verify AgentRuntime.__init__ receives redis_client= not redis=."""
        import inspect
        from shared.agent_runtime import AgentRuntime

        sig = inspect.signature(AgentRuntime.__init__)
        assert "redis_client" in sig.parameters
        assert "redis" not in sig.parameters


class TestAssetManager:
    def test_resize_image(self):
        from agents.slide_agent.asset_manager import resize_image
        from PIL import Image
        import io

        # Create a test image
        img = Image.new("RGB", (1000, 500), color="red")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        img_bytes = buf.getvalue()

        resized = resize_image(img_bytes, target_width_px=200, target_height_px=100)
        result = Image.open(io.BytesIO(resized))
        assert result.width <= 200
        assert result.height <= 100

    def test_crop_to_aspect(self):
        from agents.slide_agent.asset_manager import crop_to_aspect
        from PIL import Image
        import io

        img = Image.new("RGB", (1000, 500), color="blue")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        img_bytes = buf.getvalue()

        cropped = crop_to_aspect(img_bytes, aspect_ratio=(4, 3))
        result = Image.open(io.BytesIO(cropped))
        ratio = result.width / result.height
        assert abs(ratio - 4 / 3) < 0.1

    def test_convert_rgba_to_rgb(self):
        from agents.slide_agent.asset_manager import convert_image_format
        from PIL import Image
        import io

        img = Image.new("RGBA", (100, 100), color=(255, 0, 0, 128))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        img_bytes = buf.getvalue()

        converted = convert_image_format(img_bytes, output_format="JPEG")
        result = Image.open(io.BytesIO(converted))
        assert result.mode == "RGB"

    def test_get_image_dimensions(self):
        from agents.slide_agent.asset_manager import get_image_dimensions
        from PIL import Image
        import io

        img = Image.new("RGB", (800, 600), color="green")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        img_bytes = buf.getvalue()

        w, h = get_image_dimensions(img_bytes)
        assert w == 800
        assert h == 600


class TestSlideBuilder:
    def test_hex_to_rgb(self):
        from agents.slide_agent.slide_builder import _hex_to_rgb

        color = _hex_to_rgb("#1a2b3c")
        assert color[0] == 0x1A
        assert color[1] == 0x2B
        assert color[2] == 0x3C

    def test_hex_to_rgb_short(self):
        from agents.slide_agent.slide_builder import _hex_to_rgb

        color = _hex_to_rgb("#fff")
        assert color[0] == 0xFF
        assert color[1] == 0xFF
        assert color[2] == 0xFF

    def test_build_deck_creates_pptx(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        plan = {
            "deck": {
                "title": "Test Deck",
                "template": "blank",
                "theme": {
                    "font_family": "Calibri",
                    "font_size_title": 28,
                    "font_size_body": 16,
                },
            },
            "slides": [
                {
                    "slide_number": 1,
                    "layout": "title_slide",
                    "title": "Hello World",
                    "subtitle": "Test presentation",
                },
                {
                    "slide_number": 2,
                    "layout": "content",
                    "title": "Content Slide",
                    "content": {
                        "type": "bullets",
                        "items": ["Point 1", "Point 2", "Point 3"],
                    },
                },
            ],
        }
        output_path = tmp_path / "test.pptx"
        result = builder.build_deck(plan, output_path)
        assert result.exists()
        assert result.stat().st_size > 0

    def test_build_deck_applies_template_background_defaults(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        plan = {
            "deck": {
                "title": "Pitch Deck",
                "template": "pitch-deck",
                "theme": {
                    "text_color": "#ffffff",
                    "font_family": "Calibri",
                },
            },
            "slides": [
                {
                    "slide_number": 1,
                    "layout": "title_slide",
                    "title": "COSMIC",
                    "subtitle": "Your Personal AI OS",
                }
            ],
        }
        output_path = tmp_path / "pitch_test.pptx"
        result = builder.build_deck(plan, output_path)
        prs = Presentation(str(result))
        slide = prs.slides[0]
        assert str(slide.background.fill.fore_color.rgb) == "0F0F23"

    def test_extract_structure(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Title"

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        prs2 = builder.load_existing(pptx_path)
        structure = builder.extract_structure(prs2)
        assert structure["slide_count"] == 1
        assert structure["slides"][0]["title"] == "Test Title"


class TestSchemas:
    def test_all_schemas_exist(self):
        schemas_dir = Path(__file__).resolve().parent.parent / "schemas" / "intents"
        expected = [
            "slide.create.input.json",
            "slide.create.output.json",
            "slide.edit.input.json",
            "slide.edit.output.json",
            "slide.recall_session.input.json",
            "slide.recall_session.output.json",
        ]
        for name in expected:
            path = schemas_dir / name
            assert path.exists(), f"Missing schema: {name}"
            data = json.loads(path.read_text())
            assert "type" in data


class TestValidationHeuristics:
    def test_detect_blank_or_low_contrast_slide(self):
        from io import BytesIO

        from PIL import Image

        from agents.slide_agent.internal_llm import _detect_blank_or_low_contrast_slide

        image = Image.new("RGB", (960, 540), "white")
        buf = BytesIO()
        image.save(buf, format="PNG")

        result = _detect_blank_or_low_contrast_slide(buf.getvalue())
        assert result is not None
        assert result["pass"] is False
        assert any("blank" in issue.lower() for issue in result["issues"])


class TestAgentCard:
    def test_card_parseable(self):
        import yaml

        card_path = Path(__file__).resolve().parent.parent / "agent_card.yaml"
        data = yaml.safe_load(card_path.read_text())
        assert data["agent_id"] == "cosmic/slide-agent:1.0.0"
        assert len(data["intents"]) == 3
        intent_names = [i["name"] for i in data["intents"]]
        assert "slide.create" in intent_names
        assert "slide.edit" in intent_names
        assert "slide.recall_session" in intent_names

    def test_no_auth_requirements(self):
        import yaml

        card_path = Path(__file__).resolve().parent.parent / "agent_card.yaml"
        data = yaml.safe_load(card_path.read_text())
        assert "auth_requirements" not in data

    def test_artifact_types(self):
        import yaml

        card_path = Path(__file__).resolve().parent.parent / "agent_card.yaml"
        data = yaml.safe_load(card_path.read_text())
        assert "slide_pptx" in data["artifact_types"]
        assert "slide_pdf" in data["artifact_types"]

    def test_schemas_referenced(self):
        import yaml

        card_path = Path(__file__).resolve().parent.parent / "agent_card.yaml"
        schemas_dir = Path(__file__).resolve().parent.parent / "schemas" / "intents"
        data = yaml.safe_load(card_path.read_text())
        for intent in data["intents"]:
            input_path = schemas_dir / intent["input_schema"].split("/")[-1]
            output_path = schemas_dir / intent["output_schema"].split("/")[-1]
            assert input_path.exists(), f"Missing input schema for {intent['name']}"
            assert output_path.exists(), f"Missing output schema for {intent['name']}"

    def test_usage_hints_present_for_primary_intents(self):
        import yaml

        card_path = Path(__file__).resolve().parent.parent / "agent_card.yaml"
        data = yaml.safe_load(card_path.read_text())
        intents = {item["name"]: item for item in data["intents"]}
        assert intents["slide.create"].get("usage_hints")
        assert intents["slide.edit"].get("usage_hints")
        assert intents["slide.recall_session"].get("usage_hints")


class TestInternalLLM:
    def test_openrouter_helpers(self):
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.internal_llm import (
            _effective_temperature,
            _extra_body,
            _usage_provider_name,
        )

        cfg = SlideAgentConfig(mimo_temperature=2.0)
        assert _usage_provider_name(cfg) == "openrouter"
        assert _effective_temperature(cfg) == 1.0
        assert _extra_body(cfg) == {
            "reasoning": {"enabled": True, "max_tokens": 256}
        }

    def test_plan_deck_signature(self):
        import inspect
        from agents.slide_agent.internal_llm import plan_deck

        sig = inspect.signature(plan_deck)
        assert "description" in sig.parameters
        assert "template" in sig.parameters

    def test_validate_slide_signature(self):
        import inspect
        from agents.slide_agent.internal_llm import validate_slide

        sig = inspect.signature(validate_slide)
        assert "png_bytes" in sig.parameters
        assert "slide_number" in sig.parameters

    def test_repair_deck_exists(self):
        from agents.slide_agent.internal_llm import repair_deck
        import inspect

        sig = inspect.signature(repair_deck)
        assert "slide_plans" in sig.parameters
        assert "validation_results" in sig.parameters

    def test_plan_deck_accepts_learnings(self):
        import inspect
        from agents.slide_agent.internal_llm import plan_deck

        sig = inspect.signature(plan_deck)
        assert "learnings_context" in sig.parameters


class TestChartSandbox:
    def test_code_sandbox_module_exists(self):
        from agents.slide_agent import code_sandbox as cs

        assert hasattr(cs, "run_sandbox")
        assert hasattr(cs, "generate_chart")
        assert hasattr(cs, "provision_venv")

    def test_generate_chart(self, tmp_path):
        from agents.slide_agent import code_sandbox as cs

        result = cs.generate_chart(
            chart_code="plt.bar(['A','B','C'], [10, 20, 15])",
            output_dir=tmp_path / "chart",
        )
        assert result["success"] is True
        assert len(result.get("chart_bytes", b"")) > 0

    def test_run_sandbox_captures_files(self, tmp_path):
        from agents.slide_agent import code_sandbox as cs

        output_dir = tmp_path / "sandbox"
        result = cs.run_sandbox(
            code="from pathlib import Path\nPath('output.txt').write_text('hello')\n",
            output_dir=output_dir,
        )
        assert result["success"] is True
        assert len(result["output_files"]) > 0
        assert any("output.txt" in f["filename"] for f in result["output_files"])

    def test_run_sandbox_seaborn(self, tmp_path):
        from agents.slide_agent import code_sandbox as cs

        result = cs.run_sandbox(
            code="import seaborn as sns\nimport matplotlib.pyplot as plt\nsns.barplot(x=['A','B'], y=[1,2])\nplt.savefig('seaborn_test.png')\n",
            output_dir=tmp_path / "seaborn",
            packages=["seaborn"],
        )
        assert result["success"] is True
        assert any("seaborn_test.png" in f["filename"] for f in result["output_files"])
        assert result["output_path"].endswith(".png")
        assert len(result.get("image_bytes", b"")) > 0

    def test_generate_pie_chart(self, tmp_path):
        from agents.slide_agent import code_sandbox as cs

        result = cs.generate_chart(
            chart_code="plt.pie([30,50,20], labels=['A','B','C'])",
            output_dir=tmp_path / "pie",
        )
        assert result["success"] is True

        assert result["success"] is True


class TestSlideTransitions:
    def test_apply_transition_exists(self):
        from agents.slide_agent.slide_builder import SlideBuilder

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        assert hasattr(builder, "apply_transition")

    def test_apply_transition_does_not_crash(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        # Should not raise
        builder.apply_transition(slide, transition_type="fade", speed="med")
        prs.save(str(tmp_path / "transition_test.pptx"))
        assert (tmp_path / "transition_test.pptx").exists()


class TestUserTemplates:
    def test_load_template_supports_user_prefix(self):
        from agents.slide_agent.slide_builder import SlideBuilder

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        prs = builder.load_template("user:nonexistent")
        assert prs is not None

    def test_embed_chart_image_exists(self):
        from agents.slide_agent.slide_builder import SlideBuilder

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        assert hasattr(builder, "embed_chart_image")


class TestShapesAndConnectors:
    def test_add_shape_exists(self):
        from agents.slide_agent.slide_builder import SlideBuilder

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        assert hasattr(builder, "add_shape")
        assert hasattr(builder, "add_connector")
        assert hasattr(builder, "add_flow_diagram")

    def test_add_shape_creates_rectangle(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        builder = SlideBuilder(Path(__file__).parent.parent / "templates")

        shape = builder.add_shape(
            slide,
            shape_type="rounded_rectangle",
            x_inches=1,
            y_inches=1,
            width_inches=3,
            height_inches=1.5,
            text="Hello Box",
            fill_hex="#4472C4",
            text_color_hex="#ffffff",
            font_size=14,
            font_bold=True,
        )
        assert shape is not None
        prs.save(str(tmp_path / "shapes.pptx"))
        assert (tmp_path / "shapes.pptx").exists()

    def test_add_flow_diagram(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        builder = SlideBuilder(Path(__file__).parent.parent / "templates")

        shapes = builder.add_flow_diagram(
            slide,
            boxes=[
                {"text": "Input"},
                {"text": "Process"},
                {"text": "Output"},
            ],
            direction="horizontal",
            start_x_inches=1,
            start_y_inches=2,
            box_width_inches=2.5,
            box_height_inches=1,
            gap_inches=0.5,
            fill_hex="#4472C4",
            text_color_hex="#ffffff",
        )
        assert len(shapes) == 3
        prs.save(str(tmp_path / "flow.pptx"))
        assert (tmp_path / "flow.pptx").exists()

    def test_shape_types(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder = SlideBuilder(Path(__file__).parent.parent / "templates")

        for shape_type in [
            "rectangle",
            "rounded_rectangle",
            "ellipse",
            "diamond",
            "hexagon",
            "chevron",
        ]:
            shape = builder.add_shape(
                slide,
                shape_type=shape_type,
                x_inches=1,
                y_inches=1,
                width_inches=1.5,
                height_inches=0.8,
            )
            assert shape is not None


class TestLearningsIntegration:
    def test_learnings_file_exists(self):
        learnings_path = (
            Path(__file__).resolve().parent.parent / "store" / "learnings.md"
        )
        assert learnings_path.exists()

    def test_analyze_request_reads_learnings(self):
        """The graph's analyze_request node should read learnings.md."""
        import inspect
        from agents.slide_agent import slide_graph

        source = inspect.getsource(slide_graph)
        assert "learnings_path" in source
        assert "learnings_context" in source
        assert "User preferences from past interactions" in source


class TestLayoutEngine:
    """Test the layout engine — bounds, overlap detection, auto-spacing."""

    def test_bounding_box_overlap(self):
        from agents.slide_agent.layout_engine import BoundingBox

        a = BoundingBox(x=0, y=0, width=2, height=1, label="A")
        b = BoundingBox(x=1.5, y=0.5, width=2, height=1, label="B")
        assert a.overlaps(b) is True

    def test_bounding_box_no_overlap(self):
        from agents.slide_agent.layout_engine import BoundingBox

        a = BoundingBox(x=0, y=0, width=2, height=1, label="A")
        b = BoundingBox(x=3, y=0, width=2, height=1, label="B")
        assert a.overlaps(b) is False

    def test_bounding_box_gap_horizontal(self):
        from agents.slide_agent.layout_engine import BoundingBox

        a = BoundingBox(x=0, y=0, width=2, height=1)
        b = BoundingBox(x=3, y=0, width=2, height=1)
        assert a.gap_horizontal(b) == 1.0  # 1 inch gap

    def test_bounding_box_gap_negative_overlap(self):
        from agents.slide_agent.layout_engine import BoundingBox

        a = BoundingBox(x=0, y=0, width=2, height=1)
        b = BoundingBox(x=1.5, y=0, width=2, height=1)
        assert a.gap_horizontal(b) < 0  # overlap = negative gap

    def test_bounding_box_intersection_area(self):
        from agents.slide_agent.layout_engine import BoundingBox

        a = BoundingBox(x=0, y=0, width=2, height=2)
        b = BoundingBox(x=1, y=1, width=2, height=2)
        assert a.intersection_area(b) == 1.0  # 1x1 overlap

    def test_validate_overlap_detected(self):
        from agents.slide_agent.layout_engine import BoundingBox, LayoutEngine

        engine = LayoutEngine()
        elements = [
            BoundingBox(x=1, y=1, width=5, height=3, label="A"),
            BoundingBox(x=3, y=2, width=5, height=3, label="B"),
        ]
        report = engine.validate(elements)
        assert not report.valid
        assert any(i.code == "OVERLAP" for i in report.issues)

    def test_validate_out_of_bounds(self):
        from agents.slide_agent.layout_engine import BoundingBox, LayoutEngine

        engine = LayoutEngine()
        elements = [
            BoundingBox(x=10, y=0, width=5, height=2, label="A"),  # extends past 13.333
        ]
        report = engine.validate(elements)
        assert not report.valid
        assert any(i.code == "OUT_OF_BOUNDS" for i in report.issues)

    def test_validate_clean_slide(self):
        from agents.slide_agent.layout_engine import BoundingBox, LayoutEngine

        engine = LayoutEngine()
        elements = [
            BoundingBox(x=0.5, y=0.5, width=5, height=2, label="A"),
            BoundingBox(x=6, y=0.5, width=5, height=2, label="B"),
        ]
        report = engine.validate(elements)
        assert report.valid

    def test_validate_high_density(self):
        from agents.slide_agent.layout_engine import BoundingBox, LayoutEngine

        engine = LayoutEngine()
        # Fill 80% of the safe zone
        elements = [
            BoundingBox(x=0.5, y=0.5, width=12, height=6.5, label="huge"),
        ]
        report = engine.validate(elements)
        assert not report.valid
        assert any(i.code == "HIGH_DENSITY" for i in report.issues)

    def test_auto_layout_grid(self):
        from agents.slide_agent.layout_engine import BoundingBox, LayoutEngine

        engine = LayoutEngine()
        elements = [
            BoundingBox(x=0, y=0, width=2, height=1, label=f"box_{i}") for i in range(6)
        ]
        result = engine.auto_layout_grid(elements, columns=3)
        assert len(result) == 6
        # Verify no overlaps after layout
        report = engine.validate(result)
        assert not any(i.code == "OVERLAP" for i in report.issues)

    def test_auto_layout_stack(self):
        from agents.slide_agent.layout_engine import BoundingBox, LayoutEngine

        engine = LayoutEngine()
        elements = [
            BoundingBox(x=0, y=0, width=4, height=0.8, label=f"item_{i}")
            for i in range(4)
        ]
        result = engine.auto_layout_stack(elements, direction="vertical", gap=0.3)
        report = engine.validate(result)
        assert not any(i.code == "OVERLAP" for i in report.issues)

    def test_extract_bounding_boxes(self):
        from agents.slide_agent.slide_graph import _extract_bounding_boxes

        slide_def = {
            "slide_number": 1,
            "layout": "content",
            "title": "Test Slide",
            "content": {"type": "bullets", "items": ["A", "B"]},
        }
        boxes = _extract_bounding_boxes(slide_def)
        assert len(boxes) >= 2  # title + content
        assert any("title" in b.label for b in boxes)


class TestTemplateGuided:
    """Test template-guided design — introspection and assignment-based population."""

    def test_extract_layouts(self):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        layouts = builder.extract_layouts(prs)
        assert len(layouts) > 0
        # Title Slide should have title and subtitle placeholders
        title_slide = layouts[0]
        assert title_slide["name"] == "Title Slide"
        roles = [p["role"] for p in title_slide["placeholders"]]
        assert "title" in roles
        assert "subtitle" in roles

    def test_extract_layouts_has_zones(self):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder
        from pptx import Presentation

        prs = Presentation()
        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        layouts = builder.extract_layouts(prs)

        for layout in layouts:
            for ph in layout["placeholders"]:
                assert "zone" in ph
                assert "x_inches" in ph["zone"]
                assert "width_inches" in ph["zone"]
                assert ph["zone"]["width_inches"] > 0

    def test_build_with_assignments(self, tmp_path):
        pytest.importorskip("pptx")
        from agents.slide_agent.slide_builder import SlideBuilder

        builder = SlideBuilder(Path(__file__).parent.parent / "templates")
        plan = {
            "deck": {
                "title": "Template Guided",
                "template": "blank",
                "theme": {
                    "font_family": "Calibri",
                    "font_size_title": 28,
                    "font_size_body": 16,
                },
            },
            "slides": [
                {
                    "slide_number": 1,
                    "layout": "title_slide",
                    "title": "Guided Title",
                    "subtitle": "Guided Subtitle",
                },
                {
                    "slide_number": 2,
                    "layout": "content",
                    "title": "Content Slide",
                    "assignments": {
                        "1": {
                            "type": "body",
                            "items": ["Bullet A", "Bullet B", "Bullet C"],
                        },
                    },
                },
            ],
        }
        output_path = tmp_path / "guided.pptx"
        result = builder.build_deck(plan, output_path)
        assert result.exists()
        assert result.stat().st_size > 0


class TestGraphDelegation:
    @pytest.mark.asyncio
    async def test_langgraph_suspends_for_generated_asset(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from unittest.mock import MagicMock

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import TaskEnvelope, TaskInProgress

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        async def fake_plan_deck(**kwargs):
            return {
                "deck": {"title": "Delegated Deck", "template": "blank"},
                "slides": [
                    {
                        "slide_number": 1,
                        "layout": "content_with_image",
                        "title": "Hello",
                        "content": {"type": "bullets", "items": ["A", "B"]},
                        "image": {
                            "type": "image",
                            "source": {
                                "kind": "generate",
                                "agent": "image",
                                "prompt": "simple hero illustration",
                            },
                            "placement": {
                                "x_inches": 6.0,
                                "y_inches": 1.5,
                                "width_inches": 5.0,
                                "height_inches": 3.5,
                            },
                        },
                    }
                ],
            }

        delegate_calls: list[dict[str, object]] = []
        emitted: list[tuple[str, str, dict[str, object]]] = []

        async def fake_request_orchestrator_delegate(**kwargs):
            delegate_calls.append(kwargs)
            return {"reverse_task_id": "rev_slide_asset_1", "status": "registered"}

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            emitted.append((task_id, event_type, payload))

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_deck", fake_plan_deck)
        monkeypatch.setattr(agent, "request_orchestrator_delegate", fake_request_orchestrator_delegate)
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        task = TaskEnvelope(
            task_id="tsk_slide_suspend",
            task_list_id="tasks:slide",
            session_id="sess_slide",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={"description": "Make me a one-slide deck with a generated hero image"},
            idempotency_key="idem_slide_suspend",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await run_slide_langgraph(agent=agent, task=task)

        assert isinstance(result, TaskInProgress)
        assert len(delegate_calls) == 1
        assert delegate_calls[0]["target_intent"] == "image.generate"
        assert delegate_calls[0]["target_agent_id"] == cfg.image_agent_id
        resume_payload = delegate_calls[0]["resume_payload"]
        assert resume_payload["pending_asset_request"]["slide_number"] == 1
        assert resume_payload["pending_asset_request"]["slot"] == "image"
        assert any(event_type == "task.suspended" for _, event_type, _ in emitted)

    @pytest.mark.asyncio
    async def test_langgraph_resume_applies_delegated_asset(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from io import BytesIO
        from unittest.mock import MagicMock

        from PIL import Image

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import AgentResult, TaskEnvelope

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        png_path = tmp_path / "delegated_asset.png"
        img = Image.new("RGB", (64, 64), color="navy")
        buf = BytesIO()
        img.save(buf, format="PNG")
        png_path.write_bytes(buf.getvalue())

        built_plans: list[dict[str, object]] = []

        async def should_not_replan(**kwargs):
            raise AssertionError("resume path should not call plan_deck again")

        def fake_build_deck(self, plan, output_path):
            built_plans.append(plan)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(b"pptx")
            return output_path

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            return None

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_deck", should_not_replan)
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.SlideBuilder.build_deck",
            fake_build_deck,
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.render_slides_to_png",
            lambda *args, **kwargs: [],
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.export_to_pdf",
            lambda *args, **kwargs: None,
        )
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        task = TaskEnvelope(
            task_id="tsk_slide_resume",
            task_list_id="tasks:slide",
            session_id="sess_slide",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={
                "description": "resume",
                "_resume": {
                    "resume_state": {
                        "tool_round": 1,
                        "description": "Make me a one-slide deck with a generated hero image",
                        "template": "blank",
                        "llm_deck_plan": {
                            "deck": {"title": "Delegated Deck", "template": "blank"},
                            "slides": [
                                {
                                    "slide_number": 1,
                                    "layout": "content_with_image",
                                    "title": "Hello",
                                    "content": {
                                        "type": "bullets",
                                        "items": ["A", "B"],
                                    },
                                    "image": {
                                        "type": "image",
                                        "source": {
                                            "kind": "generate",
                                            "agent": "image",
                                            "prompt": "simple hero illustration",
                                        },
                                        "placement": {
                                            "x_inches": 6.0,
                                            "y_inches": 1.5,
                                            "width_inches": 5.0,
                                            "height_inches": 3.5,
                                        },
                                    },
                                }
                            ],
                        },
                        "pending_asset_request": {
                            "slide_number": 1,
                            "slot": "image",
                            "agent_type": "image",
                            "target_intent": "image.generate",
                        },
                    },
                    "reverse_task": {
                        "reverse_task_id": "rev_slide_asset_1",
                        "target_intent": "image.generate",
                        "target_agent_id": cfg.image_agent_id,
                    },
                    "reverse_result": {
                        "status": "completed",
                        "output": {},
                        "artifacts": [
                            {
                                "artifact_id": "art_slide_image",
                                "task_id": "tsk_img_1",
                                "mime": "image/png",
                                "path": str(png_path),
                                "sha256": "abc123",
                                "created_by_agent": cfg.image_agent_id,
                                "kind": "output",
                                "audience": "deliverable",
                            }
                        ],
                    },
                },
            },
            idempotency_key="idem_slide_resume",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await run_slide_langgraph(agent=agent, task=task)

        assert isinstance(result, AgentResult)
        assert result.status == "completed"
        assert built_plans
        image_source = built_plans[0]["slides"][0]["image"]["source"]
        assert image_source["image_bytes"] == png_path.read_bytes()
        assert any(a.mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation" for a in result.artifacts)

    @pytest.mark.asyncio
    async def test_langgraph_edit_suspends_and_resumes_generated_replace_image(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from io import BytesIO
        from unittest.mock import MagicMock

        from PIL import Image

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import AgentResult, TaskEnvelope, TaskInProgress

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        source_pptx = tmp_path / "source_edit_resume.pptx"
        source_pptx.write_bytes(b"source")

        png_path = tmp_path / "delegated_edit_asset.png"
        img = Image.new("RGB", (64, 64), color="teal")
        buf = BytesIO()
        img.save(buf, format="PNG")
        png_path.write_bytes(buf.getvalue())

        delegate_calls: list[dict[str, object]] = []
        emitted: list[tuple[str, str, dict[str, object]]] = []
        applied_operations: list[list[dict[str, object]]] = []

        async def fake_plan_edit(**kwargs):
            return {
                "action": "edit",
                "operations": [
                    {
                        "action": "replace_image",
                        "slide_number": 1,
                        "shape_name": "HeroImage",
                        "new_image": {
                            "source": {
                                "kind": "generate",
                                "agent": "image",
                                "prompt": "replacement diagram hero",
                            }
                        },
                    }
                ],
            }

        async def fake_request_orchestrator_delegate(**kwargs):
            delegate_calls.append(kwargs)
            return {"reverse_task_id": "rev_slide_edit_asset_1", "status": "registered"}

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            emitted.append((task_id, event_type, payload))

        class DummyPresentation:
            def save(self, path: str):
                Path(path).write_bytes(b"pptx")

        class DummyBuilder:
            def __init__(self, templates_dir):
                self.templates_dir = templates_dir

            def load_existing(self, pptx_path: Path):
                return DummyPresentation()

            def extract_structure(self, prs):
                return {
                    "slide_count": 1,
                    "slides": [
                        {
                            "slide_number": 1,
                            "title": "Before",
                            "shapes": [{"name": "HeroImage"}],
                        }
                    ],
                }

            def apply_edits(self, prs, operations):
                applied_operations.append(operations)
                return prs

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_edit", fake_plan_edit)
        monkeypatch.setattr("agents.slide_agent.slide_graph.SlideBuilder", DummyBuilder)
        monkeypatch.setattr(agent, "request_orchestrator_delegate", fake_request_orchestrator_delegate)
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        suspend_task = TaskEnvelope(
            task_id="tsk_slide_edit_suspend",
            task_list_id="tasks:slide",
            session_id="sess_slide",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.edit",
            input={
                "source_pptx_path": str(source_pptx),
                "edit_request": "Replace the hero image with something fresher",
            },
            idempotency_key="idem_slide_edit_suspend",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        suspended = await run_slide_langgraph(agent=agent, task=suspend_task)

        assert isinstance(suspended, TaskInProgress)
        assert len(delegate_calls) == 1
        resume_payload = delegate_calls[0]["resume_payload"]
        assert resume_payload["pending_asset_request"]["operation_index"] == 0
        assert any(event_type == "task.suspended" for _, event_type, _ in emitted)

        async def should_not_replan_edit(**kwargs):
            raise AssertionError("resume path should not call plan_edit again")

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_edit", should_not_replan_edit)
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.render_slides_to_png",
            lambda *args, **kwargs: [],
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.export_to_pdf",
            lambda *args, **kwargs: None,
        )

        resume_task = TaskEnvelope(
            task_id="tsk_slide_edit_resume",
            task_list_id="tasks:slide",
            session_id="sess_slide",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.edit",
            input={
                "source_pptx_path": str(source_pptx),
                "edit_request": "Replace the hero image with something fresher",
                "_resume": {
                    "resume_state": resume_payload,
                    "reverse_task": {
                        "reverse_task_id": "rev_slide_edit_asset_1",
                        "target_intent": "image.generate",
                        "target_agent_id": cfg.image_agent_id,
                    },
                    "reverse_result": {
                        "status": "completed",
                        "output": {},
                        "artifacts": [
                            {
                                "artifact_id": "art_slide_edit_image",
                                "task_id": "tsk_img_edit_1",
                                "mime": "image/png",
                                "path": str(png_path),
                                "sha256": "abc123",
                                "created_by_agent": cfg.image_agent_id,
                                "kind": "output",
                                "audience": "deliverable",
                            }
                        ],
                    },
                },
            },
            idempotency_key="idem_slide_edit_resume",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        resumed = await run_slide_langgraph(agent=agent, task=resume_task)

        assert isinstance(resumed, AgentResult)
        assert resumed.status == "completed"
        assert applied_operations
        replacement = applied_operations[0][0]["new_image"]
        assert replacement["image_bytes"] == png_path.read_bytes()
        assert any(
            a.mime
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            for a in resumed.artifacts
        )

    @pytest.mark.asyncio
    async def test_langgraph_suspends_for_raw_document_source_artifacts(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from unittest.mock import MagicMock

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import TaskEnvelope, TaskInProgress

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        pdf_path = tmp_path / "source.pdf"
        pdf_path.write_bytes(b"%PDF-1.4 raw")

        delegate_calls: list[dict[str, object]] = []
        emitted: list[tuple[str, str, dict[str, object]]] = []

        async def should_not_plan(**kwargs):
            raise AssertionError("raw document inputs should suspend before planning")

        async def fake_request_orchestrator_delegate(**kwargs):
            delegate_calls.append(kwargs)
            return {"reverse_task_id": "rev_slide_docs_1", "status": "registered"}

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            emitted.append((task_id, event_type, payload))

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_deck", should_not_plan)
        monkeypatch.setattr(
            agent, "request_orchestrator_delegate", fake_request_orchestrator_delegate
        )
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        task = TaskEnvelope(
            task_id="tsk_slide_docs_suspend",
            task_list_id="tasks:slide",
            session_id="sess_slide_docs",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={"description": "Make slides from this uploaded PDF"},
            input_artifacts=[
                {
                    "artifact_id": "art_pdf_1",
                    "path": str(pdf_path),
                    "mime": "application/pdf",
                    "filename": "source.pdf",
                }
            ],
            idempotency_key="idem_slide_docs_suspend",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await run_slide_langgraph(agent=agent, task=task)

        assert isinstance(result, TaskInProgress)
        assert len(delegate_calls) == 1
        assert delegate_calls[0]["target_intent"] == "docs.parse_bundle"
        assert delegate_calls[0]["target_agent_id"] == cfg.docs_parser_agent_id
        assert len(delegate_calls[0]["input_artifacts"]) == 1
        resume_payload = delegate_calls[0]["resume_payload"]
        assert resume_payload["pending_asset_request"]["request_kind"] == "docs_parse"
        assert any(event_type == "task.suspended" for _, event_type, _ in emitted)

    @pytest.mark.asyncio
    async def test_langgraph_resume_docs_parse_hydrates_source_asset(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from io import BytesIO
        from unittest.mock import MagicMock

        from PIL import Image

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import AgentResult, TaskEnvelope

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        pdf_path = tmp_path / "source.pdf"
        pdf_path.write_bytes(b"%PDF-1.4 raw")

        bundle_dir = tmp_path / "docs_bundle"
        bundle_dir.mkdir(parents=True, exist_ok=True)
        figure_path = bundle_dir / "figure-1.png"
        img = Image.new("RGB", (128, 72), color="orange")
        buf = BytesIO()
        img.save(buf, format="PNG")
        figure_bytes = buf.getvalue()
        figure_path.write_bytes(figure_bytes)

        manifest_path = bundle_dir / "manifest.json"
        manifest_path.write_text(
            json.dumps(
                {
                    "doc_id": "doc_1",
                    "source_artifact_id": "art_pdf_1",
                    "filename": "source.pdf",
                    "mime": "application/pdf",
                    "title": "Uploaded PDF",
                    "counts": {
                        "section_count": 4,
                        "chunk_count": 9,
                        "table_count": 1,
                        "figure_count": 1,
                        "page_count": 3,
                    },
                }
            ),
            encoding="utf-8",
        )
        chunk_index_path = bundle_dir / "chunk_index.json"
        chunk_index_path.write_text(
            json.dumps(
                {
                    "figures": [
                        {
                            "asset_id": "fig_asset_1",
                            "description": "A compelling architecture figure",
                            "caption": "Architecture overview",
                            "page_number": 2,
                        }
                    ],
                    "pages": [],
                    "slides": [],
                    "assets": [
                        {
                            "asset_id": "fig_asset_1",
                            "kind": "figure_image",
                            "path": "figure-1.png",
                            "mime": "image/png",
                            "description": "A compelling architecture figure",
                            "page_number": 2,
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        built_plans: list[dict[str, object]] = []

        async def fake_plan_deck(**kwargs):
            source_materials = kwargs["input_data"]["_source_materials"]
            assert source_materials["documents"][0]["artifact_id"] == "art_pdf_1"
            assert source_materials["visual_assets"][0]["asset_ref"] == "fig_asset_1"
            return {
                "deck": {"title": "Doc Deck", "template": "blank"},
                "slides": [
                    {
                        "slide_number": 1,
                        "layout": "content_with_image",
                        "title": "PDF-derived slide",
                        "content": {"type": "bullets", "items": ["Point A", "Point B"]},
                        "image": {
                            "type": "image",
                            "source": {
                                "kind": "from_asset",
                                "asset_ref": "fig_asset_1",
                            },
                            "placement": {
                                "x_inches": 6.0,
                                "y_inches": 1.5,
                                "width_inches": 5.0,
                                "height_inches": 3.5,
                            },
                        },
                    }
                ],
            }

        def fake_build_deck(self, plan, output_path):
            built_plans.append(plan)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(b"pptx")
            return output_path

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            return None

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_deck", fake_plan_deck)
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.SlideBuilder.build_deck",
            fake_build_deck,
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.render_slides_to_png",
            lambda *args, **kwargs: [],
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.export_to_pdf",
            lambda *args, **kwargs: None,
        )
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        task = TaskEnvelope(
            task_id="tsk_slide_docs_resume",
            task_list_id="tasks:slide",
            session_id="sess_slide_docs",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={
                "description": "Build slides from this uploaded PDF",
                "_resume": {
                    "resume_state": {
                        "tool_round": 1,
                        "description": "Build slides from this uploaded PDF",
                        "template": "blank",
                        "source_artifacts": [
                            {
                                "artifact_id": "art_pdf_1",
                                "path": str(pdf_path),
                                "mime": "application/pdf",
                                "filename": "source.pdf",
                            }
                        ],
                        "source_documents": [],
                        "source_visual_assets": [],
                        "source_artifacts_prepared": False,
                        "pending_asset_request": {
                            "request_kind": "docs_parse",
                            "target_intent": "docs.parse_bundle",
                            "target_agent_id": cfg.docs_parser_agent_id,
                            "source_artifact_ids": ["art_pdf_1"],
                        },
                    },
                    "reverse_task": {
                        "reverse_task_id": "rev_slide_docs_1",
                        "target_intent": "docs.parse_bundle",
                        "target_agent_id": cfg.docs_parser_agent_id,
                    },
                    "reverse_result": {
                        "status": "completed",
                        "output": {
                            "documents": [
                                {
                                    "doc_id": "doc_1",
                                    "artifact_id": "art_pdf_1",
                                    "filename": "source.pdf",
                                    "mime": "application/pdf",
                                    "title": "Uploaded PDF",
                                    "section_count": 4,
                                    "chunk_count": 9,
                                    "table_count": 1,
                                    "figure_count": 1,
                                    "artifact_refs": ["fig_asset_1"],
                                }
                            ]
                        },
                        "artifacts": [
                            {
                                "artifact_id": "art_pdf_1_manifest",
                                "task_id": "tsk_docs_parse",
                                "mime": "application/json",
                                "path": str(manifest_path),
                                "sha256": "abc",
                                "created_by_agent": cfg.docs_parser_agent_id,
                                "kind": "output",
                                "audience": "supporting",
                            },
                            {
                                "artifact_id": "art_pdf_1_chunk_index",
                                "task_id": "tsk_docs_parse",
                                "mime": "application/json",
                                "path": str(chunk_index_path),
                                "sha256": "def",
                                "created_by_agent": cfg.docs_parser_agent_id,
                                "kind": "output",
                                "audience": "supporting",
                            },
                        ],
                    },
                },
            },
            idempotency_key="idem_slide_docs_resume",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await run_slide_langgraph(agent=agent, task=task)

        assert isinstance(result, AgentResult)
        assert result.status == "completed"
        assert built_plans
        image_source = built_plans[0]["slides"][0]["image"]["source"]
        assert image_source["asset_ref"] == "fig_asset_1"
        assert image_source["image_bytes"] == figure_bytes

    def test_document_summaries_include_local_bundle_preview_and_sections(
        self, tmp_path: Path
    ):
        from agents.slide_agent.slide_graph import _document_summaries_from_artifacts

        bundle_dir = tmp_path / "bundle_local"
        bundle_dir.mkdir(parents=True, exist_ok=True)
        manifest_path = bundle_dir / "manifest.json"
        document_md_path = bundle_dir / "document.md"
        chunk_index_path = bundle_dir / "chunk_index.json"

        document_md_path.write_text(
            "# Revenue Story\n\nAcme revenue grew 41 percent year over year while margins improved.",
            encoding="utf-8",
        )
        chunk_index_path.write_text(
            json.dumps(
                {
                    "sections": [
                        {
                            "section_id": "sec_growth",
                            "title": "Growth",
                            "page_number": 2,
                            "summary": "Revenue acceleration and product-led expansion.",
                        },
                        {
                            "section_id": "sec_margin",
                            "title": "Margins",
                            "page_number": 3,
                            "summary": "Operating leverage and gross margin improvement.",
                        },
                    ]
                }
            ),
            encoding="utf-8",
        )
        manifest_path.write_text(
            json.dumps(
                {
                    "doc_id": "doc_revenue",
                    "bundle_id": "bundle_revenue",
                    "source_artifact_id": "art_pdf_revenue",
                    "filename": "revenue.pdf",
                    "mime": "application/pdf",
                    "title": "Revenue Update",
                    "outputs": {
                        "document_md": "document.md",
                        "chunk_index": "chunk_index.json",
                    },
                    "counts": {
                        "section_count": 2,
                        "chunk_count": 4,
                        "figure_count": 0,
                        "page_count": 3,
                    },
                }
            ),
            encoding="utf-8",
        )

        summaries = _document_summaries_from_artifacts(
            [
                {
                    "artifact_id": "art_manifest_revenue",
                    "path": str(manifest_path),
                    "mime": "application/json",
                    "filename": "manifest.json",
                    "created_by_agent": "cosmic/docs-parser-agent:1.0.0",
                }
            ]
        )

        assert len(summaries) == 1
        summary = summaries[0]
        assert summary["bundle_id"] == "bundle_revenue"
        assert summary["artifact_id"] == "art_pdf_revenue"
        assert "41 percent" in (summary.get("preview_excerpt") or "")
        assert summary["top_sections"][0]["section_id"] == "sec_growth"
        assert summary["paths"]["document_md"] == str(document_md_path).replace("\\", "/")

    @pytest.mark.asyncio
    async def test_langgraph_suspends_for_deeper_doc_context_request(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from unittest.mock import MagicMock

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import TaskEnvelope, TaskInProgress

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        bundle_dir = tmp_path / "bundle_context_suspend"
        bundle_dir.mkdir(parents=True, exist_ok=True)
        manifest_path = bundle_dir / "manifest.json"
        document_md_path = bundle_dir / "document.md"
        chunk_index_path = bundle_dir / "chunk_index.json"

        document_md_path.write_text(
            "# Board Memo\n\nThe company grew enterprise revenue and improved retention in Europe.",
            encoding="utf-8",
        )
        chunk_index_path.write_text(
            json.dumps({"sections": [{"section_id": "sec_1", "title": "Growth"}]}),
            encoding="utf-8",
        )
        manifest_path.write_text(
            json.dumps(
                {
                    "doc_id": "doc_board",
                    "bundle_id": "bundle_board",
                    "source_artifact_id": "art_pdf_board",
                    "filename": "board.pdf",
                    "mime": "application/pdf",
                    "title": "Board Memo",
                    "outputs": {
                        "document_md": "document.md",
                        "chunk_index": "chunk_index.json",
                    },
                    "counts": {"section_count": 1, "chunk_count": 2},
                }
            ),
            encoding="utf-8",
        )

        delegate_calls: list[dict[str, object]] = []
        emitted: list[tuple[str, str, dict[str, object]]] = []

        async def fake_plan_deck(**kwargs):
            source_materials = kwargs["input_data"]["_source_materials"]
            assert source_materials["documents"][0]["bundle_id"] == "bundle_board"
            assert source_materials["documents"][0]["preview_excerpt"]
            return {
                "action": "request_doc_context",
                "doc_request": {
                    "intent": "docs.search_bundle",
                    "bundle_id": "bundle_board",
                    "doc_id": "doc_board",
                    "query": "retention in Europe",
                    "search_kind": "sections",
                },
            }

        async def fake_request_orchestrator_delegate(**kwargs):
            delegate_calls.append(kwargs)
            return {"reverse_task_id": "rev_slide_doc_ctx_1", "status": "registered"}

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            emitted.append((task_id, event_type, payload))

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_deck", fake_plan_deck)
        monkeypatch.setattr(
            agent, "request_orchestrator_delegate", fake_request_orchestrator_delegate
        )
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        task = TaskEnvelope(
            task_id="tsk_slide_doc_ctx_suspend",
            task_list_id="tasks:slide",
            session_id="sess_slide_doc_ctx",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={"description": "Make investor slides from this board memo"},
            input_artifacts=[
                {
                    "artifact_id": "art_manifest_board",
                    "path": str(manifest_path),
                    "mime": "application/json",
                    "filename": "manifest.json",
                    "created_by_agent": cfg.docs_parser_agent_id,
                }
            ],
            idempotency_key="idem_slide_doc_ctx_suspend",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await run_slide_langgraph(agent=agent, task=task)

        assert isinstance(result, TaskInProgress)
        assert len(delegate_calls) == 1
        assert delegate_calls[0]["target_intent"] == "docs.search_bundle"
        assert delegate_calls[0]["target_agent_id"] == cfg.docs_parser_agent_id
        assert delegate_calls[0]["target_input"]["bundle_id"] == "bundle_board"
        assert delegate_calls[0]["target_input"]["doc_ids"] == ["doc_board"]
        resume_payload = delegate_calls[0]["resume_payload"]
        assert resume_payload["pending_asset_request"]["request_kind"] == "docs_context"
        assert resume_payload["doc_context_request_count"] == 1
        assert any(event_type == "task.suspended" for _, event_type, _ in emitted)

    @pytest.mark.asyncio
    async def test_langgraph_resume_docs_context_replans_with_compact_context(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from unittest.mock import MagicMock

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from agents.slide_agent.slide_graph import run_slide_langgraph
        from shared.contracts import AgentResult, TaskEnvelope

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        bundle_dir = tmp_path / "bundle_context_resume"
        bundle_dir.mkdir(parents=True, exist_ok=True)
        manifest_path = bundle_dir / "manifest.json"
        document_md_path = bundle_dir / "document.md"
        chunk_index_path = bundle_dir / "chunk_index.json"

        document_md_path.write_text(
            "# Board Memo\n\nEurope retention improved and enterprise revenue accelerated.",
            encoding="utf-8",
        )
        chunk_index_path.write_text(
            json.dumps({"sections": [{"section_id": "sec_1", "title": "Retention"}]}),
            encoding="utf-8",
        )
        manifest_path.write_text(
            json.dumps(
                {
                    "doc_id": "doc_board",
                    "bundle_id": "bundle_board",
                    "source_artifact_id": "art_pdf_board",
                    "filename": "board.pdf",
                    "mime": "application/pdf",
                    "title": "Board Memo",
                    "outputs": {
                        "document_md": "document.md",
                        "chunk_index": "chunk_index.json",
                    },
                    "counts": {"section_count": 1, "chunk_count": 2},
                }
            ),
            encoding="utf-8",
        )

        built_plans: list[dict[str, object]] = []

        async def fake_plan_deck(**kwargs):
            doc_context = kwargs["input_data"]["_document_context"]
            assert doc_context["items"][0]["kind"] == "search_hits"
            assert doc_context["items"][0]["query"] == "retention in Europe"
            return {
                "deck": {"title": "Investor Update", "template": "blank"},
                "slides": [
                    {
                        "slide_number": 1,
                        "layout": "content",
                        "title": "Retention",
                        "content": {
                            "type": "bullets",
                            "items": ["Europe retention improved", "Enterprise revenue accelerated"],
                        },
                        "speaker_notes": "Summarize retention and growth.",
                    }
                ],
            }

        def fake_build_deck(self, plan, output_path):
            built_plans.append(plan)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(b"pptx")
            return output_path

        async def fake_emit_event(task_id: str, event_type: str, payload: dict[str, object]):
            return None

        monkeypatch.setattr("agents.slide_agent.slide_graph.plan_deck", fake_plan_deck)
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.SlideBuilder.build_deck",
            fake_build_deck,
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.render_slides_to_png",
            lambda *args, **kwargs: [],
        )
        monkeypatch.setattr(
            "agents.slide_agent.slide_graph.export_to_pdf",
            lambda *args, **kwargs: None,
        )
        monkeypatch.setattr(agent, "emit_event", fake_emit_event)

        task = TaskEnvelope(
            task_id="tsk_slide_doc_ctx_resume",
            task_list_id="tasks:slide",
            session_id="sess_slide_doc_ctx",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={
                "description": "Make investor slides from this board memo",
                "_resume": {
                    "resume_state": {
                        "tool_round": 1,
                        "description": "Make investor slides from this board memo",
                        "template": "blank",
                        "source_artifacts": [
                            {
                                "artifact_id": "art_manifest_board",
                                "path": str(manifest_path),
                                "mime": "application/json",
                                "filename": "manifest.json",
                                "created_by_agent": cfg.docs_parser_agent_id,
                            }
                        ],
                        "source_documents": [],
                        "source_visual_assets": [],
                        "source_document_context": [],
                        "source_artifacts_prepared": True,
                        "doc_context_request_count": 1,
                        "max_doc_context_requests": 2,
                        "pending_asset_request": {
                            "request_kind": "docs_context",
                            "target_intent": "docs.search_bundle",
                            "target_agent_id": cfg.docs_parser_agent_id,
                            "bundle_id": "bundle_board",
                            "doc_ids": ["doc_board"],
                            "query": "retention in Europe",
                            "search_kind": "sections",
                        },
                    },
                    "reverse_task": {
                        "reverse_task_id": "rev_slide_doc_ctx_1",
                        "target_intent": "docs.search_bundle",
                        "target_agent_id": cfg.docs_parser_agent_id,
                    },
                    "reverse_result": {
                        "status": "completed",
                        "output": {
                            "bundle_id": "bundle_board",
                            "query": "retention in Europe",
                            "search_kind": "sections",
                            "matches": [
                                {
                                    "doc_id": "doc_board",
                                    "title": "Board Memo",
                                    "chunk_id": "chunk_1",
                                    "section_id": "sec_1",
                                    "section_title": "Retention",
                                    "score": 93,
                                    "excerpt": "Europe retention improved sharply in enterprise cohorts.",
                                }
                            ],
                        },
                    },
                },
            },
            idempotency_key="idem_slide_doc_ctx_resume",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await run_slide_langgraph(agent=agent, task=task)

        assert isinstance(result, AgentResult)
        assert result.status == "completed"
        assert built_plans
        assert built_plans[0]["deck"]["title"] == "Investor Update"


class TestDirectFallbackParity:
    @pytest.mark.asyncio
    async def test_slide_edit_fallback_validates_and_exports(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ):
        from io import BytesIO
        from unittest.mock import MagicMock

        from PIL import Image

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from shared.contracts import TaskEnvelope

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            slide_use_langgraph=False,
            export_pdf=True,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        source_pptx = tmp_path / "source.pptx"
        source_pptx.write_bytes(b"source")

        async def fake_plan_edit(**kwargs):
            return {"operations": [{"action": "update_text"}]}

        class DummyPresentation:
            def save(self, path: str):
                Path(path).write_bytes(b"pptx")

        class DummyBuilder:
            def __init__(self, templates_dir):
                self.templates_dir = templates_dir

            def load_existing(self, pptx_path: Path):
                return DummyPresentation()

            def extract_structure(self, prs):
                return {"slide_count": 1, "slides": [{"slide_number": 1, "title": "Before"}]}

            def apply_edits(self, prs, operations):
                return prs

        preview_dir = tmp_path / "previews"
        preview_dir.mkdir(parents=True, exist_ok=True)
        img = Image.new("RGB", (96, 54), color="purple")
        buf = BytesIO()
        img.save(buf, format="PNG")
        preview_path = preview_dir / "slide-1.png"
        preview_path.write_bytes(buf.getvalue())

        async def fake_validate_slide(**kwargs):
            return {"pass": True, "issues": []}

        def fake_export_to_pdf(pptx_path: Path, libreoffice_path: str):
            pdf_path = pptx_path.with_suffix(".pdf")
            pdf_path.write_bytes(b"%PDF-1.4")
            return pdf_path

        monkeypatch.setattr("agents.slide_agent.internal_llm.plan_edit", fake_plan_edit)
        monkeypatch.setattr("agents.slide_agent.internal_llm.validate_slide", fake_validate_slide)
        monkeypatch.setattr("agents.slide_agent.agent.SlideBuilder", DummyBuilder)
        monkeypatch.setattr(
            "agents.slide_agent.agent.render_slides_to_png",
            lambda *args, **kwargs: [preview_path],
        )
        monkeypatch.setattr("agents.slide_agent.agent.export_to_pdf", fake_export_to_pdf)

        task = TaskEnvelope(
            task_id="tsk_slide_edit_fallback",
            task_list_id="tasks:slide",
            session_id="sess_slide",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.edit",
            input={
                "source_pptx_path": str(source_pptx),
                "edit_request": "Change the title and tighten the copy",
            },
            idempotency_key="idem_slide_edit_fallback",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await agent.handle_slide_edit(task)

        assert result.status == "completed"
        assert result.output["pdf_path"].endswith(".pdf")
        assert result.output["validation_pass"] is True
        assert any(a.mime == "application/pdf" for a in result.artifacts)
        assert any(a.mime == "image/png" for a in result.artifacts)

    @pytest.mark.asyncio
    async def test_slide_create_fallback_rejects_raw_document_sources(
        self, tmp_path: Path
    ):
        from unittest.mock import MagicMock

        from agents.slide_agent.agent import SlideAgent
        from agents.slide_agent.config import SlideAgentConfig
        from shared.contracts import TaskEnvelope

        cfg = SlideAgentConfig(
            artifacts_root=tmp_path / "runs" / "artifacts",
            slide_use_langgraph=False,
            export_pdf=False,
        )
        agent = SlideAgent(redis_client=MagicMock(), config=cfg)

        pdf_path = tmp_path / "source.pdf"
        pdf_path.write_bytes(b"%PDF-1.4 raw")

        task = TaskEnvelope(
            task_id="tsk_slide_create_fallback_docs",
            task_list_id="tasks:slide",
            session_id="sess_slide_docs",
            sender="cosmic/orchestrator:1.0.0",
            recipient="cosmic/slide-agent:1.0.0",
            intent="slide.create",
            input={"description": "Build slides from this uploaded PDF"},
            input_artifacts=[
                {
                    "artifact_id": "art_pdf_1",
                    "path": str(pdf_path),
                    "mime": "application/pdf",
                    "filename": "source.pdf",
                }
            ],
            idempotency_key="idem_slide_create_fallback_docs",
            signature="sig",
            source="agent",
            channel="desktop",
        )

        result = await agent.handle_slide_create(task)

        assert result.status == "failed"
        assert result.error is not None
        assert result.error.code == "FALLBACK_SOURCE_ARTIFACT_PREPARATION_UNAVAILABLE"
