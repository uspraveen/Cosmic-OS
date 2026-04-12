"""Slide builder — cosmic-slides-2.

LangGraph-based PPTX builder. It can either:
- start from an existing `build_spec` (legacy behavior), or
- plan content + select layouts inside the graph before building.

Pipeline
────────
  plan_content → select_layouts → load → map_shapes → polish_mappings → fetch_assets
       → build_pptx → render → validate ─┐
    ├─ (all pass / max repairs / no validation) → finalize → END
    └─ (failures) → repair → build_pptx → render → validate ────────────────────────────┘

Key capabilities
────────────────
  • Vision-powered shape mapping — template slide thumbnails sent to LLM
  • Second-pass LLM polish — refines hierarchy, density, and typography before build
  • Repair loop — failed slides get re-mapped with rendered image + issues
  • Parallel LLM calls, asset fetches, and validation checks
  • Format-preserving text fill (fonts, colors, alignment)
  • Asset integration: icons (Iconify), photos (Pexels), chart placeholders
  • Image resizing (cover/contain modes)

Usage
─────
  python slide_builder.py build_spec.json templates/Startup_pitch_deck.pptx -o output/
  python slide_builder.py build_spec.json templates/Startup_pitch_deck.pptx -o output/ --no-validate
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import json
import logging
import os
import re
import shutil
import sys
import textwrap
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from copy import deepcopy
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import httpx
from dotenv import load_dotenv
from langgraph.graph import END, START, StateGraph
from PIL import Image, ImageColor, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from typing_extensions import TypedDict

from agent_tools import ToolContext, builder_tools
from asset_manager import resolve_icon, resolve_photo
from content_planner import plan_content
from layout_selector import select_layouts
from llm_tool_harness import run_json_stage_with_tools
from template_cataloger import render_template_to_pngs

# ── Config ─────────────────────────────────────────────────────────────────────

_HERE = Path(__file__).resolve().parent
load_dotenv(_HERE / ".env")

MODEL_BASE_URL: str   = os.getenv("MODEL_BASE_URL", "https://api.fireworks.ai/inference/v1").rstrip("/")
MODEL_API_KEY: str    = os.getenv("MODEL_API_KEY", "")
MODEL_NAME: str       = os.getenv("MODEL_NAME", "accounts/fireworks/models/qwen3p6-plus")
MODEL_TIMEOUT_SEC: int = int(os.getenv("MODEL_TIMEOUT_SEC", "300"))
MODEL_HTTP_RETRIES: int = int(os.getenv("MODEL_HTTP_RETRIES", "3"))
PARALLELISM: int      = int(os.getenv("BUILDER_PARALLELISM", "2"))
MAX_REPAIR: int       = int(os.getenv("BUILDER_MAX_REPAIR_ROUNDS", "2"))
VALIDATION_RENDER_DPI: int = int(os.getenv("VALIDATION_RENDER_DPI", "220"))
ASSETS_CACHE: Path    = _HERE / os.getenv("ASSETS_CACHE_DIR", "assets/cache")

logger = logging.getLogger(__name__)

_PLACEHOLDER_TEXT_PATTERNS = [
    re.compile(r"lorem ipsum", re.I),
    re.compile(r"^solution\s+\d+$", re.I),
    re.compile(r"^presented by\s*:?", re.I),
    re.compile(r"^ketut\b", re.I),
    re.compile(r"^salford$", re.I),
    re.compile(r"^thank you$", re.I),
    re.compile(r"^pitch deck presentation$", re.I),
    re.compile(r"^size of market$", re.I),
    re.compile(r"^direct competitor$", re.I),
    re.compile(r"^indirect competitor$", re.I),
    re.compile(r"^accomplishments date$", re.I),
    re.compile(r"reallygreatsite", re.I),
    re.compile(r"anywhere st", re.I),
    re.compile(r"any city", re.I),
]

_RENDER_SAFE_FONTS = [
    "Bahnschrift",
    "Georgia",
    "Verdana",
    "Arial",
    "Calibri",
]

_FONT_FALLBACKS = {
    "open sauce heavy": "Bahnschrift",
    "open sauce light": "Bahnschrift",
    "open sauce": "Bahnschrift",
}

_GENERIC_FONT_ALIASES = {
    "sans-serif": "Bahnschrift",
    "sans serif": "Bahnschrift",
    "serif": "Georgia",
    "display": "Bahnschrift",
    "mono": "Verdana",
    "monospace": "Verdana",
}


# ── State ──────────────────────────────────────────────────────────────────────

class BuilderState(TypedDict):
    # Inputs
    description: str
    max_slides: int | None
    content_plan: dict
    build_spec: dict
    template_catalog: dict
    template_path: str
    output_dir: str
    do_validate: bool
    catalog_dir: str                   # path to catalog dir (for thumbnails)

    # Intermediate
    slide_mappings: list[dict]         # per-slide fill instructions from LLM
    asset_requests: list[dict]         # collected asset requests
    asset_paths: dict[str, str]        # asset_key → local file path
    pptx_path: str                     # built PPTX
    slide_pngs: list[str]             # rendered PNGs
    build_issues: list[dict]          # structural/runtime issues found during build

    # Validation & repair
    validation_results: list[dict]
    repair_round: int
    failed_slide_indices: list[int]    # indices into spec["slides"] that failed validation

    # Output
    errors: list[str]


def _resolve_font_family_name(name: str | None, *, default: str | None = None) -> str | None:
    raw = str(name or "").strip()
    if not raw:
        return default
    normalized = raw.lower()
    if normalized in _GENERIC_FONT_ALIASES:
        return _GENERIC_FONT_ALIASES[normalized]
    if normalized in _FONT_FALLBACKS:
        return _FONT_FALLBACKS[normalized]
    return raw


def _alignment_label(value: Any) -> str | None:
    if value is None:
        return None
    try:
        return str(value).split(".")[-1].lower()
    except Exception:
        return None


def _parse_alignment(value: Any) -> Any:
    label = str(value or "").strip().lower()
    mapping = {
        "left": PP_PARAGRAPH_ALIGNMENT.LEFT,
        "center": PP_PARAGRAPH_ALIGNMENT.CENTER,
        "right": PP_PARAGRAPH_ALIGNMENT.RIGHT,
        "justify": PP_PARAGRAPH_ALIGNMENT.JUSTIFY,
    }
    return mapping.get(label)


# ═══════════════════════════════════════════════════════════════════════════════
# PYTHON-PPTX TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

def extract_shapes(slide, slide_w_emu: int, slide_h_emu: int) -> list[dict]:
    """Extract ALL shapes from a slide with stable indices.

    Each shape gets an index (0-based) that the LLM uses to reference it.
    No filtering — the LLM decides what to do with everything.
    """
    sw = slide_w_emu / 914400
    sh = slide_h_emu / 914400
    shapes = []

    for idx, shape in enumerate(slide.shapes):
        try:
            x = shape.left   / 914400
            y = shape.top    / 914400
            w = shape.width  / 914400
            h = shape.height / 914400
        except Exception:
            continue

        info: dict[str, Any] = {
            "index":     idx,
            "name":      shape.name,
            "left_in":   round(x, 2),
            "top_in":    round(y, 2),
            "width_in":  round(w, 2),
            "height_in": round(h, 2),
        }

        # Rotation
        try:
            rot = shape.rotation
            if rot and abs(rot) > 0.5:
                info["rotation_deg"] = round(rot, 1)
        except Exception:
            pass

        # Check visibility
        visible_w = min(x + w, sw) - max(x, 0)
        visible_h = min(y + h, sh) - max(y, 0)
        if visible_w <= 0 or visible_h <= 0:
            info["on_slide"] = False

        # Picture shape
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                info["type"] = "picture"
                shapes.append(info)
                continue
        except Exception:
            pass

        # Text shape — extract text + font info
        if shape.has_text_frame:
            info["type"] = "text"
            info["text"] = shape.text_frame.text.strip()[:300]
            # Font info
            font_pt = 0
            font_name = None
            font_bold = None
            font_italic = None
            font_color = None
            alignment = None
            for para in shape.text_frame.paragraphs:
                if alignment is None:
                    alignment = _alignment_label(para.alignment)
                for run in para.runs:
                    if run.font.size:
                        font_pt = max(font_pt, run.font.size.pt)
                    if font_name is None and run.font.name:
                        font_name = run.font.name
                    if font_bold is None and run.font.bold is not None:
                        font_bold = bool(run.font.bold)
                    if font_italic is None and run.font.italic is not None:
                        font_italic = bool(run.font.italic)
                    if font_color is None:
                        try:
                            if run.font.color and run.font.color.type is not None and run.font.color.rgb:
                                font_color = f"#{run.font.color.rgb}"
                        except Exception:
                            font_color = None
            info["font_pt"] = round(font_pt, 1) if font_pt else None
            if font_name:
                info["font_name"] = font_name
                resolved = _resolve_font_family_name(font_name)
                if resolved and resolved != font_name:
                    info["render_font_name"] = resolved
            if font_bold is not None:
                info["font_bold"] = font_bold
            if font_italic is not None:
                info["font_italic"] = font_italic
            if font_color:
                info["font_color"] = font_color
            if alignment:
                info["alignment"] = alignment
        else:
            info["type"] = "shape"

        try:
            fill = shape.fill
            if fill and fill.type is not None:
                try:
                    if fill.fore_color.rgb:
                        info["fill_color"] = f"#{fill.fore_color.rgb}"
                except Exception:
                    pass
                try:
                    if fill.transparency is not None:
                        info["fill_transparency"] = round(float(fill.transparency), 2)
                except Exception:
                    pass
        except Exception:
            pass

        try:
            line = shape.line
            if line:
                try:
                    if line.fill.fore_color.rgb:
                        info["line_color"] = f"#{line.fill.fore_color.rgb}"
                except Exception:
                    pass
                try:
                    if line.width:
                        info["line_width_pt"] = round(line.width.pt, 1)
                except Exception:
                    pass
                try:
                    if line.fill.transparency is not None:
                        info["line_transparency"] = round(float(line.fill.transparency), 2)
                except Exception:
                    pass
        except Exception:
            pass

        shapes.append(info)

    return shapes


def _shapes_for_prompt(shapes: list[dict]) -> str:
    """Format shape list for the LLM. Each shape has a stable [index]."""
    lines = []
    for s in shapes:
        idx  = s["index"]
        dims = f"{s['width_in']}×{s['height_in']}in at ({s['left_in']}, {s['top_in']})"
        tags = []
        if not s.get("on_slide", True):
            tags.append("OFF-SLIDE")
        if s.get("rotation_deg"):
            tags.append(f"rotated {s['rotation_deg']}°")
        tag_str = f" [{', '.join(tags)}]" if tags else ""

        if s["type"] == "text":
            font_parts = []
            if s.get("font_name"):
                font_label = str(s["font_name"])
                if s.get("render_font_name"):
                    font_label += f" -> {s['render_font_name']}"
                font_parts.append(font_label)
            if s.get("font_pt"):
                font_parts.append(f"{s['font_pt']}pt")
            if s.get("font_bold"):
                font_parts.append("bold")
            if s.get("font_italic"):
                font_parts.append("italic")
            if s.get("alignment"):
                font_parts.append(f"align={s['alignment']}")
            if s.get("font_color"):
                font_parts.append(f"color={s['font_color']}")
            if s.get("fill_color"):
                font_parts.append(f"fill={s['fill_color']}")
            if s.get("line_color"):
                font_parts.append(f"line={s['line_color']}")
            font = f", {'; '.join(font_parts)}" if font_parts else ""
            txt  = s.get("text", "") or "(empty)"
            lines.append(f'  [{idx}] "{s["name"]}" — text{font}{tag_str} — "{txt[:120]}" — {dims}')
        elif s["type"] == "picture":
            lines.append(f'  [{idx}] "{s["name"]}" — picture{tag_str} — {dims}')
        else:
            style_parts = []
            if s.get("fill_color"):
                style_parts.append(f"fill={s['fill_color']}")
            if s.get("fill_transparency") is not None:
                style_parts.append(f"fill_t={s['fill_transparency']}")
            if s.get("line_color"):
                style_parts.append(f"line={s['line_color']}")
            if s.get("line_width_pt") is not None:
                style_parts.append(f"line_w={s['line_width_pt']}pt")
            style = f", {'; '.join(style_parts)}" if style_parts else ""
            lines.append(f'  [{idx}] "{s["name"]}" — {s["type"]}{style}{tag_str} — {dims}')
    return "\n".join(lines)


def _shape_box(shape: dict, instr: dict | None = None) -> tuple[float, float, float, float]:
    left = _coerce_optional_float((instr or {}).get("left_in"))
    top = _coerce_optional_float((instr or {}).get("top_in"))
    width = _coerce_optional_float((instr or {}).get("width_in"))
    height = _coerce_optional_float((instr or {}).get("height_in"))
    x = left if left is not None else float(shape.get("left_in", 0.0) or 0.0)
    y = top if top is not None else float(shape.get("top_in", 0.0) or 0.0)
    w = width if width is not None else float(shape.get("width_in", 0.0) or 0.0)
    h = height if height is not None else float(shape.get("height_in", 0.0) or 0.0)
    return x, y, x + max(w, 0.0), y + max(h, 0.0)


def _rect_intersection(a: tuple[float, float, float, float],
                       b: tuple[float, float, float, float]) -> tuple[float, float]:
    left = max(a[0], b[0])
    top = max(a[1], b[1])
    right = min(a[2], b[2])
    bottom = min(a[3], b[3])
    if right <= left or bottom <= top:
        return 0.0, 0.0
    area = (right - left) * (bottom - top)
    a_area = max((a[2] - a[0]) * (a[3] - a[1]), 1e-6)
    b_area = max((b[2] - b[0]) * (b[3] - b[1]), 1e-6)
    return area, area / min(a_area, b_area)


def _shape_role_label(shape: dict) -> str:
    kind = str(shape.get("type") or "shape")
    width = float(shape.get("width_in", 0.0) or 0.0)
    height = float(shape.get("height_in", 0.0) or 0.0)
    font_pt = float(shape.get("font_pt", 0.0) or 0.0)
    area = width * height
    if kind == "text":
        if font_pt >= 72:
            return "hero_text"
        if font_pt >= 28:
            return "heading_text"
        if font_pt > 0 and area <= 0.5:
            return "meta_text"
        return "body_text"
    if kind == "picture":
        return "picture_slot"
    if max(width, height) >= 2.5 and min(width, height) <= 0.08:
        return "divider_line"
    if area >= 7.5:
        return "decorative_shape"
    if 0.7 <= width <= 4.0 and 0.7 <= height <= 4.0 and abs(width - height) <= 0.6:
        return "stat_container"
    return "shape"


def _instruction_summary(instr: dict, shape: dict) -> str:
    action = str(instr.get("action") or "skip")
    if action == "fill_text":
        text = str(instr.get("text") or "").replace("\n", " / ")
        parts = [text[:90] or "(empty)"]
        if instr.get("font_family"):
            parts.append(f"font={instr['font_family']}")
        if instr.get("font_size_pt"):
            parts.append(f"{instr['font_size_pt']}pt")
        if instr.get("text_alignment"):
            parts.append(f"align={instr['text_alignment']}")
        if instr.get("fill_color"):
            parts.append(f"fill={instr['fill_color']}")
        return ", ".join(parts)
    if action in {"fill_icon", "fill_image"}:
        parts = []
        if instr.get("asset_id"):
            parts.append(f"asset={instr['asset_id']}")
        elif instr.get("search_query"):
            parts.append(str(instr.get("search_query") or "").strip())
        if instr.get("image_mode"):
            parts.append(f"mode={instr['image_mode']}")
        if instr.get("image_opacity") is not None:
            parts.append(f"opacity={instr['image_opacity']}")
        if instr.get("layer_position"):
            parts.append(f"layer={instr['layer_position']}")
        return ", ".join(part for part in parts if part)
    if action == "fill_chart":
        if instr.get("asset_id"):
            return f"asset:{instr['asset_id']}"
        return f"{instr.get('chart_type', 'chart')}:{instr.get('chart_title', 'Chart')}"
    if action == "style_shape":
        parts = []
        if instr.get("fill_color"):
            parts.append(f"fill={instr['fill_color']}")
        if instr.get("line_color"):
            parts.append(f"line={instr['line_color']}")
        if instr.get("fill_transparency") is not None:
            parts.append(f"fill_t={instr['fill_transparency']}")
        if instr.get("layer_position"):
            parts.append(f"layer={instr['layer_position']}")
        return ", ".join(parts) or _shape_role_label(shape)
    if action == "clear":
        return f"clear ({instr.get('clear_mode', 'text_only')})"
    return _shape_role_label(shape)


def _mapping_for_prompt(
    mapping: dict,
    shapes: list[dict],
    *,
    limit: int = 20,
) -> str:
    lines: list[str] = []
    for instr in mapping.get("instructions", []) or []:
        if not isinstance(instr, dict):
            continue
        action = str(instr.get("action") or "skip")
        if action == "skip":
            continue
        shape_index = _coerce_shape_index(instr, shapes)
        if shape_index is None or not (0 <= shape_index < len(shapes)):
            continue
        shape = shapes[shape_index]
        box = _shape_box(shape, instr)
        lines.append(
            f"  [{shape_index}] {action} on {_shape_role_label(shape)} "
            f"at ({box[0]:.2f}, {box[1]:.2f}, {box[2] - box[0]:.2f}, {box[3] - box[1]:.2f}) "
            f"— {_instruction_summary(instr, shape)}"
        )
        if len(lines) >= limit:
            break
    return "\n".join(lines) if lines else "  (no active instructions)"


def _looks_like_display_text(shape: dict, instr: dict) -> bool:
    requested_size = _coerce_optional_float(instr.get("font_size_pt"))
    font_pt = requested_size if requested_size is not None else float(shape.get("font_pt", 0.0) or 0.0)
    text = str(instr.get("text") or "")
    return font_pt >= 60 or (len(text.split()) <= 10 and "\n" in text and font_pt >= 40)


def _layout_signals_for_prompt(
    shapes: list[dict],
    mapping: dict,
    *,
    slide_width_in: float,
    slide_height_in: float,
    limit: int = 8,
) -> str:
    signals: list[str] = []
    dividers: list[tuple[int, tuple[float, float, float, float]]] = []
    text_items: list[tuple[int, dict, dict, tuple[float, float, float, float]]] = []
    for shape in shapes:
        if _shape_role_label(shape) == "divider_line":
            dividers.append((int(shape["index"]), _shape_box(shape)))
    for instr in mapping.get("instructions", []) or []:
        if not isinstance(instr, dict):
            continue
        shape_index = _coerce_shape_index(instr, shapes)
        if shape_index is None or not (0 <= shape_index < len(shapes)):
            continue
        shape = shapes[shape_index]
        action = str(instr.get("action") or "skip")
        if action != "fill_text":
            continue
        box = _shape_box(shape, instr)
        if box[0] < -0.02 or box[1] < -0.02 or box[2] > slide_width_in + 0.02 or box[3] > slide_height_in + 0.02:
            signals.append(f"Shape [{shape_index}] text box extends beyond the slide edge.")
        for divider_idx, divider_box in dividers:
            area, ratio = _rect_intersection(box, divider_box)
            if area > 0.01 or ratio > 0.2:
                signals.append(
                    f"Shape [{shape_index}] text appears to cross divider/line shape [{divider_idx}]."
                )
        text_items.append((shape_index, shape, instr, box))
    for idx, (left_idx, left_shape, left_instr, left_box) in enumerate(text_items):
        for right_idx, right_shape, right_instr, right_box in text_items[idx + 1:]:
            area, ratio = _rect_intersection(left_box, right_box)
            if area <= 0.02 or ratio <= 0.04:
                continue
            if _looks_like_display_text(left_shape, left_instr) and _looks_like_display_text(right_shape, right_instr):
                continue
            signals.append(
                f"Shapes [{left_idx}] and [{right_idx}] may be colliding in a non-display-text region."
            )
    deduped: list[str] = []
    seen: set[str] = set()
    for signal in signals:
        if signal not in seen:
            deduped.append(signal)
            seen.add(signal)
        if len(deduped) >= limit:
            break
    return "\n".join(f"- {item}" for item in deduped) if deduped else "- No strong structural warnings detected."


# ── Clone / delete slides ──────────────────────────────────────────────────────

def _clone_slide(prs: Presentation, src_index: int) -> Any:
    """Clone slide at src_index within the same Presentation. Appends at end."""
    source = prs.slides[src_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Clear default placeholders from the new blank slide
    dst_tree = new_slide.shapes._spTree
    for child in list(dst_tree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            dst_tree.remove(child)

    # Deep-copy all elements from source
    src_tree = source.shapes._spTree
    for child in src_tree:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            dst_tree.append(deepcopy(child))

    # Copy background
    src_cSld = source._element.find(qn("p:cSld"))
    dst_cSld = new_slide._element.find(qn("p:cSld"))
    src_bg = src_cSld.find(qn("p:bg")) if src_cSld is not None else None
    if src_bg is not None:
        dst_bg = dst_cSld.find(qn("p:bg")) if dst_cSld is not None else None
        if dst_bg is not None:
            dst_cSld.remove(dst_bg)
        dst_cSld.insert(0, deepcopy(src_bg))

    # Copy relationships (images, embedded objects)
    src_part = source.part
    dst_part = new_slide.part
    for rel in src_part.rels.values():
        if not rel.is_external:
            try:
                dst_part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                pass

    return new_slide


def _delete_slide(prs: Presentation, index: int) -> None:
    """Delete slide at index."""
    sld_id_lst = prs.slides._sldIdLst
    rId = sld_id_lst[index].get(qn("r:id"))
    prs.part.drop_rel(rId)
    del sld_id_lst[index]


def _save_presentation(prs: Presentation, target: Path) -> Path:
    """Save deck, falling back to a timestamped filename if the primary path is locked."""
    try:
        prs.save(str(target))
        return target
    except PermissionError:
        alt = target.with_name(f"{target.stem}-{datetime.now().strftime('%Y%m%d-%H%M%S')}{target.suffix}")
        logger.warning("build: %s is locked; saving to %s instead", target, alt)
        prs.save(str(alt))
        return alt


# ── Text fill (format-preserving) ─────────────────────────────────────────────

def _capture_text_format(shape) -> dict:
    """Capture font/paragraph formatting from first non-empty run."""
    fmt: dict[str, Any] = {}
    if not shape.has_text_frame:
        return fmt
    for para in shape.text_frame.paragraphs:
        if not fmt.get("alignment"):
            fmt["alignment"] = para.alignment
        for run in para.runs:
            f = run.font
            fmt["name"]  = f.name
            fmt["size"]  = f.size
            fmt["bold"]  = f.bold
            fmt["italic"] = f.italic
            try:
                fmt["color"] = f.color.rgb if f.color and f.color.type is not None else None
            except Exception:
                fmt["color"] = None
            return fmt
    return fmt


def _delete_shape(shape) -> None:
    try:
        shape._element.getparent().remove(shape._element)
    except Exception:
        return


def _has_template_bullet_style(shape) -> bool:
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        p_pr = para._p.find(qn("a:pPr"))
        if p_pr is None:
            continue
        for child in p_pr:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag.startswith("bu") and tag != "buNone":
                return True
    return False


def _strip_bullet_prefix(line: str) -> str:
    return re.sub(r"^\s*(?:[-*•·▪◦]+\s+|(?:\d+[\.\)])\s+)", "", line).strip()


def _normalize_text_lines(text: str, *, use_template_bullets: bool) -> list[str]:
    lines: list[str] = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if use_template_bullets:
            line = _strip_bullet_prefix(line)
        lines.append(line)
    return lines or [""]


def _iter_text_runs(shape):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield run


def _parse_rgb(value: str | None) -> RGBColor | None:
    if not value:
        return None
    try:
        rgb = ImageColor.getrgb(value)
    except ValueError:
        return None
    return RGBColor(*rgb[:3])


def _parse_color_or_none(value: Any) -> tuple[RGBColor | None, bool]:
    raw = str(value or "").strip()
    if not raw:
        return None, False
    if raw.lower() in {"none", "transparent", "no-fill", "no fill", "clear"}:
        return None, True
    return _parse_rgb(raw), False


def _coerce_unit_ratio(value: Any) -> float | None:
    numeric = _coerce_optional_float(value)
    if numeric is None:
        return None
    if numeric > 1.0:
        numeric = numeric / 100.0
    return max(0.0, min(1.0, numeric))


def _balanced_display_text(text: str, *, max_lines: int = 3) -> str:
    words = [w for w in (text or "").split() if w]
    if len(words) < 3:
        return text.strip()

    def _partitions(total: int, groups: int, start: int = 1) -> list[list[int]]:
        if groups == 1:
            return [[total]]
        combos: list[list[int]] = []
        for size in range(start, total - groups + 2):
            for rest in _partitions(total - size, groups - 1, 1):
                combos.append([size, *rest])
        return combos

    best_lines = [" ".join(words)]
    best_score = float("inf")
    limit = min(max_lines, len(words))
    for line_count in range(2, limit + 1):
        for groups in _partitions(len(words), line_count):
            lines = []
            cursor = 0
            for size in groups:
                lines.append(" ".join(words[cursor:cursor + size]))
                cursor += size
            lengths = [len(line) for line in lines]
            score = (max(lengths) - min(lengths)) * 2 + max(lengths)
            if score < best_score:
                best_score = score
                best_lines = lines
    return "\n".join(best_lines)


def _coerce_optional_float(value: Any) -> float | None:
    try:
        if value is None or value == "":
            return None
        return float(value)
    except (TypeError, ValueError):
        return None


def _text_overrides(instr: dict, fmt: dict) -> dict[str, Any]:
    base_size = fmt.get("size").pt if getattr(fmt.get("size"), "pt", None) else None
    font_size = instr.get("font_size_pt")
    try:
        font_size = float(font_size) if font_size is not None else None
    except (TypeError, ValueError):
        font_size = None

    font_scale = instr.get("font_scale")
    try:
        font_scale = float(font_scale) if font_scale is not None else None
    except (TypeError, ValueError):
        font_scale = None

    max_size_pt = font_size or base_size or 24.0
    if font_scale:
        max_size_pt *= font_scale

    default_font = _resolve_font_family_name(fmt.get("name"))
    if not default_font:
        default_font = "Georgia" if base_size and base_size >= 42 else "Bahnschrift"

    return {
        "font_family": _resolve_font_family_name(instr.get("font_family"), default=default_font) or default_font,
        "max_size_pt": max(10.0, min(320.0, max_size_pt)),
        "font_size_pt": font_size,
        "bold": instr.get("bold") if instr.get("bold") is not None else fmt.get("bold"),
        "italic": instr.get("italic") if instr.get("italic") is not None else fmt.get("italic"),
        "text_color": _parse_rgb(instr.get("text_color")) or fmt.get("color"),
        "uppercase": bool(instr.get("uppercase")),
        "disable_fit": bool(instr.get("disable_fit")),
        "alignment": _parse_alignment(instr.get("text_alignment")) or fmt.get("alignment"),
    }


def _style_overrides(instr: dict) -> dict[str, Any]:
    fill_color, clear_fill = _parse_color_or_none(instr.get("fill_color"))
    line_color, clear_line = _parse_color_or_none(instr.get("line_color"))
    line_width_pt = _coerce_optional_float(instr.get("line_width_pt"))
    return {
        "fill_color": fill_color,
        "clear_fill": clear_fill,
        "fill_transparency": _coerce_unit_ratio(instr.get("fill_transparency")),
        "line_color": line_color,
        "clear_line": clear_line,
        "line_width_pt": max(0.0, min(24.0, line_width_pt)) if line_width_pt is not None else None,
        "line_transparency": _coerce_unit_ratio(instr.get("line_transparency")),
        "layer_position": str(instr.get("layer_position") or "").strip().lower() or None,
    }


def _fit_text_to_shape(shape, overrides: dict[str, Any]) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    font_family = overrides.get("font_family") or "Calibri"
    max_size_pt = int(round(overrides.get("max_size_pt") or 24))
    bold = bool(overrides.get("bold"))
    italic = bool(overrides.get("italic"))

    for max_size in range(max_size_pt, 9, -2):
        try:
            tf.fit_text(
                font_family=font_family,
                max_size=max_size,
                bold=bold,
                italic=italic,
            )
            return
        except Exception:
            continue
    try:
        tf.fit_text(font_family=font_family or "Bahnschrift", max_size=min(max_size_pt, 24), bold=bold, italic=italic)
    except Exception:
        return


def _apply_text_color(shape, color: RGBColor | None) -> None:
    if color is None:
        return
    for run in _iter_text_runs(shape):
        try:
            run.font.color.rgb = color
        except Exception:
            continue


def _apply_shape_fill(shape, color: RGBColor | None, *, clear_fill: bool, transparency: float | None) -> None:
    try:
        fill = shape.fill
    except Exception:
        return
    try:
        if clear_fill:
            fill.background()
            return
        if color is not None:
            fill.solid()
            fill.fore_color.rgb = color
        if transparency is not None:
            fill.transparency = transparency
    except Exception:
        return


def _apply_shape_line(shape, color: RGBColor | None, *, clear_line: bool,
                      width_pt: float | None, transparency: float | None) -> None:
    try:
        line = shape.line
    except Exception:
        return
    try:
        if clear_line:
            line.fill.background()
            return
        if color is not None:
            line.fill.solid()
            line.fill.fore_color.rgb = color
        if width_pt is not None:
            line.width = Pt(width_pt)
        if transparency is not None:
            line.fill.transparency = transparency
    except Exception:
        return


def _move_shape_to_layer(slide, shape, layer_position: str | None) -> None:
    if layer_position not in {"front", "back"}:
        return
    try:
        sp_tree = slide.shapes._spTree
        element = shape._element
        sp_tree.remove(element)
        if layer_position == "front":
            sp_tree.append(element)
            return
        insert_idx = 2 if len(sp_tree) >= 2 else len(sp_tree)
        sp_tree.insert(insert_idx, element)
    except Exception:
        return


def _apply_shape_style(slide, shape, instr: dict) -> None:
    style = _style_overrides(instr)
    _apply_shape_fill(
        shape,
        style.get("fill_color"),
        clear_fill=bool(style.get("clear_fill")),
        transparency=style.get("fill_transparency"),
    )
    _apply_shape_line(
        shape,
        style.get("line_color"),
        clear_line=bool(style.get("clear_line")),
        width_pt=style.get("line_width_pt"),
        transparency=style.get("line_transparency"),
    )
    _move_shape_to_layer(slide, shape, style.get("layer_position"))


def _set_shape_text(shape, text: str, fmt: dict | None = None, *, instr: dict | None = None) -> None:
    """Replace ALL text in shape, applying captured formatting.

    Clears every content element (runs, fields, breaks, etc.) —
    not just <a:r> — so no ghost text remains.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Nuke ALL existing paragraphs except the first, then clear the first
    txBody = tf._txBody
    p_elements = txBody.findall(qn("a:p"))
    for p_el in p_elements[1:]:
        txBody.remove(p_el)
    first_p = p_elements[0]
    # Remove ALL children except pPr (paragraph properties)
    for child in list(first_p):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag != "pPr":
            first_p.remove(child)

    if fmt is None:
        fmt = {}
    if instr is None:
        instr = {}

    overrides = _text_overrides(instr, fmt)
    use_template_bullets = _has_template_bullet_style(shape)

    if overrides.get("uppercase"):
        text = text.upper()

    lines = _normalize_text_lines(text, use_template_bullets=use_template_bullets)
    if instr.get("single_paragraph"):
        para = tf.paragraphs[0]
        if overrides.get("alignment") is not None:
            para.alignment = overrides["alignment"]
        elif fmt.get("alignment") is not None:
            para.alignment = fmt["alignment"]

        run = para.add_run()
        run.text = "\n".join(lines)
        if overrides.get("font_family"):
            run.font.name = overrides["font_family"]
        if overrides.get("font_size_pt"):
            run.font.size = Pt(float(overrides["font_size_pt"]))
        elif fmt.get("size"):
            run.font.size = fmt["size"]
        if overrides.get("bold") is not None:
            run.font.bold = bool(overrides["bold"])
        if overrides.get("italic") is not None:
            run.font.italic = bool(overrides["italic"])
        if overrides.get("text_color"):
            run.font.color.rgb = overrides["text_color"]
    else:
        for i, line in enumerate(lines):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if overrides.get("alignment") is not None:
                para.alignment = overrides["alignment"]
            elif fmt.get("alignment") is not None:
                para.alignment = fmt["alignment"]

            run = para.add_run()
            run.text = line
            if overrides.get("font_family"):
                run.font.name = overrides["font_family"]
            if overrides.get("font_size_pt"):
                run.font.size = Pt(float(overrides["font_size_pt"]))
            elif fmt.get("size"):
                run.font.size = fmt["size"]
            if overrides.get("bold") is not None:
                run.font.bold = bool(overrides["bold"])
            if overrides.get("italic") is not None:
                run.font.italic = bool(overrides["italic"])
            if overrides.get("text_color"):
                run.font.color.rgb = overrides["text_color"]

    if not overrides.get("disable_fit"):
        _fit_text_to_shape(shape, overrides)
    _apply_text_color(shape, overrides.get("text_color"))


def _looks_like_placeholder_text(text: str) -> bool:
    normalized = " ".join((text or "").split()).strip()
    if not normalized:
        return False
    return any(pattern.search(normalized) for pattern in _PLACEHOLDER_TEXT_PATTERNS)


def _shape_font_pt(shape) -> float:
    if not getattr(shape, "has_text_frame", False):
        return 0.0
    largest = 0.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                largest = max(largest, run.font.size.pt)
    return largest


def _shape_bounds_in(shape) -> tuple[float, float, float, float]:
    return (
        shape.left / 914400,
        shape.top / 914400,
        shape.width / 914400,
        shape.height / 914400,
    )


def _primary_body_copy(slide_info: dict) -> str:
    for block in slide_info.get("full_content", []):
        if block.get("type") == "text" and block.get("body"):
            return str(block["body"]).strip()
        if block.get("type") == "stat" and block.get("context"):
            return str(block["context"]).strip()
    return str(slide_info.get("deck_theme", "")).strip()


def _bullet_items(slide_info: dict) -> list[str]:
    items: list[str] = []
    for block in slide_info.get("full_content", []):
        if block.get("type") == "bullets":
            items.extend(str(item).strip() for item in block.get("items", []) if str(item).strip())
    return items


def _shorten_copy(text: str, width: int) -> str:
    clean = " ".join((text or "").split())
    if not clean:
        return ""
    return textwrap.shorten(clean, width=width, placeholder="…")


def _compress_caption(text: str, width: int = 48) -> str:
    text = re.sub(r"\([^)]*\)", "", text or "")
    replacements = {
        "How the world's third-largest automotive market is ": "",
        "Tracing India's transformation from ": "",
        "environmental aspiration": "green ambition",
        "self-reliant future": "energy independence",
        "rewriting the rules": "resetting mobility",
        "through policy momentum, market acceleration, and infrastructure challenges": "",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    text = re.sub(r"\s+", " ", text).strip(" ,.-")
    return textwrap.shorten(text, width=width, placeholder="…")


def _compress_roadmap_item(text: str, width: int = 28) -> str:
    text = re.sub(r"\([^)]*\)", "", text or "")
    replacements = {
        "Localize ": "",
        "annual ": "",
        "battery production": "battery output",
        "Achieve ": "",
        "EV share": "EV share",
        "in urban ": "in urban ",
        "Establish ": "",
        "second-life ": "2nd-life ",
        "ecosystems ": "",
        "for grid storage": "for grid storage",
        "Export hub: ": "",
        "Africa and Southeast Asia": "Africa & SE Asia",
        "public transport": "transit",
        "commercial fleets": "fleets",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    text = re.sub(r"\s+", " ", text).strip(" ,.-")
    return textwrap.shorten(text, width=width, placeholder="…")


def _hero_title_text(title: str) -> str:
    words = (title or "").split()
    if len(words) <= 2:
        return " ".join(words)
    if len(words) == 3:
        return "\n".join([
            " ".join(words[:2]),
            words[2],
        ])
    midpoint = max(1, len(words) // 2)
    return "\n".join([
        " ".join(words[:midpoint]),
        " ".join(words[midpoint:]),
    ])


def _apply_template_guardrails(
    slide_info: dict,
    shape_list: list[Any],
    original_text_by_index: dict[int, str],
    touched_shape_indices: set[int],
) -> None:
    archetype = str(slide_info.get("template_archetype") or "").strip().lower()
    if archetype not in {"cover", "closing"}:
        return

    def _textbox_like(idx: int) -> bool:
        name = getattr(shape_list[idx], "name", "")
        return name.startswith("TextBox") or bool(original_text_by_index.get(idx, "").strip())

    text_candidates: list[tuple[int, float, float, float, float, float]] = []
    for idx, shape in enumerate(shape_list):
        if not getattr(shape, "has_text_frame", False):
            continue
        left, top, width, height = _shape_bounds_in(shape)
        text_candidates.append((idx, _shape_font_pt(shape), width * height, left, top, height))

    if not text_candidates:
        return

    text_candidates.sort(key=lambda item: (item[1], item[2]), reverse=True)
    hero_idx = text_candidates[0][0]
    hero_shape = shape_list[hero_idx]
    hero_fmt = _capture_text_format(hero_shape)
    raw_title = str(slide_info.get("title") or "").strip()
    hero_title = raw_title if archetype == "cover" else _hero_title_text(raw_title)
    hero_font_size = 76.0 if archetype == "cover" else (112.0 if hero_title.count("\n") <= 1 else 92.0)
    _set_shape_text(
        hero_shape,
        hero_title,
        hero_fmt,
        instr={
            "font_size_pt": hero_font_size if archetype == "cover" else hero_font_size - 8,
            "bold": True,
            "disable_fit": True,
            "single_paragraph": True,
            "font_family": "Georgia",
            "text_color": "#3A3A3A",
        },
    )
    touched_shape_indices.add(hero_idx)

    subtitle_idx: int | None = None
    if archetype == "cover":
        body_copy = _compress_caption(_primary_body_copy(slide_info), width=62)
        subtitle_candidates = [
            item for item in text_candidates[1:]
            if _textbox_like(item[0]) and item[3] >= 14.0 and item[4] >= 7.5 and item[2] >= 2.5
        ]
        if body_copy and subtitle_candidates:
            subtitle_idx = subtitle_candidates[0][0]
            subtitle_shape = shape_list[subtitle_idx]
            subtitle_fmt = _capture_text_format(subtitle_shape)
            _set_shape_text(
                subtitle_shape,
                body_copy,
                subtitle_fmt,
                instr={
                    "font_family": "Aptos",
                    "font_size_pt": 17,
                    "italic": False,
                    "disable_fit": False,
                    "text_color": "#4B4F56",
                },
            )
            touched_shape_indices.add(subtitle_idx)

        for idx, shape in enumerate(shape_list):
            if idx in touched_shape_indices:
                continue
            if getattr(shape, "has_text_frame", False):
                if _looks_like_placeholder_text(original_text_by_index.get(idx, "")):
                    shape.text_frame.text = ""
            else:
                left, top, width, height = _shape_bounds_in(shape)
                if left >= 13.0 and top >= 8.5 and width * height <= 0.5:
                    _delete_shape(shape)

    if archetype == "closing":
        for item in text_candidates[1:]:
            idx = item[0]
            if idx == hero_idx:
                continue
            if _textbox_like(idx) and 7.0 <= item[3] <= 12.5 and item[4] >= 8.0 and item[2] >= 2.5:
                shape_list[idx].text_frame.text = ""
                touched_shape_indices.add(idx)

    if archetype == "closing":
        bullet_items = _bullet_items(slide_info)
        bullet_shapes = [
            item for item in text_candidates
            if item[0] not in {hero_idx, subtitle_idx}
            and _textbox_like(item[0])
            and item[3] >= 13.5
            and item[4] >= 6.0
            and item[5] <= 0.6
        ]
        bullet_shapes.sort(key=lambda item: item[4])
        for item, shape_meta in zip(bullet_items[:len(bullet_shapes)], bullet_shapes):
            idx = shape_meta[0]
            shape = shape_list[idx]
            fmt = _capture_text_format(shape)
            _set_shape_text(
                shape,
                _compress_roadmap_item(item, 24),
                fmt,
                instr={
                    "font_family": "Aptos",
                    "font_size_pt": 11,
                    "text_color": "#30343A",
                },
            )
            touched_shape_indices.add(idx)


# ── Image / icon embedding ─────────────────────────────────────────────────────

def _replace_shape_with_image(slide, shape, image_path: Path):
    """Replace a shape with an image at the same position / size and return the new picture."""
    picture = slide.shapes.add_picture(
        str(image_path), shape.left, shape.top, shape.width, shape.height,
    )
    shape._element.getparent().remove(shape._element)
    return picture


def _apply_shape_transform(shape, instr: dict, *, slide_width_in: float, slide_height_in: float) -> None:
    left_in = _coerce_optional_float(instr.get("left_in"))
    top_in = _coerce_optional_float(instr.get("top_in"))
    width_in = _coerce_optional_float(instr.get("width_in"))
    height_in = _coerce_optional_float(instr.get("height_in"))

    if left_in is not None:
        shape.left = Inches(max(-1.0, min(slide_width_in + 1.0, left_in)))
    if top_in is not None:
        shape.top = Inches(max(-1.0, min(slide_height_in + 1.0, top_in)))
    if width_in is not None:
        shape.width = Inches(max(0.2, min(slide_width_in + 2.0, width_in)))
    if height_in is not None:
        shape.height = Inches(max(0.1, min(slide_height_in + 2.0, height_in)))


def _set_speaker_notes(slide, text: str) -> None:
    """Set speaker notes on a slide."""
    if not text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ── Image utilities ────────────────────────────────────────────────────────────

def resize_image_to_fit(
    image_path: Path,
    width_px: int,
    height_px: int,
    mode: str = "cover",
    *,
    focus_x: float | None = None,
    focus_y: float | None = None,
    opacity: float | None = None,
) -> Path:
    """Resize image to exact dimensions.

    mode='cover'  : crop center to fill (no letterboxing)
    mode='contain' : fit within bounds, may have transparency around edges
    """
    img = Image.open(image_path).convert("RGBA")
    fx = _coerce_unit_ratio(focus_x)
    fy = _coerce_unit_ratio(focus_y)
    alpha_scale = _coerce_unit_ratio(opacity)
    if alpha_scale is None:
        alpha_scale = 1.0
    focus_key = f"{fx:.2f}_{fy:.2f}" if fx is not None or fy is not None else "center"
    out_path = image_path.parent / f"{image_path.stem}_{mode}_{width_px}x{height_px}_{focus_key}_{alpha_scale:.2f}.png"

    if mode == "cover":
        img_r = img.width / img.height
        tgt_r = width_px / max(height_px, 1)
        if img_r > tgt_r:
            new_h = img.height
            new_w = int(new_h * tgt_r)
            crop_left_max = max(img.width - new_w, 0)
            left = int(round((fx if fx is not None else 0.5) * crop_left_max))
            img = img.crop((left, 0, left + new_w, new_h))
        else:
            new_w = img.width
            new_h = int(new_w / tgt_r)
            crop_top_max = max(img.height - new_h, 0)
            top = int(round((fy if fy is not None else 0.5) * crop_top_max))
            img = img.crop((0, top, new_w, top + new_h))
        img = img.resize((width_px, height_px), Image.LANCZOS)
    else:
        contained = img.copy()
        contained.thumbnail((width_px, height_px), Image.LANCZOS)
        canvas = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
        left_span = max(width_px - contained.width, 0)
        top_span = max(height_px - contained.height, 0)
        left = int(round((fx if fx is not None else 0.5) * left_span))
        top = int(round((fy if fy is not None else 0.5) * top_span))
        canvas.paste(contained, (left, top), contained)
        img = canvas

    if alpha_scale < 1.0:
        alpha = img.getchannel("A")
        alpha = alpha.point(lambda px: int(px * alpha_scale))
        img.putalpha(alpha)

    img.save(str(out_path), "PNG")
    return out_path


def _load_chart_font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "arialbd.ttf" if bold else "arial.ttf",
        "timesbd.ttf" if bold else "times.ttf",
        "calibrib.ttf" if bold else "calibri.ttf",
    ]
    for name in candidates:
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _as_number(value: Any) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _extract_series_points(series: dict) -> list[tuple[str, float]]:
    points: list[tuple[str, float]] = []
    for item in series.get("data", []):
        if isinstance(item, (list, tuple)) and len(item) >= 2:
            label = str(item[0])
            value = _as_number(item[1])
        elif isinstance(item, dict):
            label = str(item.get("label") or item.get("x") or "")
            value = _as_number(item.get("value") or item.get("y"))
        else:
            continue
        if label and value is not None:
            points.append((label, value))
    return points


def _find_chart_block(slide_info: dict, instr: dict) -> dict | None:
    blocks = slide_info.get("full_content", []) or []
    idx = instr.get("content_block_index")
    if isinstance(idx, int) and 0 <= idx < len(blocks) and blocks[idx].get("type") == "chart":
        return blocks[idx]

    chart_blocks = [block for block in blocks if block.get("type") == "chart"]
    if not chart_blocks:
        return None

    target_title = (instr.get("chart_title") or "").strip().lower()
    for block in chart_blocks:
        title = str(block.get("title") or "").strip().lower()
        if target_title and title and (title == target_title or target_title in title):
            return block

    target_type = (instr.get("chart_type") or "").strip().lower()
    for block in chart_blocks:
        if target_type and str(block.get("chart_type") or "").strip().lower() == target_type:
            return block

    return chart_blocks[0]


def generate_chart_image(
    chart_block: dict,
    *,
    fallback_title: str,
    fallback_type: str,
    width_px: int = 800,
    height_px: int = 500,
) -> Path:
    payload = {
        "chart_block": chart_block,
        "fallback_title": fallback_title,
        "fallback_type": fallback_type,
        "width_px": width_px,
        "height_px": height_px,
    }
    cache = ASSETS_CACHE / "charts"
    cache.mkdir(parents=True, exist_ok=True)
    digest = hashlib.md5(json.dumps(payload, sort_keys=True, ensure_ascii=False).encode("utf-8")).hexdigest()[:12]
    out = cache / f"chart_{digest}.png"
    if out.exists():
        return out

    title = chart_block.get("title") or fallback_title or "Chart"
    chart_type = str(chart_block.get("chart_type") or fallback_type or "bar").lower()
    series = chart_block.get("series") or []
    parsed_series = []
    for idx, item in enumerate(series):
        points = _extract_series_points(item)
        if points:
            parsed_series.append({
                "name": item.get("name") or f"Series {idx + 1}",
                "points": points,
            })

    if not parsed_series:
        return generate_chart_placeholder(title, chart_type, width_px=width_px, height_px=height_px)

    img = Image.new("RGB", (width_px, height_px), (248, 249, 252))
    draw = ImageDraw.Draw(img)
    panel = [18, 18, width_px - 18, height_px - 18]
    draw.rounded_rectangle(panel, radius=24, fill=(250, 251, 255), outline=(205, 211, 222), width=2)

    font_title = _load_chart_font(24, bold=True)
    font_axis = _load_chart_font(14)
    font_tick = _load_chart_font(12)
    font_legend = _load_chart_font(12, bold=True)

    draw.text((width_px // 2, 34), str(title), fill=(47, 52, 64), font=font_title, anchor="ma")

    plot_left = 86
    plot_top = 78
    plot_right = width_px - 44
    plot_bottom = height_px - 82
    plot_w = plot_right - plot_left
    plot_h = plot_bottom - plot_top

    palette = [
        (64, 99, 222),
        (33, 157, 143),
        (230, 130, 55),
        (196, 78, 140),
        (108, 92, 231),
    ]

    all_points = [point for item in parsed_series for point in item["points"]]
    y_values = [value for _, value in all_points]
    y_max = max(y_values) if y_values else 1.0
    y_max = max(1.0, y_max * 1.1)

    for step in range(6):
        y = plot_top + int(plot_h * step / 5)
        value = y_max * (1 - step / 5)
        draw.line([(plot_left, y), (plot_right, y)], fill=(229, 233, 240), width=1)
        draw.text((plot_left - 12, y), f"{value:.0f}", fill=(116, 124, 141), font=font_tick, anchor="rm")

    draw.line([(plot_left, plot_top), (plot_left, plot_bottom)], fill=(160, 168, 182), width=2)
    draw.line([(plot_left, plot_bottom), (plot_right, plot_bottom)], fill=(160, 168, 182), width=2)

    x_labels = [label for label, _ in parsed_series[0]["points"]]
    step_count = max(1, len(x_labels) - 1)
    x_positions = {
        label: plot_left + int(plot_w * idx / step_count)
        for idx, label in enumerate(x_labels)
    }

    if chart_type == "pie":
        first = parsed_series[0]
        values = [value for _, value in first["points"]]
        labels = [label for label, _ in first["points"]]
        total = sum(values) or 1.0
        radius = min(plot_w, plot_h) // 2 - 10
        center = (plot_left + plot_w // 2, plot_top + plot_h // 2)
        start = -90.0
        for idx, (label, value) in enumerate(zip(labels, values)):
            span = 360.0 * value / total
            color = palette[idx % len(palette)]
            draw.pieslice(
                [center[0] - radius, center[1] - radius, center[0] + radius, center[1] + radius],
                start=start,
                end=start + span,
                fill=color,
                outline=(255, 255, 255),
                width=2,
            )
            start += span
    elif chart_type in {"bar", "area"}:
        group_width = max(18, plot_w // max(1, len(x_labels) * 2))
        inner_width = max(10, int(group_width / max(1, len(parsed_series))))
        for series_idx, item in enumerate(parsed_series):
            color = palette[series_idx % len(palette)]
            for point_idx, (label, value) in enumerate(item["points"]):
                center_x = x_positions.get(label, plot_left + point_idx * group_width)
                left = center_x - group_width // 2 + series_idx * inner_width
                right = left + inner_width - 2
                top = plot_bottom - int((value / y_max) * plot_h)
                draw.rounded_rectangle([left, top, right, plot_bottom], radius=5, fill=color)
    else:
        for series_idx, item in enumerate(parsed_series):
            color = palette[series_idx % len(palette)]
            pts = []
            for label, value in item["points"]:
                x = x_positions.get(label, plot_left)
                y = plot_bottom - int((value / y_max) * plot_h)
                pts.append((x, y))
            if len(pts) >= 2:
                draw.line(pts, fill=color, width=4, joint="curve")
            for x, y in pts:
                draw.ellipse([x - 5, y - 5, x + 5, y + 5], fill=(255, 255, 255), outline=color, width=3)

    for label, x in x_positions.items():
        draw.text((x, plot_bottom + 18), str(label), fill=(116, 124, 141), font=font_tick, anchor="ma")

    legend_x = plot_right
    legend_y = 56
    for idx, item in enumerate(parsed_series):
        color = palette[idx % len(palette)]
        top = legend_y + idx * 18
        draw.rounded_rectangle([legend_x - 140, top - 7, legend_x - 126, top + 7], radius=4, fill=color)
        draw.text((legend_x - 118, top), str(item["name"]), fill=(78, 86, 101), font=font_legend, anchor="lm")

    if chart_block.get("x_label"):
        draw.text((plot_left + plot_w // 2, height_px - 38), str(chart_block["x_label"]),
                  fill=(95, 103, 119), font=font_axis, anchor="ma")
    if chart_block.get("y_label"):
        draw.text((28, plot_top - 6), str(chart_block["y_label"]),
                  fill=(95, 103, 119), font=font_axis, anchor="la")

    img.save(str(out), "PNG")
    return out


def generate_chart_placeholder(
    chart_title: str,
    chart_type: str,
    width_px: int = 800,
    height_px: int = 500,
) -> Path:
    """Create a simple placeholder chart image with Pillow."""
    img  = Image.new("RGB", (width_px, height_px), (248, 248, 252))
    draw = ImageDraw.Draw(img)

    try:
        font_title = ImageFont.truetype("arialbd.ttf", 20)
        font_label = ImageFont.truetype("arial.ttf", 14)
    except OSError:
        font_title = ImageFont.load_default()
        font_label = font_title

    draw.text((width_px // 2, 25), chart_title, fill=(50, 50, 50),
              font=font_title, anchor="mt")

    # Sketch bars / lines based on chart_type
    base_y = height_px - 60
    left_m = 60
    right_m = width_px - 40

    if chart_type in ("bar", "area"):
        bars = [0.55, 0.78, 0.42, 0.91, 0.64, 0.83]
        bw = (right_m - left_m) // (len(bars) * 2)
        for i, h in enumerate(bars):
            x = left_m + i * bw * 2
            bar_h = int(h * (base_y - 70))
            draw.rectangle([x, base_y - bar_h, x + bw, base_y],
                           fill=(90, 120, 200), outline=(70, 100, 180))
    elif chart_type == "line":
        pts = [(left_m + i * (right_m - left_m) // 5,
                base_y - int(h * (base_y - 70)))
               for i, h in enumerate([0.3, 0.55, 0.48, 0.72, 0.65, 0.88])]
        draw.line(pts, fill=(90, 120, 200), width=3)
    elif chart_type == "pie":
        cx, cy, r = width_px // 2, height_px // 2 + 10, min(width_px, height_px) // 3
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=(90, 120, 200),
                     outline=(70, 100, 180))
        draw.pieslice([cx - r, cy - r, cx + r, cy + r], 0, 130,
                      fill=(200, 120, 90), outline=(180, 100, 70))
        draw.pieslice([cx - r, cy - r, cx + r, cy + r], 130, 210,
                      fill=(120, 200, 140), outline=(100, 180, 120))

    draw.line([(left_m, base_y), (right_m, base_y)], fill=(180, 180, 180), width=1)
    draw.text((width_px // 2, height_px - 20),
              "[chart placeholder — code execution pending]",
              fill=(170, 170, 170), font=font_label, anchor="mb")

    cache = ASSETS_CACHE / "charts"
    cache.mkdir(parents=True, exist_ok=True)
    h = hashlib.md5(f"{chart_title}:{chart_type}".encode()).hexdigest()[:10]
    out = cache / f"chart_{h}.png"
    img.save(str(out), "PNG")
    return out


# ═══════════════════════════════════════════════════════════════════════════════
# LLM HELPERS — non-streaming, with rate-limit backoff
# ═══════════════════════════════════════════════════════════════════════════════

# Global rate-limit lock — serializes calls when rate limited
import threading
_rate_limit_lock = threading.Lock()
_rate_limit_until = 0.0  # time.time() when we can call again


def _call_llm(messages: list[dict], *, temp: float = 0.25) -> str:
    """Call the LLM (non-streaming) and return the content string.

    Handles rate limiting with backoff. Non-streaming avoids the
    GeneratorExit issues and lets us read response headers properly.
    """
    global _rate_limit_until
    payload = {
        "model":       MODEL_NAME,
        "messages":    messages,
        "temperature": temp,
        "stream":      False,
    }
    headers = {
        "Authorization": f"Bearer {MODEL_API_KEY}",
        "Content-Type":  "application/json",
    }
    url = f"{MODEL_BASE_URL}/chat/completions"

    # Wait if rate limited
    now = time.time()
    if now < _rate_limit_until:
        wait = _rate_limit_until - now
        logger.info("  rate-limit: waiting %.1fs …", wait)
        time.sleep(wait)

    last_exc: Exception | None = None
    for attempt in range(MODEL_HTTP_RETRIES):
        try:
            with httpx.Client(timeout=MODEL_TIMEOUT_SEC) as client:
                resp = client.post(url, json=payload, headers=headers)

            # Check rate limit headers
            remaining_gen = int(resp.headers.get("x-ratelimit-remaining-tokens-generated", "99999"))
            over_limit = resp.headers.get("x-ratelimit-over-limit", "no")
            if over_limit == "yes" or remaining_gen <= 0:
                _rate_limit_until = time.time() + 15
                logger.warning("  rate-limited (remaining_gen=%d) — backing off 15s", remaining_gen)

            if resp.status_code == 429:
                _rate_limit_until = time.time() + 30
                raise ValueError("Rate limited (429)")

            if resp.status_code >= 500:
                raise ValueError(f"LLM API {resp.status_code}: {resp.text[:300]}")

            if resp.status_code >= 400:
                raise ValueError(f"LLM API {resp.status_code}: {resp.text[:300]}")

            data = resp.json()
            msg = (data.get("choices") or [{}])[0].get("message", {})
            content = msg.get("content", "") or ""
            reasoning = msg.get("reasoning_content", "") or ""

            if not content and not reasoning:
                logger.warning("  LLM returned empty content — possible rate limit")
                raise ValueError("Empty LLM response")

            if not content and reasoning:
                logger.debug("  content empty, using reasoning_content")
                content = reasoning

            logger.debug("  LLM response length: content=%d, reasoning=%d",
                         len(content), len(reasoning))
            return content.strip()
        except (httpx.TimeoutException, httpx.TransportError, ValueError) as exc:
            last_exc = exc
            retryable = isinstance(exc, (httpx.TimeoutException, httpx.TransportError))
            if isinstance(exc, ValueError):
                retryable = "Rate limited" in str(exc) or "LLM API 5" in str(exc) or "Empty LLM response" in str(exc)
            if attempt >= MODEL_HTTP_RETRIES - 1 or not retryable:
                raise
            wait = min(20, 5 * (attempt + 1))
            logger.warning("  LLM call attempt %d/%d failed: %s — retrying in %ss",
                           attempt + 1, MODEL_HTTP_RETRIES, exc, wait)
            time.sleep(wait)

    if last_exc is not None:
        raise last_exc
    raise RuntimeError("LLM call failed without raising a specific exception")


def _parse_json(raw: str) -> dict:
    """Extract JSON from LLM response — handles think tags, fences, prose.

    Kimi k2p5 often returns: <think>reasoning...</think> followed by JSON.
    May also include markdown fences, trailing commas, or multiple objects.
    """
    import re

    if not raw:
        raise ValueError("Empty response")

    # Strip <think>...</think> blocks (Kimi reasoning)
    raw = re.sub(r"<think>.*?</think>", "", raw, flags=re.DOTALL).strip()

    # Strip markdown fences
    if "```" in raw:
        raw = re.sub(r"```(?:json)?\s*", "", raw).strip()

    # Find the JSON object that contains "instructions"
    # Try to find all top-level { ... } blocks
    depth = 0
    start = None
    candidates = []
    for i, ch in enumerate(raw):
        if ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start is not None:
                candidates.append(raw[start:i + 1])
                start = None

    # Try candidates that look like our expected format first
    for candidate in candidates:
        if "instructions" in candidate:
            try:
                return json.loads(candidate)
            except json.JSONDecodeError:
                # Try fixing trailing commas
                fixed = re.sub(r",\s*([}\]])", r"\1", candidate)
                try:
                    return json.loads(fixed)
                except json.JSONDecodeError:
                    continue

    # Fall back: try any candidate
    for candidate in candidates:
        try:
            return json.loads(candidate)
        except json.JSONDecodeError:
            fixed = re.sub(r",\s*([}\]])", r"\1", candidate)
            try:
                return json.loads(fixed)
            except json.JSONDecodeError:
                continue

    # Last resort: original approach
    if "{" in raw:
        blob = raw[raw.index("{"):raw.rindex("}") + 1]
        blob = re.sub(r",\s*([}\]])", r"\1", blob)
        return json.loads(blob)

    raise ValueError(f"No JSON found in response (len={len(raw)}, start={raw[:100]!r})")


def _run_builder_stage_with_tools(
    messages: list[dict[str, Any]],
    *,
    output_dir: Path,
    temperature: float,
    final_hint: str,
) -> tuple[dict[str, Any], dict[str, str]]:
    tools = builder_tools()
    if not tools:
        return _parse_json(_call_llm(messages, temp=temperature)), {}
    return run_json_stage_with_tools(
        messages,
        call_llm=lambda msgs, temp: _call_llm(msgs, temp=temp),
        parse_json=_parse_json,
        is_final_result=lambda payload: isinstance(payload.get("instructions"), list),
        tools=tools,
        tool_context=ToolContext(output_dir=output_dir),
        final_hint=final_hint,
        temperature=temperature,
        max_tool_rounds=3,
        allow_asset_ids=True,
    )


# ═══════════════════════════════════════════════════════════════════════════════
# SHAPE MAPPING — LLM decides what goes where per slide
# ═══════════════════════════════════════════════════════════════════════════════

_MAP_SYSTEM = """\
You are a world-class presentation designer. You look at a template slide (image)
and its extracted shapes, then decide exactly how to fill each shape with the
given content to produce a polished, professional slide.

IMPORTANT: Each shape is identified by its [index] number. Use "shape_index" (integer)
in your response — NOT shape names. This ensures exact matching.

Actions for each shape:

  fill_text     — Replace text. Provide "text". Use \\n for line breaks, "• " for bullets.
                  Optional controls: "font_family", "font_size_pt", "font_scale",
                  "text_color", "bold", "italic", "uppercase", "single_paragraph",
                  "disable_fit", "style_source_index", "text_alignment",
                  plus shape styling like "fill_color", "line_color", "fill_transparency".
  fill_image    — Replace with a stock photo or generated asset. Provide "search_query" or "asset_id".
                  Optional: "image_mode" = "cover" or "contain", "image_focus_x",
                  "image_focus_y", "image_opacity", "layer_position".
  fill_icon     — Replace with an icon. Provide "search_query" + "icon_color" (hex).
                  Icons render with transparent background, so choose colors with contrast.
  fill_chart    — Replace with a real chart. Provide "chart_title" + "chart_type" (bar/line/pie),
                  or use "asset_id" for a generated chart/diagram.
                  Optional: "content_block_index" if a specific chart block should drive the render.
  style_shape   — Keep the existing shape but restyle it. Use for theme, panels, dividers,
                  glass cards, dark/light blocks, and decorative shapes.
  clear         — Remove template content. Optional: "clear_mode" = "text_only" or "delete_shape".
  skip          — Leave untouched.

Optional on ANY instruction:
  "left_in", "top_in", "width_in", "height_in" — move or resize the EXISTING template shape in inches.
  Use this to align scattered placeholders, reclaim whitespace, and keep the template structure strong.
  "fill_color", "fill_transparency", "line_color", "line_width_pt", "line_transparency",
  "layer_position" ("front" or "back") — use these to control theme, panel styling, rules,
  and background-image placement while still staying grounded to the template.

Look at the template slide image carefully. Cross-reference what you see visually
with the shape list (positions, text content, font sizes). The shape with the
biggest font and the template's main text (e.g. "Startup.") is the HERO TITLE.
Smaller text shapes with lorem ipsum or placeholder text are BODY areas to fill.

You are allowed to adapt typography and layout usage to make the slide actually
work. If text would overflow, collide, or become unreadable, shorten the copy,
split it across multiple body shapes, reduce the font size, change the font family,
change case/weight/color, or clear low-value placeholders. Never leave template
brand text, lorem ipsum, contact info, or stray placeholders visible unless the
content explicitly needs them. Never create micro-text just to fill every slot.
Preserve the template's typography and hierarchy by default. Only override the font
family if the current one is clearly wrong for the content or render-safe substitution
is needed. Avoid generic labels like "sans-serif"; use concrete families.
If the brief implies a theme shift such as black-dominant, dark luxury, or brighter editorial,
you may restyle template shapes with fill/line colors and move large image shapes behind content.

Return ONLY valid JSON — every shape must have an instruction:
{
  "speaker_notes": "...",
  "instructions": [
    {"shape_index": 0, "action": "...", ...},
    {"shape_index": 1, "action": "...", ...},
    ...
  ]
}
"""

_MAP_USER = """\
TEMPLATE SLIDE LAYOUT:
{visual_description}

SHAPES (extracted via python-pptx — use [index] to reference):
{shapes_text}

CONTENT TO PLACE:
  Title: {title}
  Deck theme: {deck_theme}
  Role: {content_role}
  Layout reasoning: {layout_reasoning}
  Content blocks: {content_json}
  Speaker notes: {speaker_notes}
  Render-safe font families: Bahnschrift, Georgia, Verdana, Arial, Calibri

Study the template slide image above. Match each visual element you see to a shape
in the list by its position and text. Then map the content to the right shapes.
Return JSON with shape_index (integer), not shape names.
"""


def _get_template_thumbnail_b64(catalog_dir: str, template_slide_num: int) -> str | None:
    """Load the catalog thumbnail PNG for a template slide as base64."""
    if not catalog_dir:
        return None
    thumb = Path(catalog_dir) / "thumbnails" / f"slide-{template_slide_num:02d}.png"
    if thumb.exists():
        return base64.b64encode(thumb.read_bytes()).decode("utf-8")
    return None


_VALID_ACTIONS = {"fill_text", "fill_image", "fill_icon", "fill_chart", "style_shape", "clear", "skip"}
_VALID_CLEAR_MODES = {"text_only", "delete_shape"}
_VALID_LAYER_POSITIONS = {"front", "back"}


def _slide_payload(spec: dict, slide_info: dict) -> dict:
    payload = dict(slide_info)
    payload["deck_theme"] = spec.get("deck_theme", "")
    return payload


def _extract_template_shapes(prs: Presentation, template_slide_number: int, sw: int, sh: int) -> list[dict]:
    idx = template_slide_number - 1
    if 0 <= idx < len(prs.slides):
        return extract_shapes(prs.slides[idx], sw, sh)
    return []


def _coerce_shape_index(instr: dict, shapes: list[dict]) -> int | None:
    raw = instr.get("shape_index")
    if isinstance(raw, float) and raw.is_integer():
        raw = int(raw)
    elif isinstance(raw, str):
        raw = raw.strip()
        if raw.isdigit():
            raw = int(raw)
    if isinstance(raw, int) and 0 <= raw < len(shapes):
        return raw

    shape_name = str(instr.get("shape_name") or "").strip()
    if shape_name:
        for shape in shapes:
            if shape.get("name") == shape_name:
                return int(shape["index"])
    return None


def _normalize_mapping(
    result: dict,
    shapes: list[dict],
    *,
    fallback_speaker_notes: str = "",
    base_mapping: dict | None = None,
) -> dict:
    instructions_by_index: dict[int, dict] = {}
    base_by_index: dict[int, dict] = {}
    for base_instr in (base_mapping or {}).get("instructions", []) or []:
        if isinstance(base_instr, dict):
            base_idx = _coerce_shape_index(base_instr, shapes)
            if base_idx is not None:
                base_by_index[base_idx] = dict(base_instr)

    for raw_instr in result.get("instructions", []) or []:
        if not isinstance(raw_instr, dict):
            continue

        shape_index = _coerce_shape_index(raw_instr, shapes)
        if shape_index is None:
            continue

        instr = dict(raw_instr)
        instr["shape_index"] = shape_index

        action = str(instr.get("action") or "skip").strip().lower()
        if action not in _VALID_ACTIONS:
            action = "skip"
        instr["action"] = action

        clear_mode = str(instr.get("clear_mode") or "").strip().lower()
        if clear_mode in _VALID_CLEAR_MODES:
            instr["clear_mode"] = clear_mode
        else:
            instr.pop("clear_mode", None)

        if action == "fill_text":
            text = str(instr.get("text") or "").strip()
            if not text:
                instr = {"shape_index": shape_index, "action": "clear"}
            else:
                shape_meta = shapes[shape_index] if 0 <= shape_index < len(shapes) else {}
                if (
                    "\n" not in text
                    and len(text.split()) >= 3
                    and (shape_meta.get("font_pt") or 0) >= 96
                ):
                    text = _balanced_display_text(text)
                    instr["single_paragraph"] = False
                instr["text"] = text
        elif action in {"fill_image", "fill_icon"}:
            asset_id = str(instr.get("asset_id") or "").strip()
            query = str(instr.get("search_query") or "").strip()
            if action == "fill_icon":
                if not query:
                    instr = {"shape_index": shape_index, "action": "skip"}
                else:
                    instr["search_query"] = query
            else:
                if not asset_id and not query:
                    instr = {"shape_index": shape_index, "action": "skip"}
                else:
                    if asset_id:
                        instr["asset_id"] = asset_id
                    if query:
                        instr["search_query"] = query
                image_mode = str(instr.get("image_mode") or "cover").strip().lower()
                instr["image_mode"] = image_mode if image_mode in {"cover", "contain"} else "cover"
        elif action == "fill_chart":
            asset_id = str(instr.get("asset_id") or "").strip()
            chart_title = str(instr.get("chart_title") or "Chart").strip()
            chart_type = str(instr.get("chart_type") or "bar").strip().lower()
            if asset_id:
                instr["asset_id"] = asset_id
            instr["chart_title"] = chart_title or "Chart"
            instr["chart_type"] = chart_type or "bar"
        elif action == "style_shape":
            pass

        for field in ("left_in", "top_in", "width_in", "height_in"):
            numeric = _coerce_optional_float(instr.get(field))
            if numeric is not None:
                instr[field] = numeric
            else:
                instr.pop(field, None)

        text_alignment = str(instr.get("text_alignment") or "").strip().lower()
        if text_alignment not in {"left", "center", "right", "justify"}:
            instr.pop("text_alignment", None)
        else:
            instr["text_alignment"] = text_alignment

        layer_position = str(instr.get("layer_position") or "").strip().lower()
        if layer_position not in _VALID_LAYER_POSITIONS:
            instr.pop("layer_position", None)
        else:
            instr["layer_position"] = layer_position

        for field in ("fill_transparency", "line_transparency", "image_focus_x", "image_focus_y", "image_opacity"):
            numeric = _coerce_unit_ratio(instr.get(field))
            if numeric is None:
                instr.pop(field, None)
            else:
                instr[field] = numeric

        for field in ("fill_color", "line_color"):
            raw = str(instr.get(field) or "").strip()
            if not raw:
                instr.pop(field, None)
            elif raw.lower() in {"none", "transparent", "no-fill", "no fill", "clear"}:
                instr[field] = "none"
            elif _parse_rgb(raw):
                instr[field] = raw
            else:
                instr.pop(field, None)

        line_width_pt = _coerce_optional_float(instr.get("line_width_pt"))
        if line_width_pt is None:
            instr.pop("line_width_pt", None)
        else:
            instr["line_width_pt"] = max(0.0, min(24.0, line_width_pt))

        style_source_index = _coerce_shape_index({"shape_index": instr.get("style_source_index")}, shapes)
        if style_source_index is not None:
            instr["style_source_index"] = style_source_index
        else:
            instr.pop("style_source_index", None)

        base_instr = base_by_index.get(shape_index, {})
        if base_instr and base_instr.get("action") == instr.get("action"):
            for field in (
                "left_in", "top_in", "width_in", "height_in",
                "style_source_index", "text_alignment", "fill_color", "fill_transparency",
                "line_color", "line_width_pt", "line_transparency", "layer_position",
                "image_mode", "image_focus_x", "image_focus_y", "image_opacity",
                "font_family", "font_size_pt", "font_scale", "text_color",
                "bold", "italic", "uppercase", "single_paragraph", "disable_fit", "icon_color",
            ):
                if field not in instr and field in base_instr:
                    instr[field] = base_instr[field]

        instructions_by_index[shape_index] = instr

    for shape in shapes:
        idx = int(shape["index"])
        if idx in instructions_by_index:
            continue
        if idx in base_by_index:
            preserved = dict(base_by_index[idx])
            preserved["shape_index"] = idx
            instructions_by_index[idx] = preserved
        else:
            instructions_by_index[idx] = {"shape_index": idx, "action": "skip"}

    return {
        "speaker_notes": str(result.get("speaker_notes") or fallback_speaker_notes or ""),
        "instructions": [instructions_by_index[idx] for idx in sorted(instructions_by_index)],
    }


def _map_one_slide(
    slide_info: dict,
    shapes: list[dict],
    *,
    output_dir: Path,
    template_thumb_b64: str | None = None,
) -> tuple[dict, dict[str, str]]:
    """Call LLM to map content to shapes for one slide.

    If template_thumb_b64 is provided, sends the template slide image
    to the vision model for more accurate shape-to-content mapping.

    Returns: {"speaker_notes": str, "instructions": [...]}
    """
    user_text = _MAP_USER.format(
        visual_description=slide_info.get("template_visual_description", ""),
        shapes_text=_shapes_for_prompt(shapes),
        title=slide_info.get("title", ""),
        deck_theme=slide_info.get("deck_theme", ""),
        content_role=slide_info.get("content_role", ""),
        layout_reasoning=slide_info.get("layout_reasoning", ""),
        content_json=json.dumps(slide_info.get("full_content", []), ensure_ascii=False),
        speaker_notes=slide_info.get("speaker_notes", ""),
    )

    # Build user message — with or without vision
    if template_thumb_b64:
        user_content: Any = [
            {"type": "image_url",
             "image_url": {"url": f"data:image/png;base64,{template_thumb_b64}"}},
            {"type": "text",
             "text": f"Above is the template slide you are filling into.\n\n{user_text}"},
        ]
    else:
        user_content = user_text

    messages = [
        {"role": "system", "content": _MAP_SYSTEM},
        {"role": "user",   "content": user_content},
    ]

    for attempt in range(3):
        raw = ""
        try:
            result, tool_assets = _run_builder_stage_with_tools(
                messages,
                output_dir=output_dir,
                temperature=0.2,
                final_hint='{"speaker_notes":"...","instructions":[{"shape_index":0,"action":"fill_text|fill_image|fill_icon|fill_chart|style_shape|clear|skip"}]}',
            )
            if "instructions" not in result:
                raise ValueError("Missing 'instructions' key")
            return result, tool_assets
        except Exception as exc:
            logger.warning("map_shapes slide '%s' attempt %d: %s",
                           slide_info.get("title", "?"), attempt + 1, exc)
            if raw:
                logger.debug("  raw response (first 500 chars): %s", raw[:500])

    # Fallback — fill the biggest text shape with the title
    logger.error("map_shapes fallback for '%s'", slide_info.get("title"))
    text_shapes = [s for s in shapes if s.get("type") == "text"]
    # Pick largest text shape by area
    text_shapes.sort(key=lambda s: s.get("width_in", 0) * s.get("height_in", 0), reverse=True)
    fallback = [{"shape_index": s["index"], "action": "skip"} for s in shapes]
    if text_shapes:
        fallback.append({"shape_index": text_shapes[0]["index"], "action": "fill_text",
                        "text": slide_info.get("title", "")})
    return {"speaker_notes": slide_info.get("speaker_notes", ""), "instructions": fallback}, {}


_POLISH_SYSTEM = """\
You are a presentation design director doing a SECOND PASS over an initial slide mapping.

Your job is to turn a merely plausible mapping into the strongest final mapping for the
template and content. You are not required to preserve the first draft. Rebalance it.

What to optimize:
1. Visual hierarchy — the slide must have a clear hero, support copy, and whitespace
2. Legibility — no micro-text, no overloaded sidebars, no forcing content into tiny slots
3. Design sense — better font choices, contrast, icon usage, and emphasis when needed
4. Template intelligence — use the template's intended structure, but do not obey bad placeholders
5. Content judgment — decide what to condense, split, move, clear, or emphasize

Allowed actions per shape:
  fill_text     — with optional "font_family", "font_size_pt", "font_scale",
                  "text_color", "bold", "italic", "uppercase", "single_paragraph",
                  "disable_fit", "style_source_index", "text_alignment",
                  plus "fill_color", "line_color", "fill_transparency"
  fill_image    — with "search_query" or "asset_id", plus optional "image_mode", "image_focus_x",
                  "image_focus_y", "image_opacity", "layer_position"
  fill_icon     — with "search_query" and "icon_color"
  fill_chart    — with "chart_title", "chart_type", optional "content_block_index", or "asset_id"
  style_shape   — use for recoloring panels, dividers, text-box backgrounds, and decorative shapes
  clear         — optional "clear_mode" = "text_only" or "delete_shape"
  skip

Optional on any instruction:
  "left_in", "top_in", "width_in", "height_in" to reposition or resize an existing template shape.
  "fill_color", "fill_transparency", "line_color", "line_width_pt", "line_transparency",
  "layer_position" for theme control and background treatments.

Principles:
- Prefer fewer strong elements over many weak ones.
- If a contact-info slot, footer line, or tiny text box cannot carry meaningful readable content,
  clear it, often with "delete_shape".
- If a title or body area is cramped, rewrite or redistribute the copy instead of shrinking it into noise.
- Use icon placeholders only when the result adds clarity and remains visually balanced.
- Preserve the template's mood, but upgrade typography if the default style is weak.
- Stay grounded to the template. Rework bad placements by moving/resizing template shapes before inventing new styling.
- Use shape styling when the slide needs a stronger palette, darker blocks, tinted containers,
  or an image pushed into the background behind text.

Return ONLY valid JSON:
{
  "speaker_notes": "...",
  "instructions": [
    {"shape_index": 0, "action": "...", ...}
  ]
}
"""

_POLISH_USER = """\
TEMPLATE SLIDE LAYOUT:
{visual_description}

SHAPES:
{shapes_text}

SLIDE CONTENT:
  Title: {title}
  Deck theme: {deck_theme}
  Role: {content_role}
  Layout reasoning: {layout_reasoning}
  Content blocks: {content_json}
  Speaker notes: {speaker_notes}
  Render-safe font families: Bahnschrift, Georgia, Verdana, Arial, Calibri

INITIAL MAPPING DRAFT:
{initial_mapping}

Produce the best FINAL mapping for this slide. Fix weak hierarchy, overcrowding,
tiny text, poor icon usage, placeholder leakage, and bad typography choices.
Return JSON only.
"""


def _polish_one_slide(
    slide_info: dict,
    shapes: list[dict],
    initial_mapping: dict,
    *,
    output_dir: Path,
    template_thumb_b64: str | None = None,
) -> tuple[dict, dict[str, str]]:
    user_text = _POLISH_USER.format(
        visual_description=slide_info.get("template_visual_description", ""),
        shapes_text=_shapes_for_prompt(shapes),
        title=slide_info.get("title", ""),
        deck_theme=slide_info.get("deck_theme", ""),
        content_role=slide_info.get("content_role", ""),
        layout_reasoning=slide_info.get("layout_reasoning", ""),
        content_json=json.dumps(slide_info.get("full_content", []), ensure_ascii=False),
        speaker_notes=slide_info.get("speaker_notes", ""),
        initial_mapping=json.dumps(initial_mapping, indent=2, ensure_ascii=False),
    )

    if template_thumb_b64:
        user_content: Any = [
            {"type": "image_url",
             "image_url": {"url": f"data:image/png;base64,{template_thumb_b64}"}} ,
            {"type": "text",
             "text": f"Above is the template slide you are finalizing.\n\n{user_text}"},
        ]
    else:
        user_content = user_text

    messages = [
        {"role": "system", "content": _POLISH_SYSTEM},
        {"role": "user", "content": user_content},
    ]

    for attempt in range(2):
        raw = ""
        try:
            result, tool_assets = _run_builder_stage_with_tools(
                messages,
                output_dir=output_dir,
                temperature=0.15,
                final_hint='{"speaker_notes":"...","instructions":[{"shape_index":0,"action":"fill_text|fill_image|fill_icon|fill_chart|style_shape|clear|skip"}]}',
            )
            if "instructions" not in result:
                raise ValueError("Missing 'instructions' key")
            return result, tool_assets
        except Exception as exc:
            logger.warning("polish slide '%s' attempt %d: %s",
                           slide_info.get("title", "?"), attempt + 1, exc)
            if raw:
                logger.debug("  raw response (first 500 chars): %s", raw[:500])

    logger.error("polish fallback for '%s' — keeping initial mapping", slide_info.get("title"))
    return initial_mapping, {}


# ═══════════════════════════════════════════════════════════════════════════════
# VALIDATION — vision LLM checks rendered slides
# ═══════════════════════════════════════════════════════════════════════════════

_VAL_SYSTEM = """\
You are a presentation quality reviewer. Review the rendered slide against the template slide,
the template shape inventory, the active placement summary, and the soft layout signals.

Important:
- Distinguish DELIBERATE overlap from accidental collision.
- Text inside a stat circle, on a glass card, or as a stacked hero treatment can be correct.
- Only flag overlap when readability, ordering, containment, or hierarchy is harmed.
- Treat the supplied layout signals as hints, not as ground truth.

Check for:
1. Text overflow or truncation (content cut off or extending beyond its intended area)
2. Accidental collisions (body text crossing dividers, labels colliding, scattered captions)
3. Empty zones that clearly should have content (unfilled placeholder text still visible)
4. Poor contrast (light text on light background or vice versa)
5. Professional appearance and hierarchy — does it look intentional and well-composed?
6. Distorted icons/photos, visible white boxes, or missing media
7. Chart placeholders or obviously fake data visuals where real chart data should exist
8. Style drift — typography or spacing that breaks the template's design language

Return ONLY JSON:
{
  "verdict": "pass" or "fail",
  "score": <1-10>,
  "summary": "short overall judgment",
  "issues": [
    {
      "category": "overflow|collision|alignment|hierarchy|contrast|missing_asset|placeholder|chart|style_drift|other",
      "severity": "low|medium|high",
      "confidence": 0.0,
      "description": "what is wrong",
      "target_shape_indices": [1, 4],
      "accidental_overlap": true,
      "fix_hint": "what to change"
    }
  ]
}
"""


_VAL_USER = """\
SLIDE NUMBER: {slide_num}

TEMPLATE VISUAL DESCRIPTION:
{visual_description}

TEMPLATE SHAPES:
{shapes_text}

ACTIVE PLACEMENT SUMMARY:
{mapping_text}

SOFT LAYOUT SIGNALS:
{layout_signals}

Review the rendered slide using the context above. Focus on final visual quality, not just whether text exists.
Return JSON only.
"""


def _normalize_validation_issue(issue: Any, shapes: list[dict]) -> dict:
    if isinstance(issue, str):
        return {
            "category": "other",
            "severity": "medium",
            "confidence": 0.7,
            "description": issue.strip(),
            "target_shape_indices": [],
            "accidental_overlap": None,
            "fix_hint": "",
        }
    if not isinstance(issue, dict):
        return {
            "category": "other",
            "severity": "medium",
            "confidence": 0.5,
            "description": str(issue).strip(),
            "target_shape_indices": [],
            "accidental_overlap": None,
            "fix_hint": "",
        }
    category = str(issue.get("category") or "other").strip().lower()
    if category not in {
        "overflow", "collision", "alignment", "hierarchy", "contrast",
        "missing_asset", "placeholder", "chart", "style_drift", "other",
    }:
        category = "other"
    severity = str(issue.get("severity") or "medium").strip().lower()
    if severity not in {"low", "medium", "high"}:
        severity = "medium"
    confidence = _coerce_optional_float(issue.get("confidence"))
    if confidence is None:
        confidence = 0.6
    confidence = max(0.0, min(1.0, confidence))
    targets: list[int] = []
    for raw in issue.get("target_shape_indices", []) or []:
        idx = _coerce_shape_index({"shape_index": raw}, shapes)
        if idx is not None and idx not in targets:
            targets.append(idx)
    accidental_overlap = issue.get("accidental_overlap")
    if accidental_overlap not in {True, False, None}:
        accidental_overlap = None
    description = str(issue.get("description") or "").strip()
    if not description:
        description = "Validation flagged a presentation quality issue."
    fix_hint = str(issue.get("fix_hint") or "").strip()
    return {
        "category": category,
        "severity": severity,
        "confidence": round(confidence, 2),
        "description": description,
        "target_shape_indices": targets,
        "accidental_overlap": accidental_overlap,
        "fix_hint": fix_hint,
    }


def _validation_issue_text(issue: dict) -> str:
    parts = [f"[{issue.get('severity', 'medium')}/{issue.get('category', 'other')}] {issue.get('description', '')}"]
    targets = issue.get("target_shape_indices") or []
    if targets:
        parts.append(f"targets={targets}")
    if issue.get("accidental_overlap") is True:
        parts.append("accidental_overlap=true")
    elif issue.get("accidental_overlap") is False:
        parts.append("accidental_overlap=false")
    if issue.get("fix_hint"):
        parts.append(f"hint={issue['fix_hint']}")
    return " | ".join(parts)


def _normalize_validation_result(result: dict, slide_num: int, shapes: list[dict]) -> dict:
    verdict = str(result.get("verdict") or "fail").strip().lower()
    if verdict not in {"pass", "fail"}:
        verdict = "fail"
    score = result.get("score", 0)
    try:
        score = int(float(score))
    except Exception:
        score = 0
    score = max(0, min(10, score))
    issues = [_normalize_validation_issue(issue, shapes) for issue in result.get("issues", []) or []]
    if verdict == "pass" and issues:
        verdict = "fail"
    summary = str(result.get("summary") or "").strip()
    if not summary:
        summary = "Validation passed." if verdict == "pass" else "Validation found issues."
    return {
        "slide_number": slide_num,
        "verdict": verdict,
        "score": score,
        "summary": summary,
        "issues": issues,
    }


def _merge_validation_issues(existing: list[dict], incoming: list[dict]) -> list[dict]:
    merged: list[dict] = []
    seen: set[tuple[str, tuple[int, ...]]] = set()
    for issue in list(existing) + list(incoming):
        if not isinstance(issue, dict):
            continue
        key = (str(issue.get("description") or "").strip(), tuple(issue.get("target_shape_indices") or []))
        if key in seen:
            continue
        merged.append(issue)
        seen.add(key)
    return merged


def _validate_one_slide(
    png_path: Path,
    slide_num: int,
    *,
    slide_info: dict,
    shapes: list[dict],
    mapping: dict,
    slide_width_in: float,
    slide_height_in: float,
    template_thumb_b64: str | None = None,
) -> dict:
    """Call vision LLM to validate one rendered slide."""
    b64 = base64.b64encode(png_path.read_bytes()).decode("utf-8")
    user_text = _VAL_USER.format(
        slide_num=slide_num,
        visual_description=slide_info.get("template_visual_description", ""),
        shapes_text=_shapes_for_prompt(shapes),
        mapping_text=_mapping_for_prompt(mapping, shapes),
        layout_signals=_layout_signals_for_prompt(
            shapes,
            mapping,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
        ),
    )
    content_parts: list[Any] = [
        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
        {"type": "text", "text": "Above: the rendered slide to review."},
    ]
    if template_thumb_b64:
        content_parts.extend([
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{template_thumb_b64}"}},
            {"type": "text", "text": "Above: the original template slide."},
        ])
    content_parts.append({"type": "text", "text": user_text})
    messages = [
        {"role": "system", "content": _VAL_SYSTEM},
        {"role": "user", "content": content_parts},
    ]
    try:
        raw = _call_llm(messages, temp=0.1)
        return _normalize_validation_result(_parse_json(raw), slide_num, shapes)
    except Exception as exc:
        return _normalize_validation_result({
            "verdict": "fail",
            "score": 0,
            "summary": "Validator could not produce reliable JSON.",
            "issues": [f"Validator could not produce reliable JSON: {exc}"],
        }, slide_num, shapes)


# ═══════════════════════════════════════════════════════════════════════════════
# LANGGRAPH NODES
# ═══════════════════════════════════════════════════════════════════════════════

def load_node(state: BuilderState) -> dict:
    """Load and validate inputs."""
    spec = state["build_spec"]
    tpl  = Path(state["template_path"])
    out  = Path(state["output_dir"])
    out.mkdir(parents=True, exist_ok=True)

    if not tpl.exists():
        return {"errors": [f"Template not found: {tpl}"]}
    if not spec.get("slides"):
        return {"errors": ["build_spec has no slides"]}

    logger.info("load: %d slides, template=%s, output=%s",
                len(spec["slides"]), tpl.name, out)
    return {"errors": []}


def _write_json_artifact(output_dir: Path, filename: str, payload: dict) -> str:
    path = output_dir / filename
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    return str(path)


def plan_content_node(state: BuilderState) -> dict:
    """Generate a deck content plan when starting from a description."""
    if state.get("content_plan", {}).get("slides"):
        return {}
    if state.get("build_spec", {}).get("slides"):
        return {}

    description = str(state.get("description") or "").strip()
    if not description:
        return {"errors": ["Missing description for planning stage."]}

    logger.info("plan_content: generating deck plan …")
    plan = plan_content(description, num_slides=state.get("max_slides"))
    out_dir = Path(state["output_dir"])
    plan_path = _write_json_artifact(out_dir, "plan.json", plan)
    logger.info("plan_content: %d slides → %s", len(plan.get("slides", [])), plan_path)
    return {"content_plan": plan}


def select_layouts_node(state: BuilderState) -> dict:
    """Create a build spec from the plan and external template catalog."""
    if state.get("build_spec", {}).get("slides"):
        return {}

    plan = state.get("content_plan", {})
    if not plan.get("slides"):
        return {"errors": ["No content plan available for layout selection."]}

    catalog = state.get("template_catalog", {})
    if not catalog:
        return {"errors": ["Template catalog missing; cataloging must run before the builder graph."]}

    logger.info("select_layouts: matching plan to template slides …")
    build_spec = select_layouts(plan, catalog, max_slides=state.get("max_slides"))
    spec_path = _write_json_artifact(Path(state["output_dir"]), "build_spec.json", build_spec)
    logger.info("select_layouts: %d slides → %s", len(build_spec.get("slides", [])), spec_path)
    return {"build_spec": build_spec}


def _collect_asset_requests(mappings: list[dict]) -> list[dict]:
    """Extract asset requests from shape mapping instructions."""
    asset_reqs: list[dict] = []
    for mapping in mappings:
        for instr in mapping.get("instructions", []):
            act = instr.get("action", "")
            if act == "fill_image" and instr.get("asset_id"):
                continue
            if act == "fill_chart" and instr.get("asset_id"):
                continue
            if act == "fill_image" and instr.get("search_query"):
                asset_reqs.append({"type": "photo", "query": instr["search_query"],
                                   "key": f"photo:{instr['search_query']}"})
            elif act == "fill_icon" and instr.get("search_query"):
                color = instr.get("icon_color", "#222222")
                asset_reqs.append({"type": "icon", "query": instr["search_query"],
                                   "color": color,
                                   "key": f"icon:{instr['search_query']}:{color}"})
            elif act == "fill_chart":
                ct = instr.get("chart_title", "Chart")
                ty = instr.get("chart_type", "bar")
                asset_reqs.append({"type": "chart", "title": ct, "chart_type": ty,
                                   "key": f"chart:{ct}:{ty}"})
    return asset_reqs


def map_shapes_node(state: BuilderState) -> dict:
    """For each slide, extract shapes from the template and call LLM to decide fill mapping.

    Sends the template slide thumbnail (vision) so the LLM can see exactly
    what the slide looks like and make better fill decisions.
    """
    spec = state["build_spec"]
    tpl  = Path(state["template_path"])
    catalog_dir = state.get("catalog_dir", "")
    prs  = Presentation(str(tpl))
    total = len(spec["slides"])
    sw = prs.slide_width
    sh = prs.slide_height

    def _process(slide_info: dict) -> tuple[int, dict, dict[str, str]]:
        dsn = slide_info["deck_slide_number"]
        tsn = slide_info["template_slide_number"]
        slide_payload = _slide_payload(spec, slide_info)
        shapes = _extract_template_shapes(prs, tsn, sw, sh)
        thumb_b64 = _get_template_thumbnail_b64(catalog_dir, tsn)
        raw_mapping, tool_assets = _map_one_slide(
            slide_payload,
            shapes,
            output_dir=Path(state["output_dir"]),
            template_thumb_b64=thumb_b64,
        )
        mapping = _normalize_mapping(
            raw_mapping,
            shapes,
            fallback_speaker_notes=slide_payload.get("speaker_notes", ""),
        )
        return dsn, mapping, tool_assets

    logger.info("map_shapes: mapping %d slides (parallelism=%d) …", total, PARALLELISM)
    results: dict[int, dict] = {}
    tool_asset_paths = dict(state.get("asset_paths", {}))
    with ThreadPoolExecutor(max_workers=PARALLELISM) as pool:
        futures = {pool.submit(_process, si): si["deck_slide_number"]
                   for si in spec["slides"]}
        for fut in as_completed(futures):
            dsn, mapping, new_tool_assets = fut.result()
            results[dsn] = mapping
            tool_asset_paths.update(new_tool_assets)
            logger.info("  [%d/%d] mapped", dsn, total)

    # Ordered by deck_slide_number
    mappings = [results[si["deck_slide_number"]] for si in spec["slides"]]

    return {
        "slide_mappings": mappings,
        "asset_requests": _collect_asset_requests(mappings),
        "asset_paths": tool_asset_paths,
    }


def polish_mappings_node(state: BuilderState) -> dict:
    """Run a second LLM pass to refine hierarchy, density, and typography choices."""
    spec = state["build_spec"]
    tpl = Path(state["template_path"])
    catalog_dir = state.get("catalog_dir", "")
    mappings = state.get("slide_mappings", [])
    prs = Presentation(str(tpl))
    total = len(spec["slides"])
    sw = prs.slide_width
    sh = prs.slide_height

    def _process(index: int, slide_info: dict) -> tuple[int, dict, dict[str, str]]:
        dsn = slide_info["deck_slide_number"]
        tsn = slide_info["template_slide_number"]
        slide_payload = _slide_payload(spec, slide_info)
        shapes = _extract_template_shapes(prs, tsn, sw, sh)
        thumb_b64 = _get_template_thumbnail_b64(catalog_dir, tsn)
        initial_mapping = mappings[index] if index < len(mappings) else {"instructions": []}
        polished_raw, tool_assets = _polish_one_slide(
            slide_payload,
            shapes,
            initial_mapping,
            output_dir=Path(state["output_dir"]),
            template_thumb_b64=thumb_b64,
        )
        polished = _normalize_mapping(
            polished_raw,
            shapes,
            fallback_speaker_notes=slide_payload.get("speaker_notes", ""),
            base_mapping=initial_mapping,
        )
        return dsn, polished, tool_assets

    logger.info("polish_mappings: refining %d slides (parallelism=%d) …", total, PARALLELISM)
    results: dict[int, dict] = {}
    tool_asset_paths = dict(state.get("asset_paths", {}))
    with ThreadPoolExecutor(max_workers=PARALLELISM) as pool:
        futures = {
            pool.submit(_process, idx, slide_info): slide_info["deck_slide_number"]
            for idx, slide_info in enumerate(spec["slides"])
        }
        for fut in as_completed(futures):
            dsn, polished, new_tool_assets = fut.result()
            results[dsn] = polished
            tool_asset_paths.update(new_tool_assets)
            logger.info("  [%d/%d] polished", dsn, total)

    polished_mappings = [results[si["deck_slide_number"]] for si in spec["slides"]]
    return {
        "slide_mappings": polished_mappings,
        "asset_requests": _collect_asset_requests(polished_mappings),
        "asset_paths": tool_asset_paths,
    }


def fetch_assets_node(state: BuilderState) -> dict:
    """Download all icons, photos, generate chart placeholders — in parallel."""
    reqs = state.get("asset_requests", [])
    existing_paths = dict(state.get("asset_paths", {}))
    if not reqs:
        return {"asset_paths": existing_paths}

    # Deduplicate by key
    unique: dict[str, dict] = {r["key"]: r for r in reqs}
    logger.info("fetch_assets: %d unique assets to fetch …", len(unique))

    def _fetch(req: dict) -> tuple[str, str | None]:
        key = req["key"]
        try:
            if req["type"] == "photo":
                path = resolve_photo(req["query"])
                return key, str(path) if path else None
            elif req["type"] == "icon":
                path = resolve_icon(req["query"], color=req.get("color"))
                return key, str(path) if path else None
            elif req["type"] == "chart":
                path = generate_chart_placeholder(req["title"], req["chart_type"])
                return key, str(path)
        except Exception as exc:
            logger.warning("fetch_assets '%s': %s", key, exc)
        return key, None

    paths: dict[str, str] = dict(existing_paths)
    with ThreadPoolExecutor(max_workers=PARALLELISM) as pool:
        futures = {pool.submit(_fetch, r): r["key"] for r in unique.values()}
        for fut in as_completed(futures):
            key, path = fut.result()
            if path:
                paths[key] = path
                logger.info("  asset: %s → %s", key, Path(path).name)

    return {"asset_paths": paths}


def build_pptx_node(state: BuilderState) -> dict:
    """Clone template slides and fill content. Sequential — python-pptx is single-threaded."""
    spec     = state["build_spec"]
    tpl      = Path(state["template_path"])
    out_dir  = Path(state["output_dir"])
    mappings = state["slide_mappings"]
    assets   = state.get("asset_paths", {})
    errors   = list(state.get("errors", []))
    build_issues: list[dict] = []

    prs = Presentation(str(tpl))
    original_count = len(prs.slides)
    slide_width_in = prs.slide_width / 914400
    slide_height_in = prs.slide_height / 914400

    # Phase 1: Clone needed template slides in deck order
    logger.info("build: cloning %d slides …", len(spec["slides"]))
    for slide_info in spec["slides"]:
        tsn = slide_info["template_slide_number"]
        idx = max(0, min(tsn - 1, original_count - 1))
        _clone_slide(prs, idx)

    # Phase 2: Delete all original template slides
    for _ in range(original_count):
        _delete_slide(prs, 0)

    # Now prs has only the cloned slides in deck order
    logger.info("build: filling %d slides …", len(prs.slides))

    for i, slide in enumerate(prs.slides):
        slide_info = spec["slides"][i]
        mapping    = mappings[i] if i < len(mappings) else {}
        instrs     = mapping.get("instructions", [])
        dsn        = slide_info.get("deck_slide_number", i + 1)

        # Build shape lookups — index-based (primary) + name-based (fallback)
        shape_list = list(slide.shapes)
        shape_by_name: dict[str, Any] = {s.name: s for s in shape_list}
        format_by_index = {
            idx: _capture_text_format(shape)
            for idx, shape in enumerate(shape_list)
            if getattr(shape, "has_text_frame", False)
        }
        original_text_by_index = {
            idx: shape.text_frame.text
            for idx, shape in enumerate(shape_list)
            if getattr(shape, "has_text_frame", False)
        }
        touched_shape_indices: set[int] = set()

        for instr in instrs:
            action = instr.get("action", "skip")
            if action == "skip":
                continue

            # Resolve shape: prefer shape_index, fall back to shape_name
            shape = None
            si = instr.get("shape_index")
            if si is not None and 0 <= si < len(shape_list):
                shape = shape_list[si]
            else:
                sname = instr.get("shape_name", "")
                shape = shape_by_name.get(sname)

            if not shape:
                ref = si if si is not None else instr.get("shape_name", "?")
                logger.warning("build: slide %d — shape %s not found, skipping", dsn, ref)
                continue
            resolved_idx = si if si is not None and 0 <= si < len(shape_list) else None
            if resolved_idx is None:
                for idx, candidate in enumerate(shape_list):
                    if candidate == shape:
                        resolved_idx = idx
                        break
            if resolved_idx is not None:
                touched_shape_indices.add(resolved_idx)

            try:
                _apply_shape_transform(
                    shape,
                    instr,
                    slide_width_in=slide_width_in,
                    slide_height_in=slide_height_in,
                )

                if action == "fill_text":
                    style_source_idx = instr.get("style_source_index")
                    fmt = format_by_index.get(style_source_idx) if isinstance(style_source_idx, int) else None
                    if not fmt:
                        fmt = format_by_index.get(resolved_idx) or _capture_text_format(shape)
                    _apply_shape_style(slide, shape, instr)
                    _set_shape_text(shape, instr.get("text", ""), fmt, instr=instr)

                elif action == "fill_image":
                    asset_id = str(instr.get("asset_id") or "").strip()
                    if asset_id:
                        img_path = assets.get(asset_id)
                    else:
                        key = f"photo:{instr.get('search_query', '')}"
                        img_path = assets.get(key)
                    if img_path and Path(img_path).exists():
                        w_px = int(shape.width / 914400 * 150)
                        h_px = int(shape.height / 914400 * 150)
                        resized = resize_image_to_fit(
                            Path(img_path),
                            w_px,
                            h_px,
                            mode=instr.get("image_mode", "cover"),
                            focus_x=instr.get("image_focus_x"),
                            focus_y=instr.get("image_focus_y"),
                            opacity=instr.get("image_opacity"),
                        )
                        picture = _replace_shape_with_image(slide, shape, resized)
                        _move_shape_to_layer(slide, picture, str(instr.get("layer_position") or "").strip().lower())
                    else:
                        build_issues.append({
                            "slide_number": dsn,
                            "issue": (
                                f"Missing generated image asset '{asset_id}'"
                                if asset_id else
                                f"Missing photo asset for '{instr.get('search_query', '')}'"
                            ),
                        })

                elif action == "fill_icon":
                    color = instr.get("icon_color", "#222222")
                    key = f"icon:{instr.get('search_query', '')}:{color}"
                    icon_path = assets.get(key)
                    if icon_path and Path(icon_path).exists():
                        width_in = shape.width / 914400
                        height_in = shape.height / 914400
                        if width_in <= 3.0 and height_in <= 3.0:
                            w_px = int(shape.width / 914400 * 150)
                            h_px = int(shape.height / 914400 * 150)
                            prepared = resize_image_to_fit(
                                Path(icon_path),
                                w_px,
                                h_px,
                                mode="contain",
                                focus_x=instr.get("image_focus_x"),
                                focus_y=instr.get("image_focus_y"),
                                opacity=instr.get("image_opacity"),
                            )
                            picture = _replace_shape_with_image(slide, shape, prepared)
                            _move_shape_to_layer(slide, picture, str(instr.get("layer_position") or "").strip().lower())
                        else:
                            logger.info(
                                "build: slide %d skipped oversized icon target %.2fx%.2fin for '%s'",
                                dsn,
                                width_in,
                                height_in,
                                instr.get("search_query", ""),
                            )
                    else:
                        build_issues.append({
                            "slide_number": dsn,
                            "issue": f"Missing icon asset for '{instr.get('search_query', '')}'",
                        })

                elif action == "fill_chart":
                    asset_id = str(instr.get("asset_id") or "").strip()
                    ct = instr.get("chart_title", "Chart")
                    ty = instr.get("chart_type", "bar")
                    chart_block = _find_chart_block(slide_info, instr)
                    w_px = int(shape.width / 914400 * 150)
                    h_px = int(shape.height / 914400 * 150)
                    chart_path: Path | None = None
                    if asset_id:
                        generated_chart = assets.get(asset_id)
                        if generated_chart and Path(generated_chart).exists():
                            chart_path = Path(generated_chart)
                    elif chart_block:
                        chart_path = generate_chart_image(
                            chart_block,
                            fallback_title=ct,
                            fallback_type=ty,
                            width_px=w_px,
                            height_px=h_px,
                        )
                    else:
                        key = f"chart:{ct}:{ty}"
                        cached_chart = assets.get(key)
                        if cached_chart and Path(cached_chart).exists():
                            chart_path = Path(cached_chart)
                    if chart_path and chart_path.exists():
                        picture = _replace_shape_with_image(slide, shape, chart_path)
                        _move_shape_to_layer(slide, picture, str(instr.get("layer_position") or "").strip().lower())
                    else:
                        build_issues.append({
                            "slide_number": dsn,
                            "issue": (
                                f"Missing generated chart asset '{asset_id}'"
                                if asset_id else
                                f"Missing chart asset for '{ct}'"
                            ),
                        })

                elif action == "style_shape":
                    _apply_shape_style(slide, shape, instr)

                elif action == "clear":
                    clear_mode = str(instr.get("clear_mode") or "text_only").strip().lower()
                    if clear_mode == "delete_shape":
                        _delete_shape(shape)
                    elif shape.has_text_frame:
                        shape.text_frame.text = ""

            except Exception as exc:
                ref = si if si is not None else shape.name
                errors.append(f"slide {dsn} shape {ref}: {exc}")
                logger.warning("build: slide %d shape %s: %s", dsn, ref, exc)

        for idx, shape in enumerate(shape_list):
            if idx in touched_shape_indices or not getattr(shape, "has_text_frame", False):
                continue
            if _looks_like_placeholder_text(original_text_by_index.get(idx, "")):
                try:
                    shape.text_frame.text = ""
                except Exception:
                    continue

        # Speaker notes
        notes = mapping.get("speaker_notes") or slide_info.get("speaker_notes", "")
        if notes:
            try:
                _set_speaker_notes(slide, notes)
            except Exception:
                pass

    pptx_path = out_dir / "deck.pptx"
    pptx_path = _save_presentation(prs, pptx_path)
    logger.info("build: saved → %s", pptx_path)
    return {"pptx_path": str(pptx_path), "errors": errors, "build_issues": build_issues}


def render_node(state: BuilderState) -> dict:
    """Render built PPTX to PNGs for visual inspection and validation."""
    pptx_path = Path(state.get("pptx_path", ""))
    out_dir   = Path(state["output_dir"]) / "slides"

    if not pptx_path.exists():
        return {"slide_pngs": []}

    try:
        pngs = render_template_to_pngs(pptx_path, out_dir, dpi=VALIDATION_RENDER_DPI)
        logger.info("render: %d slide PNGs → %s", len(pngs), out_dir)
        return {"slide_pngs": [str(p) for p in pngs]}
    except Exception as exc:
        logger.warning("render: failed (%s) — validation will be skipped", exc)
        return {"slide_pngs": []}


def validate_node(state: BuilderState) -> dict:
    """Vision LLM checks each rendered slide for quality issues.

    Identifies failed slides and stores their indices for the repair loop.
    """
    if not state.get("do_validate"):
        return {"validation_results": [], "failed_slide_indices": []}

    pngs = [Path(p) for p in state.get("slide_pngs", []) if Path(p).exists()]
    if not pngs:
        return {"validation_results": [], "failed_slide_indices": []}

    spec = state["build_spec"]
    tpl = Path(state["template_path"])
    mappings = state.get("slide_mappings", [])
    catalog_dir = state.get("catalog_dir", "")
    prs = Presentation(str(tpl))
    slide_width_in = prs.slide_width / 914400
    slide_height_in = prs.slide_height / 914400

    logger.info("validate: checking %d slides (parallelism=%d) …", len(pngs), PARALLELISM)

    results: list[dict] = []

    def _validate(idx: int, png: Path) -> dict:
        slide_info = spec["slides"][idx]
        dsn = int(slide_info.get("deck_slide_number", idx + 1))
        tsn = slide_info["template_slide_number"]
        shapes = _extract_template_shapes(prs, tsn, prs.slide_width, prs.slide_height)
        mapping = mappings[idx] if idx < len(mappings) else {"instructions": []}
        template_b64 = _get_template_thumbnail_b64(catalog_dir, tsn)
        return _validate_one_slide(
            png,
            dsn,
            slide_info=_slide_payload(spec, slide_info),
            shapes=shapes,
            mapping=mapping,
            slide_width_in=slide_width_in,
            slide_height_in=slide_height_in,
            template_thumb_b64=template_b64,
        )

    with ThreadPoolExecutor(max_workers=PARALLELISM) as pool:
        futures = {pool.submit(_validate, i, png): i
                   for i, png in enumerate(pngs)}
        for fut in as_completed(futures):
            results.append(fut.result())

    results.sort(key=lambda r: r.get("slide_number", 0))

    issue_map: dict[int, list[dict]] = {}
    for item in state.get("build_issues", []):
        slide_number = int(item.get("slide_number", 0) or 0)
        issue = str(item.get("issue") or "").strip()
        if slide_number and issue:
            issue_map.setdefault(slide_number, []).append({
                "category": "missing_asset" if "Missing " in issue else "other",
                "severity": "high" if "Missing " in issue else "medium",
                "confidence": 1.0,
                "description": issue,
                "target_shape_indices": [],
                "accidental_overlap": None,
                "fix_hint": "Fill or replace the missing asset slot." if "Missing " in issue else "",
            })

    results_by_slide = {r.get("slide_number", 0): r for r in results}
    for slide_number, issues in issue_map.items():
        result = results_by_slide.get(slide_number)
        if result is None:
            result = {
                "slide_number": slide_number,
                "verdict": "fail",
                "score": 0,
                "summary": "Build found structural issues before validation completed.",
                "issues": [],
            }
            results.append(result)
            results_by_slide[slide_number] = result
        result["verdict"] = "fail"
        result["issues"] = _merge_validation_issues(list(result.get("issues", [])), issues)

    results.sort(key=lambda r: r.get("slide_number", 0))

    passed = sum(1 for r in results if r.get("verdict") == "pass")
    index_by_slide_number = {
        int(slide_info.get("deck_slide_number", idx + 1)): idx
        for idx, slide_info in enumerate(spec["slides"])
    }
    failed_indices = [
        index_by_slide_number.get(int(r.get("slide_number", 0) or 0), int(r.get("slide_number", 1) or 1) - 1)
        for r in results
        if r.get("verdict") != "pass"
    ]

    logger.info("validate: %d/%d passed", passed, len(results))
    if failed_indices:
        logger.info("validate: failed slides (0-indexed): %s", failed_indices)

    return {"validation_results": results, "failed_slide_indices": failed_indices}


# ═══════════════════════════════════════════════════════════════════════════════
# REPAIR — re-map failed slides with vision feedback
# ═══════════════════════════════════════════════════════════════════════════════

_REPAIR_SYSTEM = """\
You are a senior presentation designer FIXING a slide that failed quality review.

You will receive:
1. The rendered slide image showing what went WRONG
2. The template slide image showing the ORIGINAL design to follow
3. The shapes on the slide and the content that should be placed
4. The structured validation issues that were found

Your job: produce a CORRECTED fill mapping that fixes the issues.
Common fixes:
- Text overflow → condense text, use shorter phrasing, remove excess bullets
- Empty zones → fill missing shapes that were skipped
- Poor contrast → change icon_color to contrast with background
- Unfilled placeholders → map content to shapes that were missed
- Weak typography → override font family, size, weight, case, or color
- Placeholder charts → point fill_chart at the correct chart data block or use a generated asset_id
- Broken layout → keep or adjust good geometry overrides; use "left_in", "top_in",
  "width_in", "height_in" to reposition existing template shapes instead of scattering text
- Style drift → preserve the template hierarchy and concrete font families; use
  "style_source_index" or "text_alignment" when needed
- Theme mismatch → use "style_shape", "fill_color", "line_color", "fill_transparency",
  "image_opacity", and "layer_position" to restore the intended visual mood
- Deliberate overlap can stay if it reads as intentional. Fix accidental collisions, not designed overlays.

Return ONLY the same JSON format as before — use shape_index (integer).
You may return ONLY the instructions you want to change; omitted shapes will keep their previous mapping:
{
  "speaker_notes": "...",
  "instructions": [
    {"shape_index": 0, "action": "fill_text|fill_image|fill_icon|fill_chart|style_shape|clear|skip", "asset_id": "gen_ab12cd34ef56", "clear_mode": "text_only|delete_shape", "style_source_index": 3, "left_in": 1.2, "top_in": 4.0, "width_in": 6.5, "height_in": 2.0, "fill_color": "#111111", "layer_position": "back", ...}
  ]
}
"""

_REPAIR_USER = """\
VALIDATION SUMMARY:
{validation_summary}

STRUCTURED VALIDATION ISSUES:
{issues_json}

ISSUES AS TEXT:
{issues_text}

TEMPLATE SLIDE VISUAL DESCRIPTION:
{visual_description}

SHAPES ON THIS SLIDE:
{shapes_text}

PREVIOUS ACTIVE MAPPING SUMMARY:
{mapping_summary}

CONTENT TO PLACE:
  Title: {title}
  Deck theme: {deck_theme}
  Role: {content_role}
  Layout reasoning: {layout_reasoning}
  Content blocks: {content_json}
  Speaker notes: {speaker_notes}
  Render-safe font families: Bahnschrift, Georgia, Verdana, Arial, Calibri

PREVIOUS MAPPING THAT FAILED:
{previous_mapping}

Fix the issues above. Preserve strong regions and deliberate overlap treatments. Prefer localized edits on the target shapes rather than remapping the whole slide. Return corrected JSON only.
"""


def _repair_one_slide(
    slide_info: dict,
    shapes: list[dict],
    prev_mapping: dict,
    validation_result: dict,
    *,
    output_dir: Path,
    rendered_slide_b64: str | None = None,
    template_thumb_b64: str | None = None,
) -> tuple[dict, dict[str, str]]:
    """Re-map a single failed slide with validation feedback + vision."""
    structured_issues = validation_result.get("issues", []) or []
    issues_text = "\n".join(f"- {_validation_issue_text(iss)}" for iss in structured_issues)
    if not issues_text:
        issues_text = f"Score: {validation_result.get('score', '?')}/10 — below threshold"

    user_text = _REPAIR_USER.format(
        validation_summary=validation_result.get("summary", ""),
        issues_json=json.dumps(structured_issues, indent=2, ensure_ascii=False),
        issues_text=issues_text,
        visual_description=slide_info.get("template_visual_description", ""),
        shapes_text=_shapes_for_prompt(shapes),
        mapping_summary=_mapping_for_prompt(prev_mapping, shapes),
        title=slide_info.get("title", ""),
        deck_theme=slide_info.get("deck_theme", ""),
        content_role=slide_info.get("content_role", ""),
        layout_reasoning=slide_info.get("layout_reasoning", ""),
        content_json=json.dumps(slide_info.get("full_content", []), ensure_ascii=False),
        speaker_notes=slide_info.get("speaker_notes", ""),
        previous_mapping=json.dumps(prev_mapping, indent=2, ensure_ascii=False),
    )

    # Build vision content — include both rendered (broken) and template (target)
    content_parts: list[Any] = []
    if rendered_slide_b64:
        content_parts.append(
            {"type": "image_url",
             "image_url": {"url": f"data:image/png;base64,{rendered_slide_b64}"}})
        content_parts.append(
            {"type": "text", "text": "Above: the RENDERED slide that FAILED validation."})
    if template_thumb_b64:
        content_parts.append(
            {"type": "image_url",
             "image_url": {"url": f"data:image/png;base64,{template_thumb_b64}"}})
        content_parts.append(
            {"type": "text", "text": "Above: the TEMPLATE slide to match."})
    content_parts.append({"type": "text", "text": user_text})

    messages = [
        {"role": "system", "content": _REPAIR_SYSTEM},
        {"role": "user",   "content": content_parts if content_parts[:-1] else user_text},
    ]

    for attempt in range(2):
        try:
            result, tool_assets = _run_builder_stage_with_tools(
                messages,
                output_dir=output_dir,
                temperature=0.15,
                final_hint='{"speaker_notes":"...","instructions":[{"shape_index":0,"action":"fill_text|fill_image|fill_icon|fill_chart|style_shape|clear|skip"}]}',
            )
            if "instructions" not in result:
                raise ValueError("Missing 'instructions' key")
            return result, tool_assets
        except Exception as exc:
            logger.warning("repair slide '%s' attempt %d: %s",
                           slide_info.get("title", "?"), attempt + 1, exc)

    # If repair fails, return original mapping unchanged
    logger.error("repair fallback for '%s' — keeping original mapping", slide_info.get("title"))
    return prev_mapping, {}


def repair_node(state: BuilderState) -> dict:
    """Re-map only the slides that failed validation, using vision feedback.

    Sends both the broken rendered slide and the target template slide to the LLM
    along with the validation issues, so it can make targeted fixes.
    """
    spec        = state["build_spec"]
    tpl         = Path(state["template_path"])
    catalog_dir = state.get("catalog_dir", "")
    mappings    = list(state["slide_mappings"])  # copy — we'll mutate specific indices
    failed      = state.get("failed_slide_indices", [])
    val_results = state.get("validation_results", [])
    slide_pngs  = state.get("slide_pngs", [])
    repair_round = state.get("repair_round", 0) + 1

    prs = Presentation(str(tpl))
    sw = prs.slide_width
    sh = prs.slide_height

    # Build validation lookup: slide_number → result
    val_map = {r["slide_number"]: r for r in val_results}

    logger.info("repair: round %d — fixing %d failed slides …", repair_round, len(failed))

    def _repair(idx: int) -> tuple[int, dict, dict[str, str]]:
        slide_info = spec["slides"][idx]
        slide_payload = _slide_payload(spec, slide_info)
        dsn = slide_info["deck_slide_number"]
        tsn = slide_info["template_slide_number"]

        # Extract shapes from template
        shapes = _extract_template_shapes(prs, tsn, sw, sh)

        # Load images
        rendered_b64 = None
        if idx < len(slide_pngs) and Path(slide_pngs[idx]).exists():
            rendered_b64 = base64.b64encode(Path(slide_pngs[idx]).read_bytes()).decode()
        template_b64 = _get_template_thumbnail_b64(catalog_dir, tsn)

        result, tool_assets = _repair_one_slide(
            slide_payload, shapes,
            mappings[idx] if idx < len(mappings) else {},
            val_map.get(dsn, {}),
            output_dir=Path(state["output_dir"]),
            rendered_slide_b64=rendered_b64,
            template_thumb_b64=template_b64,
        )
        normalized = _normalize_mapping(
            result,
            shapes,
            fallback_speaker_notes=slide_payload.get("speaker_notes", ""),
            base_mapping=mappings[idx] if idx < len(mappings) else {},
        )
        return idx, normalized, tool_assets

    existing_assets = state.get("asset_paths", {})
    # Parallel repair of failed slides
    repaired_tool_assets = dict(existing_assets)
    with ThreadPoolExecutor(max_workers=PARALLELISM) as pool:
        futures = {pool.submit(_repair, idx): idx for idx in failed}
        for fut in as_completed(futures):
            idx, new_mapping, new_tool_assets = fut.result()
            mappings[idx] = new_mapping
            repaired_tool_assets.update(new_tool_assets)
            logger.info("  repaired slide %d", idx + 1)

    # Collect any new asset requests from repaired mappings
    new_asset_reqs = _collect_asset_requests(mappings)
    existing_assets = repaired_tool_assets

    # Fetch only NEW assets that aren't already cached
    new_unique = {r["key"]: r for r in new_asset_reqs if r["key"] not in existing_assets}
    if new_unique:
        logger.info("repair: fetching %d new assets …", len(new_unique))

        def _fetch(req: dict) -> tuple[str, str | None]:
            key = req["key"]
            try:
                if req["type"] == "photo":
                    path = resolve_photo(req["query"])
                    return key, str(path) if path else None
                elif req["type"] == "icon":
                    path = resolve_icon(req["query"], color=req.get("color"))
                    return key, str(path) if path else None
                elif req["type"] == "chart":
                    path = generate_chart_placeholder(req["title"], req["chart_type"])
                    return key, str(path)
            except Exception as exc:
                logger.warning("repair fetch '%s': %s", key, exc)
            return key, None

        updated_assets = dict(existing_assets)
        with ThreadPoolExecutor(max_workers=PARALLELISM) as pool:
            futs = {pool.submit(_fetch, r): r["key"] for r in new_unique.values()}
            for fut in as_completed(futs):
                key, path = fut.result()
                if path:
                    updated_assets[key] = path
    else:
        updated_assets = existing_assets

    return {
        "slide_mappings": mappings,
        "asset_paths": updated_assets,
        "repair_round": repair_round,
        "failed_slide_indices": [],  # reset — will be re-populated by next validate
    }


def finalize_node(state: BuilderState) -> dict:
    """Write build report and optionally export PDF."""
    out_dir   = Path(state["output_dir"])
    pptx_path = state.get("pptx_path", "")
    errors    = state.get("errors", [])
    val       = state.get("validation_results", [])

    report = {
        "built_at":     datetime.now(timezone.utc).isoformat(),
        "template":     state.get("build_spec", {}).get("template_name", ""),
        "deck_title":   state.get("build_spec", {}).get("deck_title", ""),
        "total_slides": len(state.get("build_spec", {}).get("slides", [])),
        "pptx_path":    pptx_path,
        "slide_pngs":   state.get("slide_pngs", []),
        "slide_mappings": state.get("slide_mappings", []),
        "build_issues": state.get("build_issues", []),
        "validation":   val,
        "errors":       errors,
    }

    report_path = out_dir / "build_report.json"
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")
    logger.info("finalize: report → %s", report_path)

    # PDF export via LibreOffice (best effort)
    pdf_path = ""
    if pptx_path:
        try:
            from template_cataloger import _run, LIBREOFFICE_PATH
            result = _run([
                LIBREOFFICE_PATH, "--headless",
                "--convert-to", "pdf:impress_pdf_Export",
                "--outdir", str(out_dir),
                pptx_path,
            ])
            candidate = out_dir / (Path(pptx_path).stem + ".pdf")
            if candidate.exists():
                pdf_path = str(candidate)
                logger.info("finalize: PDF → %s", pdf_path)
        except Exception as exc:
            logger.warning("finalize: PDF export failed: %s", exc)

    return {"errors": errors}


# ═══════════════════════════════════════════════════════════════════════════════
# GRAPH — with validate → repair loop
# ═══════════════════════════════════════════════════════════════════════════════
#
#  load → map_shapes → polish_mappings → fetch_assets → build_pptx → render → validate ─┐
#                                                                                          │
#    ┌─────────────────────────────────────────────────────────────────────────────────────┘
#    │
#    ├─ (all pass OR max repairs OR no validation) → finalize → END
#    │
#    └─ (failures AND repair_round < MAX_REPAIR)
#         → repair → build_pptx → render → validate → (loop)
#


def _route_after_validate(state: BuilderState) -> str:
    """Conditional edge: decide whether to repair failed slides or finalize."""
    if not state.get("do_validate"):
        return "finalize"

    failed = state.get("failed_slide_indices", [])
    repair_round = state.get("repair_round", 0)

    if not failed:
        return "finalize"

    if repair_round >= MAX_REPAIR:
        logger.warning("route: %d slides still failing after %d repair rounds — finalizing anyway",
                        len(failed), repair_round)
        return "finalize"

    logger.info("route: %d slides failed → entering repair round %d/%d",
                len(failed), repair_round + 1, MAX_REPAIR)
    return "repair"


def build_graph() -> Any:
    g = StateGraph(BuilderState)
    g.add_node("plan_content", plan_content_node)
    g.add_node("select_layouts", select_layouts_node)
    g.add_node("load",         load_node)
    g.add_node("map_shapes",   map_shapes_node)
    g.add_node("polish_mappings", polish_mappings_node)
    g.add_node("fetch_assets", fetch_assets_node)
    g.add_node("build_pptx",   build_pptx_node)
    g.add_node("render",       render_node)
    g.add_node("validate",     validate_node)
    g.add_node("repair",       repair_node)
    g.add_node("finalize",     finalize_node)

    # Main pipeline
    g.add_edge(START,           "plan_content")
    g.add_edge("plan_content",  "select_layouts")
    g.add_edge("select_layouts", "load")
    g.add_edge("load",          "map_shapes")
    g.add_edge("map_shapes",    "polish_mappings")
    g.add_edge("polish_mappings", "fetch_assets")
    g.add_edge("fetch_assets",  "build_pptx")
    g.add_edge("build_pptx",    "render")
    g.add_edge("render",        "validate")

    # Conditional: validate → repair loop or finalize
    g.add_conditional_edges("validate", _route_after_validate, {
        "repair":   "repair",
        "finalize": "finalize",
    })

    # Repair feeds back into build → render → validate
    g.add_edge("repair",   "build_pptx")

    g.add_edge("finalize", END)

    return g.compile()


def _resolve_catalog_dir(template_path: Path) -> str:
    """Find the catalog directory for a template (for thumbnails)."""
    catalogs_base = _HERE / os.getenv("CATALOGS_DIR", "catalogs")
    cat_dir = catalogs_base / template_path.stem
    if cat_dir.exists():
        return str(cat_dir)
    return ""


def run_builder(
    build_spec: dict,
    template_path: str | Path,
    output_dir: str | Path,
    *,
    validate: bool = True,
) -> dict:
    """Run the full build pipeline. Returns the final state."""
    graph = build_graph()
    tpl = Path(template_path).resolve()
    initial: BuilderState = {
        "description":          "",
        "max_slides":           None,
        "content_plan":         {},
        "build_spec":           build_spec,
        "template_catalog":     {},
        "template_path":        str(tpl),
        "output_dir":           str(Path(output_dir).resolve()),
        "do_validate":          validate,
        "catalog_dir":          _resolve_catalog_dir(tpl),
        "slide_mappings":       [],
        "asset_requests":       [],
        "asset_paths":          {},
        "pptx_path":            "",
        "slide_pngs":           [],
        "build_issues":         [],
        "validation_results":   [],
        "repair_round":         0,
        "failed_slide_indices": [],
        "errors":               [],
    }
    result = graph.invoke(initial)
    return result


def run_deck_builder(
    description: str,
    template_path: str | Path,
    output_dir: str | Path,
    *,
    template_catalog: dict,
    max_slides: int | None = None,
    validate: bool = True,
) -> dict:
    """Run planning, layout selection, and slide building inside one graph."""
    graph = build_graph()
    tpl = Path(template_path).resolve()
    initial: BuilderState = {
        "description":          description,
        "max_slides":           max_slides,
        "content_plan":         {},
        "build_spec":           {},
        "template_catalog":     template_catalog or {},
        "template_path":        str(tpl),
        "output_dir":           str(Path(output_dir).resolve()),
        "do_validate":          validate,
        "catalog_dir":          _resolve_catalog_dir(tpl),
        "slide_mappings":       [],
        "asset_requests":       [],
        "asset_paths":          {},
        "pptx_path":            "",
        "slide_pngs":           [],
        "build_issues":         [],
        "validation_results":   [],
        "repair_round":         0,
        "failed_slide_indices": [],
        "errors":               [],
    }
    result = graph.invoke(initial)
    return result


# ═══════════════════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════════════════

def _cli() -> None:
    parser = argparse.ArgumentParser(
        description="Build a PPTX presentation from a build spec and template."
    )
    parser.add_argument("build_spec", type=Path, help="Path to build_spec.json.")
    parser.add_argument("template",   type=Path, help="Path to the .pptx template.")
    parser.add_argument("--out", "-o", type=Path, default=Path("output"),
                        help="Output directory (default: ./output).")
    parser.add_argument("--no-validate", action="store_true",
                        help="Skip visual validation and repair.")
    parser.add_argument("--verbose", "-v", action="store_true")
    args = parser.parse_args()

    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(levelname)s  %(message)s",
    )

    if not args.build_spec.exists():
        print(f"Error: {args.build_spec} not found", file=sys.stderr)
        sys.exit(1)
    if not args.template.exists():
        print(f"Error: {args.template} not found", file=sys.stderr)
        sys.exit(1)

    spec = json.loads(args.build_spec.read_text(encoding="utf-8"))

    result = run_builder(spec, args.template, args.out, validate=not args.no_validate)

    print(f"\n{'=' * 60}")
    print(f"  Deck: {spec.get('deck_title', '?')}")
    print(f"  PPTX: {result.get('pptx_path', 'n/a')}")
    pngs = result.get("slide_pngs", [])
    if pngs:
        print(f"  PNGs: {len(pngs)} slides → {Path(pngs[0]).parent}")
    val = result.get("validation_results", [])
    if val:
        passed = sum(1 for v in val if v.get("verdict") == "pass")
        print(f"  Validation: {passed}/{len(val)} passed")
    errs = result.get("errors", [])
    if errs:
        print(f"  Errors: {len(errs)}")
        for e in errs[:5]:
            print(f"    - {e}")
    print(f"{'=' * 60}\n")


if __name__ == "__main__":
    _cli()
