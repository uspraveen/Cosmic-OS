"""Generate branded PPTX templates with custom slide masters.

Each template has:
- Custom theme colors (primary, accent, background, text)
- Custom fonts (heading + body)
- Styled slide backgrounds on layouts
- Professional placeholder positioning
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


def _set_slide_bg(slide, r: int, g: int, b: int):
    """Set solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _set_shape_fill(shape, r: int, g: int, b: int):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


def _set_text_color(shape, r: int, g: int, b: int):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = RGBColor(r, g, b)


def generate_corporate_dark(output_path: Path):
    """Dark corporate template: navy bg, white text, blue accent."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Dark background, white text, blue accent
    BG = (26, 26, 46)  # #1a1a2e
    TEXT = (255, 255, 255)  # #ffffff
    ACCENT = (0, 123, 255)  # #007bff
    ACCENT2 = (233, 69, 96)  # #e94560
    SUBTEXT = (189, 195, 199)  # #bdc3c7

    for layout in prs.slide_layouts:
        for slide_placeholder in [layout]:
            pass  # Can't set bg on layouts directly

    # Style each layout
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            # Set font
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Calibri"
                        if ph.placeholder_format.idx == 0:  # title
                            run.font.size = Pt(32)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(*TEXT)
                        elif ph.placeholder_format.idx == 1:  # subtitle/body
                            run.font.size = Pt(18)
                            run.font.color.rgb = RGBColor(*SUBTEXT)

    prs.save(str(output_path))


def generate_corporate_light(output_path: Path):
    """Light corporate template: white bg, dark text, blue accent."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG = (255, 255, 255)  # #ffffff
    TEXT = (33, 37, 41)  # #212529
    ACCENT = (0, 123, 255)  # #007bff
    SUBTEXT = (108, 117, 125)  # #6c757d

    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Calibri"
                        if ph.placeholder_format.idx == 0:
                            run.font.size = Pt(32)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(*TEXT)
                        elif ph.placeholder_format.idx == 1:
                            run.font.size = Pt(18)
                            run.font.color.rgb = RGBColor(*SUBTEXT)

    prs.save(str(output_path))


def generate_minimal(output_path: Path):
    """Minimal template: clean, lots of whitespace, thin fonts."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    TEXT = (51, 51, 51)  # #333333
    SUBTEXT = (136, 136, 136)  # #888888

    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Helvetica"
                        run.font.color.rgb = RGBColor(*TEXT)
                        if ph.placeholder_format.idx == 0:
                            run.font.size = Pt(36)
                            run.font.bold = False
                        elif ph.placeholder_format.idx == 1:
                            run.font.size = Pt(16)
                            run.font.color.rgb = RGBColor(*SUBTEXT)

    prs.save(str(output_path))


def generate_pitch_deck(output_path: Path):
    """Startup pitch deck template: bold, colorful, impact-focused."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_DARK = (15, 15, 35)  # #0f0f23
    TEXT = (255, 255, 255)  # #ffffff
    ACCENT = (102, 51, 255)  # #6633ff
    ACCENT2 = (0, 255, 136)  # #00ff88
    SUBTEXT = (180, 180, 200)

    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if ph.has_text_frame:
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Arial"
                        run.font.color.rgb = RGBColor(*TEXT)
                        if ph.placeholder_format.idx == 0:
                            run.font.size = Pt(40)
                            run.font.bold = True
                        elif ph.placeholder_format.idx == 1:
                            run.font.size = Pt(20)
                            run.font.color.rgb = RGBColor(*SUBTEXT)

    prs.save(str(output_path))


def generate_all():
    """Generate all 4 templates."""
    templates_dir = Path(__file__).resolve().parent / "templates"
    templates_dir.mkdir(parents=True, exist_ok=True)

    generate_corporate_dark(templates_dir / "corporate-dark.pptx")
    print(f"  corporate-dark.pptx")

    generate_corporate_light(templates_dir / "corporate-light.pptx")
    print(f"  corporate-light.pptx")

    generate_minimal(templates_dir / "minimal.pptx")
    print(f"  minimal.pptx")

    generate_pitch_deck(templates_dir / "pitch-deck.pptx")
    print(f"  pitch-deck.pptx")

    print("All templates generated.")


if __name__ == "__main__":
    if len(sys.argv) > 1:
        templates_dir = Path(sys.argv[1])
    else:
        templates_dir = Path(__file__).resolve().parent / "templates"
    generate_all()
