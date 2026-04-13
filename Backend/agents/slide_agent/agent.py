"""COSMIC adapter for the copied cosmic-slides-2 implementation."""

from __future__ import annotations

import asyncio
import hashlib
import json
import logging
import mimetypes
import re
import shutil
import sys
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from uuid import uuid4

from shared.agent_runtime import AgentRuntime
from shared.contracts import AgentError, AgentResult, ArtifactManifest, TaskEnvelope, TaskInProgress, utcnow
from shared.sqlite_client import connect_sync

from .config import AGENT_ROOT, BACKEND_ROOT, SlideAgentConfig
from .source_retriever import SlideSourceCollection


logger = logging.getLogger(__name__)

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_DOCUMENT_EXTENSIONS = {".pdf", ".doc", ".docx", ".ppt", ".pptx"}
_DOCUMENT_MIMES = {
    "application/pdf",
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.ms-powerpoint",
    _PPTX_MIME,
}

_SLIDE_SESSIONS_SQL = """
CREATE TABLE IF NOT EXISTS slide_sessions (
    session_id TEXT,
    task_id TEXT,
    intent TEXT NOT NULL,
    workflow TEXT,
    editable INTEGER NOT NULL DEFAULT 0,
    title TEXT,
    template_name TEXT,
    pptx_path TEXT,
    slide_count INTEGER NOT NULL DEFAULT 0,
    created_at TEXT NOT NULL,
    PRIMARY KEY (session_id, task_id)
);
CREATE INDEX IF NOT EXISTS idx_slide_sessions_session_created
ON slide_sessions (session_id, created_at DESC);
"""

_ROLE_TO_ARCHETYPES: dict[str, set[str]] = {
    "opening": {"cover"},
    "agenda": {"title-body", "two-column"},
    "section_break": {"section-break", "cover"},
    "narrative": {"title-body", "two-column"},
    "data_story": {"chart-focus", "big-stat", "two-column"},
    "comparison": {"comparison", "two-column"},
    "highlight": {"big-stat", "quote", "cover"},
    "timeline": {"timeline"},
    "steps": {"timeline", "title-body", "two-column"},
    "visual": {"full-bleed-image", "cover"},
    "people": {"people-showcase", "two-column"},
    "closing": {"closing", "cover"},
}


@dataclass(frozen=True)
class _DeckPart:
    index: int
    total_parts: int
    start_slide: int
    end_slide: int
    max_slides: int | None
    description: str
    output_dir: Path
    content_plan: dict[str, Any] | None = None


class SlideAgent(AgentRuntime):
    """Slide specialist wrapper that preserves the standalone slides core."""

    def __init__(self, redis_client, config: SlideAgentConfig | None = None):
        self.config = config or SlideAgentConfig.from_env()
        self.config.apply_to_environment()
        if str(AGENT_ROOT) not in sys.path:
            sys.path.insert(0, str(AGENT_ROOT))

        super().__init__(
            agent_card_path=AGENT_ROOT / "agent_card.yaml",
            redis_client=redis_client,
            gateway_url=self.config.gateway_url,
            gateway_internal_token=self.config.gateway_internal_token,
        )
        self.artifacts_root = self.config.artifacts_root.resolve()
        self.db = None

    async def on_startup(self) -> None:
        for path in (
            self.config.templates_dir,
            self.config.catalogs_dir,
            self.config.assets_cache_dir,
            AGENT_ROOT / "runtime" / "cache",
            AGENT_ROOT / "runtime" / "logs",
            AGENT_ROOT / "store" / "data",
        ):
            path.mkdir(parents=True, exist_ok=True)

        self.db = connect_sync(AGENT_ROOT / "store" / "data" / "slide_sessions.db")
        self.db.executescript(_SLIDE_SESSIONS_SQL)
        self.db.commit()

    async def execute(self, task: TaskEnvelope) -> AgentResult | TaskInProgress:
        handler = getattr(self, f"handle_{task.intent.replace('.', '_')}", None)
        if handler is None:
            return self._failed(
                "INVALID_INPUT",
                f"Unknown intent: {task.intent}",
                next_action="escalate",
            )
        try:
            return await handler(task)
        except Exception as exc:
            logger.exception("slide_agent.task_failed task_id=%s intent=%s", task.task_id, task.intent)
            return self._failed(
                "INTERNAL_ERROR",
                str(exc).strip()[:500] or "Slide Agent execution failed.",
                retryable=False,
                next_action="escalate",
            )

    async def handle_slide_create(self, task: TaskEnvelope) -> AgentResult | TaskInProgress:
        description = self._text(task.input.get("description")) or self._text(task.input.get("query"))
        if not description:
            return self._failed("INVALID_INPUT", "slide.create requires a description.")

        requested_slides = self._requested_slide_count(task.input)
        validate = self._bool(task.input.get("validate"), self.config.validate_outputs)
        force_catalog = self._bool(task.input.get("force_catalog"), self.config.force_catalog_default)
        workflow = self._normalize_workflow(task.input.get("workflow") or task.input.get("mode"))
        explicit_template = self._resolve_template_from_input(task, persist_uploaded=True)

        if not workflow:
            workflow = self._normalize_workflow(self.config.default_workflow)
        if workflow == "auto":
            workflow = "template" if explicit_template else "html"
        if not workflow:
            return self._failed(
                "NEEDS_WORKFLOW_CHOICE",
                (
                    "Choose a slide workflow first: 'html' is fast and image-backed/non-editable; "
                    "'template' is slower and creates an editable PPTX from a template."
                ),
                retryable=False,
                next_action="ask_user",
            )

        source_context = await self._prepare_source_materials_for_generation(
            task,
            exclude_pptx_artifacts=explicit_template is not None and workflow == "template",
        )
        if isinstance(source_context, (AgentResult, TaskInProgress)):
            return source_context

        output_dir = self._task_artifact_dir(task.task_id)
        output_dir.mkdir(parents=True, exist_ok=True)
        parts = self._deck_parts_for_generation(
            description=description,
            output_dir=output_dir,
            source_context=source_context,
            requested_slides=requested_slides,
        )

        if workflow == "html":
            results = []
            for part in parts:
                results.append(
                    await asyncio.to_thread(
                        self._run_html,
                        part.description,
                        part.output_dir,
                        part.max_slides,
                        validate,
                        part.content_plan,
                    )
                )
            if len(results) == 1:
                return self._result_from_pipeline(
                    task=task,
                    result=results[0],
                    workflow="html",
                    editable=False,
                    template_path=None,
                )
            return self._result_from_split_pipeline(
                task=task,
                results=results,
                parts=parts,
                workflow="html",
                editable=False,
                template_path=None,
            )

        if workflow != "template":
            return self._failed("INVALID_INPUT", f"Unknown slide workflow: {workflow}")

        selected_template = explicit_template
        if selected_template is None:
            first_part = parts[0]
            if first_part.content_plan is not None:
                selected_template = await asyncio.to_thread(
                    self._select_template_for_plan,
                    first_part.content_plan,
                    force_catalog,
                )
            else:
                selected_template, planned_content = await asyncio.to_thread(
                    self._select_template_for_request,
                    first_part.description,
                    first_part.max_slides,
                    force_catalog,
                )
                if len(parts) == 1:
                    parts[0] = _DeckPart(
                        index=first_part.index,
                        total_parts=first_part.total_parts,
                        start_slide=first_part.start_slide,
                        end_slide=first_part.end_slide,
                        max_slides=first_part.max_slides,
                        description=first_part.description,
                        output_dir=first_part.output_dir,
                        content_plan=planned_content,
                    )

        results = []
        for part in parts:
            results.append(
                await asyncio.to_thread(
                    self._run_template,
                    part.description,
                    selected_template,
                    part.output_dir,
                    part.max_slides,
                    validate,
                    force_catalog,
                    part.content_plan,
                )
            )
        if len(results) == 1:
            return self._result_from_pipeline(
                task=task,
                result=results[0],
                workflow="template",
                editable=True,
                template_path=selected_template,
            )
        return self._result_from_split_pipeline(
            task=task,
            results=results,
            parts=parts,
            workflow="template",
            editable=True,
            template_path=selected_template,
        )

    async def handle_slide_edit(self, task: TaskEnvelope) -> AgentResult | TaskInProgress:
        edit_request = self._text(task.input.get("edit_request")) or self._text(task.input.get("description"))
        if not edit_request:
            return self._failed("INVALID_INPUT", "slide.edit requires edit_request.")

        await self._emit_progress(task.task_id, "Preparing the slide edit.")
        source_deck = self._resolve_source_deck(task)
        if source_deck is None:
            return self._failed(
                "MISSING_ARTIFACT",
                "slide.edit requires a source PPTX path or uploaded PPTX artifact.",
                next_action="ask_user",
            )

        output_dir = self._task_artifact_dir(task.task_id)
        output_dir.mkdir(parents=True, exist_ok=True)
        max_slides = self._max_slides(task.input) or self._count_pptx_slides(source_deck)
        validate = self._bool(task.input.get("validate"), self.config.validate_outputs)
        force_catalog = self._bool(task.input.get("force_catalog"), self.config.force_catalog_default)
        source_context = await self._prepare_source_materials_for_generation(
            task,
            exclude_pptx_artifacts=True,
        )
        if isinstance(source_context, (AgentResult, TaskInProgress)):
            return source_context
        description = (
            "Revise this existing presentation using the provided deck as the editable template backbone.\n\n"
            f"Edit request:\n{edit_request}"
        )
        description = self._augment_description_with_source_context(description, source_context)

        source_profile = self._source_deck_profile(source_deck)
        if self._should_regenerate_html_for_edit(source_deck, source_profile):
            html_description = self._html_edit_regeneration_description(
                source_deck=source_deck,
                edit_request=edit_request,
                source_profile=source_profile,
            )
            html_description = self._augment_description_with_source_context(html_description, source_context)
            result = await self._run_blocking_stage(
                task,
                "Regenerating the non-editable HTML deck with the requested changes.",
                self._run_html,
                html_description,
                output_dir,
                max_slides,
                validate,
                None,
            )
            await self._emit_progress(task.task_id, "Packaging the regenerated HTML deck.")
            return self._result_from_pipeline(
                task=task,
                result=result,
                workflow="html",
                editable=False,
                template_path=None,
                edit_mode="html_regeneration_from_non_editable_source",
            )

        result = await self._run_blocking_stage(
            task,
            "Running template-backed deck regeneration.",
            self._run_template,
            description,
            source_deck,
            output_dir,
            max_slides,
            validate,
            force_catalog,
            None,
        )
        await self._emit_progress(task.task_id, "Packaging the regenerated template deck.")
        return self._result_from_pipeline(
            task=task,
            result=result,
            workflow="template",
            editable=True,
            template_path=source_deck,
            edit_mode="template_backed_regeneration",
        )

    async def handle_slide_catalog_template(self, task: TaskEnvelope) -> AgentResult:
        template_path = self._resolve_template_from_input(task, persist_uploaded=True)
        if template_path is None:
            return self._failed(
                "MISSING_ARTIFACT",
                "slide.catalog_template requires a template_path or uploaded PPTX artifact.",
                next_action="ask_user",
            )

        force = self._bool(task.input.get("force"), self._bool(task.input.get("force_catalog"), False))
        catalog = await asyncio.to_thread(self._ensure_catalog, template_path, force)
        artifacts = self._copy_catalog_artifacts(task.task_id, template_path, catalog)
        output = {
            "response": f"Cataloged template {template_path.name}.",
            "template_name": template_path.stem,
            "template_path": str(template_path),
            "catalog_path": str(self._catalog_path(template_path)),
            "slide_count": int(catalog.get("slide_count") or len(catalog.get("slides") or [])),
            "force": force,
        }
        return AgentResult(status="completed", output=output, artifacts=artifacts, error=None)

    async def handle_slide_recall_session(self, task: TaskEnvelope) -> AgentResult:
        session_id = self._text(task.input.get("session_id")) or self._text(task.session_id)
        limit = max(1, min(50, self._int(task.input.get("limit"), 10)))
        if self.db is None:
            return AgentResult(status="completed", output={"history": [], "count": 0}, artifacts=[], error=None)
        rows = self.db.execute(
            """SELECT * FROM slide_sessions
               WHERE session_id = ?
               ORDER BY created_at DESC LIMIT ?""",
            [session_id, limit],
        ).fetchall()
        history = [dict(row) for row in rows]
        return AgentResult(status="completed", output={"history": history, "count": len(history)}, artifacts=[], error=None)

    async def _prepare_source_materials_for_generation(
        self,
        task: TaskEnvelope,
        *,
        exclude_pptx_artifacts: bool,
    ) -> dict[str, Any] | AgentResult | TaskInProgress:
        """Parse uploaded document artifacts and build compact context for the slides core."""
        resume_block = task.input.get("_resume") if isinstance(task.input.get("_resume"), dict) else {}
        resume_state = (
            resume_block.get("resume_state")
            if isinstance(resume_block.get("resume_state"), dict)
            else {}
        )
        reverse_result = (
            resume_block.get("reverse_result")
            if isinstance(resume_block.get("reverse_result"), dict)
            else {}
        )
        pending_request = (
            resume_state.get("pending_asset_request")
            if isinstance(resume_state.get("pending_asset_request"), dict)
            else {}
        )

        source_artifacts = self._merge_artifact_descriptors(
            self._input_artifacts(task),
            resume_state.get("source_artifacts"),
        )
        source_documents = self._merge_document_summaries(
            task.input.get("source_documents"),
            resume_state.get("source_documents"),
            self._document_summaries_from_artifacts(source_artifacts),
        )

        if reverse_result and str(pending_request.get("request_kind") or "").strip().lower() == "docs_parse":
            parse_result = self._apply_docs_parse_reverse_result(
                source_artifacts=source_artifacts,
                source_documents=source_documents,
                reverse_result=reverse_result,
            )
            if isinstance(parse_result, AgentResult):
                return parse_result
            source_artifacts, source_documents = parse_result

        parsed_source_artifact_ids = {
            str(item.get("artifact_id") or "").strip()
            for item in source_documents
            if str(item.get("artifact_id") or "").strip()
        }
        raw_document_artifacts = [
            artifact
            for artifact in source_artifacts
            if self._is_document_artifact(artifact)
            and not self._is_docs_parser_artifact(artifact)
            and not (exclude_pptx_artifacts and self._is_pptx_artifact(artifact))
            and str(artifact.get("artifact_id") or "").strip() not in parsed_source_artifact_ids
        ]

        if raw_document_artifacts:
            return await self._request_docs_parse(task, raw_document_artifacts, source_artifacts, source_documents)

        parsed_documents = self._merge_document_summaries(
            source_documents,
            self._document_summaries_from_artifacts(source_artifacts),
        )
        return {
            "source_artifacts": source_artifacts,
            "documents": parsed_documents,
            "visual_assets": self._visual_asset_summaries_from_artifacts(source_artifacts),
        }

    async def _request_docs_parse(
        self,
        task: TaskEnvelope,
        raw_document_artifacts: list[dict[str, Any]],
        source_artifacts: list[dict[str, Any]],
        source_documents: list[dict[str, Any]],
    ) -> TaskInProgress | AgentResult:
        pending_asset = {
            "request_kind": "docs_parse",
            "target_intent": "docs.parse_bundle",
            "target_agent_id": self.config.docs_parser_agent_id,
            "source_artifact_ids": [
                str(item.get("artifact_id") or "").strip()
                for item in raw_document_artifacts
                if str(item.get("artifact_id") or "").strip()
            ],
        }
        resume_payload = {
            "source_artifacts": source_artifacts,
            "source_documents": source_documents,
            "pending_asset_request": pending_asset,
        }
        label = (
            self._text(task.input.get("bundle_label"))
            or self._text(task.input.get("description") or task.input.get("query"))[:120]
            or f"slide_agent_{task.task_id}"
        )
        try:
            result = await self.request_orchestrator_delegate(
                current_task=task,
                target_intent="docs.parse_bundle",
                target_input={
                    "bundle_label": label,
                    "generate_page_images": True,
                    "generate_picture_images": True,
                },
                target_agent_id=self.config.docs_parser_agent_id,
                input_artifacts=raw_document_artifacts,
                resume_payload=resume_payload,
                reason="slide_agent needs parsed document bundles before slide generation",
            )
            reverse_task_id = self._text(result.get("reverse_task_id"))
            if not reverse_task_id:
                return self._failed(
                    "DOCS_PARSE_DELEGATION_FAILED",
                    "Docs parse delegation did not return a reverse_task_id.",
                    next_action="retry",
                    retryable=True,
                )
            await self._emit_suspension_event(
                task,
                reverse_task_id=reverse_task_id,
                source_artifact_count=len(raw_document_artifacts),
            )
        except Exception as exc:
            return self._failed(
                "DOCS_PARSE_DELEGATION_FAILED",
                str(exc).strip()[:500] or "Docs parse delegation failed.",
                retryable=True,
                next_action="retry",
            )
        return TaskInProgress(
            task_id=task.task_id,
            idempotency_key=task.idempotency_key,
            executing_since=utcnow(),
            check_after_sec=10,
        )

    async def _emit_suspension_event(
        self,
        task: TaskEnvelope,
        *,
        reverse_task_id: str,
        source_artifact_count: int,
    ) -> None:
        try:
            await self.emit_event(
                task.task_id,
                "task.suspended",
                {
                    "reason": "slide_prepare_source_materials",
                    "reverse_task_id": reverse_task_id,
                    "target_intent": "docs.parse_bundle",
                    "target_agent_id": self.config.docs_parser_agent_id,
                    "resume_intent": "agent.resume",
                    "source_artifact_count": source_artifact_count,
                },
            )
        except Exception as exc:
            logger.warning("slide_agent.suspension_event_failed task_id=%s error=%s", task.task_id, exc)

    async def _emit_progress(self, task_id: str, message: str, **payload: Any) -> None:
        body = {"message": message, **payload}
        try:
            await self.emit_event(task_id, "task.progress", body)
        except Exception as exc:
            logger.warning("slide_agent.progress_emit_failed task_id=%s error=%s", task_id, exc)

    async def _run_blocking_stage(
        self,
        task: TaskEnvelope,
        message: str,
        func: Any,
        *args: Any,
    ) -> dict[str, Any]:
        await self._emit_progress(task.task_id, message)
        worker = asyncio.create_task(asyncio.to_thread(func, *args))
        elapsed_sec = 0
        heartbeat_sec = 20
        while not worker.done():
            done, _pending = await asyncio.wait({worker}, timeout=heartbeat_sec)
            if done:
                break
            elapsed_sec += heartbeat_sec
            await self._emit_progress(
                task.task_id,
                f"{message} Still working after {elapsed_sec}s.",
                stage="running",
                elapsed_sec=elapsed_sec,
            )
        result = await worker
        return result if isinstance(result, dict) else {}

    def _apply_docs_parse_reverse_result(
        self,
        *,
        source_artifacts: list[dict[str, Any]],
        source_documents: list[dict[str, Any]],
        reverse_result: dict[str, Any],
    ) -> tuple[list[dict[str, Any]], list[dict[str, Any]]] | AgentResult:
        status = self._text(reverse_result.get("status")).lower()
        if status and status != "completed":
            err = reverse_result.get("error") if isinstance(reverse_result.get("error"), dict) else {}
            return self._failed(
                self._text(err.get("code")) or "DOCS_PARSE_FAILED",
                self._text(err.get("message")) or "Delegated docs parsing failed.",
                next_action="retry",
                retryable=True,
            )
        merged_artifacts = self._merge_artifact_descriptors(
            source_artifacts,
            reverse_result.get("artifacts"),
        )
        output = reverse_result.get("output") if isinstance(reverse_result.get("output"), dict) else {}
        merged_documents = self._merge_document_summaries(
            source_documents,
            self._enrich_document_summaries(output.get("documents")),
            self._document_summaries_from_artifacts(merged_artifacts),
        )
        return merged_artifacts, merged_documents

    def _augment_description_with_source_context(self, description: str, source_context: dict[str, Any]) -> str:
        documents = source_context.get("documents") if isinstance(source_context.get("documents"), list) else []
        visual_assets = source_context.get("visual_assets") if isinstance(source_context.get("visual_assets"), list) else []
        if not documents and not visual_assets:
            return description

        lines = [
            description.rstrip(),
            "",
            "Source document context:",
            "Use the following parsed source material as the factual basis for the deck. Preserve important names, figures, dates, and section structure; do not invent facts not supported by this context.",
        ]
        for index, doc in enumerate(documents[:5], start=1):
            if not isinstance(doc, dict):
                continue
            title = self._text(doc.get("title")) or self._text(doc.get("filename")) or f"Document {index}"
            lines.append("")
            lines.append(f"[Document {index}] {title}")
            meta_parts = []
            for label, key in (
                ("filename", "filename"),
                ("pages", "page_count"),
                ("slides", "slide_count"),
                ("sections", "section_count"),
                ("tables", "table_count"),
                ("figures", "figure_count"),
            ):
                value = doc.get(key)
                if value not in {None, "", 0}:
                    meta_parts.append(f"{label}: {value}")
            if meta_parts:
                lines.append("Metadata: " + "; ".join(meta_parts))
            preview = self._compact_text(doc.get("preview_excerpt"), limit=1800)
            if preview:
                lines.append("Preview excerpt: " + preview)
            sections = doc.get("top_sections") if isinstance(doc.get("top_sections"), list) else []
            if sections:
                lines.append("Relevant sections:")
                for section in sections[:8]:
                    if not isinstance(section, dict):
                        continue
                    section_title = self._text(section.get("title")) or self._text(section.get("section_id")) or "Untitled section"
                    section_preview = self._compact_text(section.get("preview"), limit=260) or ""
                    if section_preview:
                        lines.append(f"- {section_title}: {section_preview}")
                    else:
                        lines.append(f"- {section_title}")

        if visual_assets:
            lines.append("")
            lines.append("Detected source visuals:")
            for asset in visual_assets[:10]:
                if not isinstance(asset, dict):
                    continue
                label = self._text(asset.get("label")) or self._text(asset.get("filename")) or self._text(asset.get("asset_ref")) or "visual"
                detail = self._compact_text(asset.get("description") or asset.get("caption"), limit=220)
                lines.append(f"- {label}" + (f": {detail}" if detail else ""))

        return "\n".join(lines)[:20000]

    def _deck_parts_for_generation(
        self,
        *,
        description: str,
        output_dir: Path,
        source_context: dict[str, Any],
        requested_slides: int | None,
    ) -> list[_DeckPart]:
        source_collection = self._source_collection(source_context)
        one_per_source_unit = (
            self._is_one_slide_per_source_unit_request(description)
            and source_collection is not None
            and source_collection.unit_count > 0
        )
        effective_slides = requested_slides
        if one_per_source_unit and source_collection is not None:
            effective_slides = source_collection.unit_count if effective_slides is None else min(effective_slides, source_collection.unit_count)

        base_description = self._augment_description_with_source_context(description, source_context)
        deck_limit = max(1, self.config.max_slides_per_deck)
        if effective_slides is None:
            return [
                _DeckPart(
                    index=1,
                    total_parts=1,
                    start_slide=1,
                    end_slide=0,
                    max_slides=None,
                    description=base_description,
                    output_dir=output_dir,
                    content_plan=None,
                )
            ]

        parts: list[_DeckPart] = []
        total_parts = (effective_slides + deck_limit - 1) // deck_limit
        base_title = self._deck_title_from_description(description, source_context)
        for index, start_slide in enumerate(range(1, effective_slides + 1, deck_limit), start=1):
            end_slide = min(start_slide + deck_limit - 1, effective_slides)
            part_output_dir = output_dir if total_parts == 1 else output_dir / f"part-{index:03d}"
            content_plan = None
            window_prompt = ""
            if one_per_source_unit and source_collection is not None:
                window_prompt = source_collection.window_prompt(start_slide, end_slide)
                part_title = base_title if total_parts == 1 else f"{base_title} (Part {index} of {total_parts})"
                content_plan = source_collection.one_unit_per_slide_plan(
                    start_ordinal=start_slide,
                    end_ordinal=end_slide,
                    deck_title=part_title,
                    deck_theme=(
                        "One slide per source page/slide, preserving the source ordering "
                        "and using segmented source retrieval as the factual basis."
                    ),
                )
            part_description = self._part_description(
                base_description=base_description,
                part_index=index,
                total_parts=total_parts,
                start_slide=start_slide,
                end_slide=end_slide,
                total_slides=effective_slides,
                deck_limit=deck_limit,
                source_window=window_prompt,
                one_per_source_unit=one_per_source_unit,
            )
            parts.append(
                _DeckPart(
                    index=index,
                    total_parts=total_parts,
                    start_slide=start_slide,
                    end_slide=end_slide,
                    max_slides=end_slide - start_slide + 1,
                    description=part_description,
                    output_dir=part_output_dir,
                    content_plan=content_plan,
                )
            )
        return parts

    def _source_collection(self, source_context: dict[str, Any]) -> SlideSourceCollection | None:
        documents = source_context.get("documents") if isinstance(source_context.get("documents"), list) else []
        if not documents:
            return None
        collection = SlideSourceCollection(documents, self._resolve_path)
        return collection if collection.unit_count else None

    def _part_description(
        self,
        *,
        base_description: str,
        part_index: int,
        total_parts: int,
        start_slide: int,
        end_slide: int,
        total_slides: int,
        deck_limit: int,
        source_window: str,
        one_per_source_unit: bool,
    ) -> str:
        lines = [
            base_description.rstrip(),
            "",
            f"Deck part {part_index} of {total_parts}.",
            f"Create only slides {start_slide}-{end_slide} of {total_slides}.",
            f"Hard limit: each produced PPTX must contain at most {deck_limit} slides.",
        ]
        if one_per_source_unit:
            lines.append("This is a one-slide-per-source-page/slide request; preserve the source unit order exactly for this part.")
        elif total_parts > 1:
            lines.append("Do not duplicate earlier/later parts. Cover only this slice of the full requested deck.")
        if source_window:
            lines.extend(["", source_window])
        return "\n".join(lines)

    def _deck_title_from_description(self, description: str, source_context: dict[str, Any]) -> str:
        documents = source_context.get("documents") if isinstance(source_context.get("documents"), list) else []
        for doc in documents:
            if not isinstance(doc, dict):
                continue
            title = self._text(doc.get("title")) or self._text(doc.get("filename"))
            if title:
                return self._compact_text(title, limit=90)
        first_line = next((line.strip() for line in description.splitlines() if line.strip()), "")
        return self._compact_text(first_line or "Slide Deck", limit=90)

    @staticmethod
    def _is_one_slide_per_source_unit_request(description: str) -> bool:
        text = re.sub(r"\s+", " ", description or "").lower()
        patterns = (
            r"\bone\s+slide\s+per\s+(page|slide)\b",
            r"\b1\s+slide\s+per\s+(page|slide)\b",
            r"\b(each|every)\s+(page|slide)\s+(as|into|gets|should have)\s+(one\s+)?slide\b",
            r"\b(page|slide)[-\s]+by[-\s]+(page|slide)\b",
        )
        return any(re.search(pattern, text) for pattern in patterns)

    def _run_html(
        self,
        description: str,
        output_dir: Path,
        max_slides: int | None,
        validate: bool,
        content_plan: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        from html_workflow import run_html_pipeline

        result = run_html_pipeline(
            description,
            output_dir=output_dir,
            max_slides=max_slides,
            validate=validate,
            content_plan=content_plan,
        )
        result.setdefault("workflow", "html")
        return result

    def _run_template(
        self,
        description: str,
        template_path: Path,
        output_dir: Path,
        max_slides: int | None,
        validate: bool,
        force_catalog: bool,
        content_plan: dict[str, Any] | None,
    ) -> dict[str, Any]:
        from content_planner import plan_content
        from layout_selector import select_layouts
        from slide_builder import run_builder

        catalog = self._ensure_catalog(template_path, force_catalog)
        plan = content_plan or plan_content(description, num_slides=max_slides)
        self._write_json(output_dir / "plan.json", plan)

        build_spec = select_layouts(plan, catalog, max_slides=max_slides)
        build_spec.setdefault("deck_title", plan.get("deck_title") or template_path.stem)
        self._write_json(output_dir / "build_spec.json", build_spec)

        result = run_builder(
            build_spec,
            template_path,
            output_dir,
            validate=validate,
        )
        result.setdefault("workflow", "template")
        result["content_plan"] = result.get("content_plan") or plan
        result["build_spec"] = result.get("build_spec") or build_spec
        result["template_catalog"] = catalog
        return result

    def _select_template_for_request(
        self,
        description: str,
        max_slides: int | None,
        force_catalog: bool,
    ) -> tuple[Path, dict[str, Any]]:
        from content_planner import plan_content

        templates = sorted(self.config.templates_dir.glob("*.pptx"))
        if not templates:
            raise FileNotFoundError(f"No PPTX templates found in {self.config.templates_dir}")
        plan = plan_content(description, num_slides=max_slides)
        if len(templates) == 1:
            return templates[0], plan

        scored: list[tuple[float, Path]] = []
        for template in templates:
            catalog = self._ensure_catalog(template, force_catalog)
            scored.append((self._score_template(plan, catalog, template), template))
        scored.sort(key=lambda item: (item[0], item[1].name.lower()), reverse=True)
        logger.info("slide_agent.template_selected template=%s score=%.2f", scored[0][1], scored[0][0])
        return scored[0][1], plan

    def _select_template_for_plan(self, plan: dict[str, Any], force_catalog: bool) -> Path:
        templates = sorted(self.config.templates_dir.glob("*.pptx"))
        if not templates:
            raise FileNotFoundError(f"No PPTX templates found in {self.config.templates_dir}")
        if len(templates) == 1:
            return templates[0]
        scored: list[tuple[float, Path]] = []
        for template in templates:
            catalog = self._ensure_catalog(template, force_catalog)
            scored.append((self._score_template(plan, catalog, template), template))
        scored.sort(key=lambda item: (item[0], item[1].name.lower()), reverse=True)
        logger.info("slide_agent.template_selected template=%s score=%.2f", scored[0][1], scored[0][0])
        return scored[0][1]

    def _score_template(self, plan: dict[str, Any], catalog: dict[str, Any], template: Path) -> float:
        slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
        desired_count = len(slides)
        role_counts = Counter(
            str(slide.get("content_role") or "").strip().lower()
            for slide in slides
            if isinstance(slide, dict)
        )
        archetype_counts = Counter(
            str(slide.get("layout_archetype") or "").strip().lower()
            for slide in catalog.get("slides", [])
            if isinstance(slide, dict)
        )

        score = 0.0
        for role, count in role_counts.items():
            preferred = _ROLE_TO_ARCHETYPES.get(role) or set()
            for archetype, archetype_count in archetype_counts.items():
                if archetype in preferred:
                    score += 3.0 * count * archetype_count
                elif any(token and token in archetype for item in preferred for token in item.split("-")):
                    score += 0.75 * count * archetype_count

        slide_count = int(catalog.get("slide_count") or len(catalog.get("slides") or []) or 0)
        if desired_count and slide_count:
            score += max(0.0, 5.0 - abs(slide_count - desired_count))

        request_tokens = {
            token
            for token in re.findall(
                r"[a-z0-9]+",
                f"{plan.get('deck_title', '')} {plan.get('deck_theme', '')}".lower(),
            )
            if len(token) >= 4
        }
        template_tokens = set(re.findall(r"[a-z0-9]+", template.stem.lower()))
        score += 0.5 * len(request_tokens & template_tokens)
        return score

    def _ensure_catalog(self, template_path: Path, force: bool) -> dict[str, Any]:
        from template_cataloger import catalog_template, load_catalog

        template_path = template_path.resolve()
        catalog = None if force else load_catalog(template_path)
        if catalog is None:
            catalog = catalog_template(template_path, force=force)
        return self._localize_catalog_paths(template_path, catalog)

    def _localize_catalog_paths(self, template_path: Path, catalog: dict[str, Any]) -> dict[str, Any]:
        normalized = json.loads(json.dumps(catalog))
        catalog_dir = self._catalog_dir(template_path)
        normalized["template_name"] = template_path.stem
        normalized["template_path"] = str(template_path.resolve())
        normalized["collage_path"] = str((catalog_dir / "collage.png").resolve())
        for slide in normalized.get("slides", []) or []:
            if isinstance(slide, dict):
                slide_number = self._int(slide.get("slide_number"), 0)
                if slide_number:
                    slide["thumbnail_path"] = str((catalog_dir / "thumbnails" / f"slide-{slide_number:02d}.png").resolve())
        return normalized

    def _source_deck_profile(self, source_deck: Path) -> dict[str, Any]:
        profile: dict[str, Any] = {}
        report = self._safe_json_load(source_deck.parent / "build_report.json")
        if report:
            profile["workflow"] = self._text(report.get("workflow")).lower()
            profile["description"] = self._text(report.get("description"))
            profile["report"] = report
        plan = self._safe_json_load(source_deck.parent / "plan.json")
        if plan:
            profile["plan"] = plan
            profile["title"] = self._text(plan.get("deck_title"))

        db_profile = self._source_deck_profile_from_db(source_deck)
        for key, value in db_profile.items():
            if self._emptyish(profile.get(key)) and not self._emptyish(value):
                profile[key] = value
        return profile

    def _source_deck_profile_from_db(self, source_deck: Path) -> dict[str, Any]:
        if self.db is None:
            return {}
        try:
            source_resolved = source_deck.resolve()
            row = self.db.execute(
                """SELECT workflow, editable, title
                   FROM slide_sessions
                   WHERE pptx_path = ?
                   ORDER BY created_at DESC
                   LIMIT 1""",
                [str(source_resolved)],
            ).fetchone()
            if row is None:
                return {}
            return {
                "workflow": self._text(row["workflow"]).lower(),
                "editable": bool(row["editable"]),
                "title": self._text(row["title"]),
            }
        except Exception:
            return {}

    def _should_regenerate_html_for_edit(self, source_deck: Path, source_profile: dict[str, Any]) -> bool:
        workflow = self._text(source_profile.get("workflow")).lower()
        if workflow == "html":
            return True
        if source_profile.get("editable") is False:
            return True
        return self._looks_like_image_backed_deck(source_deck)

    def _html_edit_regeneration_description(
        self,
        *,
        source_deck: Path,
        edit_request: str,
        source_profile: dict[str, Any],
    ) -> str:
        parts = [
            "Regenerate the existing presentation as a new HTML-mode, image-backed deck.",
            "Preserve the prior deck's subject, slide count, and core content unless the edit request overrides them.",
            "",
            f"Source deck: {source_deck.name}",
        ]
        title = self._text(source_profile.get("title"))
        if title:
            parts.append(f"Previous deck title: {title}")
        previous_description = self._text(source_profile.get("description"))
        if previous_description:
            parts.extend(["", "Previous deck brief:", previous_description])
        previous_plan = source_profile.get("plan")
        if isinstance(previous_plan, dict) and previous_plan:
            parts.extend(
                [
                    "",
                    "Previous deck content plan JSON:",
                    self._compact_json(previous_plan, limit=12000),
                ]
            )
        parts.extend(["", "Requested edit:", edit_request])
        return "\n".join(parts)

    @staticmethod
    def _compact_json(value: Any, *, limit: int) -> str:
        text = json.dumps(value, ensure_ascii=False, indent=2, default=str)
        if limit <= 0 or len(text) <= limit:
            return text
        return text[: max(0, limit - 1)].rstrip() + "..."

    def _looks_like_image_backed_deck(self, source_deck: Path) -> bool:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            presentation = Presentation(str(source_deck))
            if not presentation.slides:
                return False
            for slide in presentation.slides:
                has_text = any(
                    bool(getattr(shape, "has_text_frame", False))
                    and bool(getattr(shape.text_frame, "text", "").strip())
                    for shape in slide.shapes
                )
                picture_count = sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
                if has_text or picture_count != 1:
                    return False
            return True
        except Exception:
            return False

    def _copy_catalog_artifacts(
        self,
        task_id: str,
        template_path: Path,
        catalog: dict[str, Any],
    ) -> list[ArtifactManifest]:
        output_dir = self._task_artifact_dir(task_id) / "catalog"
        output_dir.mkdir(parents=True, exist_ok=True)
        artifacts: list[ArtifactManifest] = []

        catalog_path = output_dir / f"{template_path.stem}.catalog.json"
        self._write_json(catalog_path, catalog)
        artifacts.append(
            self._artifact_manifest(
                task_id=task_id,
                path=catalog_path,
                mime="application/json",
                audience="supporting",
            )
        )

        collage_path = Path(str(catalog.get("collage_path") or ""))
        if collage_path.exists() and collage_path.is_file():
            copied = output_dir / f"{template_path.stem}.collage.png"
            shutil.copy2(collage_path, copied)
            artifacts.append(
                self._artifact_manifest(
                    task_id=task_id,
                    path=copied,
                    mime="image/png",
                    audience="supporting",
                )
            )
        return artifacts

    def _result_from_pipeline(
        self,
        *,
        task: TaskEnvelope,
        result: dict[str, Any],
        workflow: str,
        editable: bool,
        template_path: Path | None,
        edit_mode: str | None = None,
    ) -> AgentResult:
        artifacts: list[ArtifactManifest] = []
        seen: set[Path] = set()
        deliverable_candidates = [
            result.get("pptx_path"),
            result.get("pdf_path"),
            result.get("contact_sheet"),
        ]
        supporting_candidates = [
            result.get("plan_path"),
            result.get("theme_path"),
            self._task_artifact_dir(task.task_id) / "build_report.json",
            self._task_artifact_dir(task.task_id) / "build_spec.json",
            self._task_artifact_dir(task.task_id) / "plan.json",
        ]

        for raw in deliverable_candidates:
            path = self._maybe_path(raw)
            if path is None or path in seen:
                continue
            seen.add(path)
            artifacts.append(
                self._artifact_manifest(
                    task_id=task.task_id,
                    path=path,
                    mime=self._mime_for_path(path),
                    audience="deliverable",
                )
            )

        for raw in supporting_candidates:
            path = self._maybe_path(raw)
            if path is None or path in seen:
                continue
            seen.add(path)
            artifacts.append(
                self._artifact_manifest(
                    task_id=task.task_id,
                    path=path,
                    mime=self._mime_for_path(path),
                    audience="supporting",
                )
            )

        slide_count = self._slide_count_from_result(result)
        title = self._result_title(result, fallback=Path(str(result.get("pptx_path") or "deck.pptx")).stem)
        pptx_path = self._text(result.get("pptx_path"))
        pdf_path = self._text(result.get("pdf_path"))
        validation_results = result.get("validation_results") if isinstance(result.get("validation_results"), list) else []
        validation_issues = [
            item
            for item in validation_results
            if isinstance(item, dict) and item.get("verdict") not in {None, "pass"}
        ]
        output = {
            "response": self._summary_response(workflow=workflow, editable=editable, slide_count=slide_count, title=title),
            "workflow": workflow,
            "editable": editable,
            "pptx_path": pptx_path,
            "pdf_path": pdf_path,
            "slide_count": slide_count,
            "template": template_path.name if template_path else None,
            "template_path": str(template_path) if template_path else None,
            "title": title,
            "validation_pass": not validation_issues,
            "validation_issues": validation_issues,
            "edit_mode": edit_mode,
        }
        self._record_session(
            task=task,
            workflow=workflow,
            editable=editable,
            title=title,
            template_path=template_path,
            pptx_path=pptx_path,
            slide_count=slide_count,
        )
        return AgentResult(status="completed", output=output, artifacts=artifacts, error=None)

    def _result_from_split_pipeline(
        self,
        *,
        task: TaskEnvelope,
        results: list[dict[str, Any]],
        parts: list[_DeckPart],
        workflow: str,
        editable: bool,
        template_path: Path | None,
    ) -> AgentResult:
        artifacts: list[ArtifactManifest] = []
        seen: set[Path] = set()
        deck_parts: list[dict[str, Any]] = []
        validation_issues: list[dict[str, Any]] = []

        for part, result in zip(parts, results, strict=False):
            deliverable_candidates = [
                result.get("pptx_path"),
                result.get("pdf_path"),
                result.get("contact_sheet"),
            ]
            supporting_candidates = [
                result.get("plan_path"),
                result.get("theme_path"),
                part.output_dir / "build_report.json",
                part.output_dir / "build_spec.json",
                part.output_dir / "plan.json",
            ]
            for raw in deliverable_candidates:
                path = self._maybe_path(raw)
                if path is None or path in seen:
                    continue
                seen.add(path)
                artifacts.append(
                    self._artifact_manifest(
                        task_id=task.task_id,
                        path=path,
                        mime=self._mime_for_path(path),
                        audience="deliverable",
                    )
                )
            for raw in supporting_candidates:
                path = self._maybe_path(raw)
                if path is None or path in seen:
                    continue
                seen.add(path)
                artifacts.append(
                    self._artifact_manifest(
                        task_id=task.task_id,
                        path=path,
                        mime=self._mime_for_path(path),
                        audience="supporting",
                    )
                )

            part_slide_count = self._slide_count_from_result(result)
            part_validation = result.get("validation_results") if isinstance(result.get("validation_results"), list) else []
            part_issues = [
                {**item, "deck_part": part.index}
                for item in part_validation
                if isinstance(item, dict) and item.get("verdict") not in {None, "pass"}
            ]
            validation_issues.extend(part_issues)
            deck_parts.append(
                {
                    "part_index": part.index,
                    "total_parts": part.total_parts,
                    "start_slide": part.start_slide,
                    "end_slide": part.end_slide,
                    "slide_count": part_slide_count,
                    "pptx_path": self._text(result.get("pptx_path")),
                    "pdf_path": self._text(result.get("pdf_path")),
                    "contact_sheet": self._text(result.get("contact_sheet")),
                    "title": self._result_title(
                        result,
                        fallback=Path(str(result.get("pptx_path") or f"deck_part_{part.index}.pptx")).stem,
                    ),
                    "validation_pass": not part_issues,
                }
            )

        total_slide_count = sum(self._int(item.get("slide_count"), 0) for item in deck_parts)
        first_part = deck_parts[0] if deck_parts else {}
        title = self._text(first_part.get("title")) or "slide deck"
        pptx_path = self._text(first_part.get("pptx_path"))
        pdf_path = self._text(first_part.get("pdf_path"))
        output = {
            "response": self._summary_split_response(
                workflow=workflow,
                editable=editable,
                slide_count=total_slide_count,
                part_count=len(deck_parts),
                title=title,
            ),
            "workflow": workflow,
            "editable": editable,
            "pptx_path": pptx_path,
            "pdf_path": pdf_path,
            "slide_count": total_slide_count,
            "split": True,
            "max_slides_per_deck": self.config.max_slides_per_deck,
            "deck_parts": deck_parts,
            "template": template_path.name if template_path else None,
            "template_path": str(template_path) if template_path else None,
            "title": title,
            "validation_pass": not validation_issues,
            "validation_issues": validation_issues,
            "edit_mode": None,
        }
        self._record_session(
            task=task,
            workflow=workflow,
            editable=editable,
            title=title,
            template_path=template_path,
            pptx_path=pptx_path,
            slide_count=total_slide_count,
        )
        return AgentResult(status="completed", output=output, artifacts=artifacts, error=None)

    def _summary_response(self, *, workflow: str, editable: bool, slide_count: int, title: str) -> str:
        editability = "editable PPTX" if editable else "image-backed non-editable PPTX"
        target = title or "slide deck"
        return f"Created {target} with {slide_count} slides via {workflow} workflow ({editability})."

    def _summary_split_response(
        self,
        *,
        workflow: str,
        editable: bool,
        slide_count: int,
        part_count: int,
        title: str,
    ) -> str:
        editability = "editable PPTX" if editable else "image-backed non-editable PPTX"
        target = title or "slide deck"
        return (
            f"Created {target} with {slide_count} slides via {workflow} workflow "
            f"as {part_count} PPTX files ({editability}, max {self.config.max_slides_per_deck} slides per file)."
        )


    def _record_session(
        self,
        *,
        task: TaskEnvelope,
        workflow: str,
        editable: bool,
        title: str,
        template_path: Path | None,
        pptx_path: str,
        slide_count: int,
    ) -> None:
        if self.db is None:
            return
        try:
            self.db.execute(
                """INSERT OR REPLACE INTO slide_sessions
                   (session_id, task_id, intent, workflow, editable, title, template_name, pptx_path, slide_count, created_at)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                [
                    task.session_id,
                    task.task_id,
                    task.intent,
                    workflow,
                    1 if editable else 0,
                    title,
                    template_path.name if template_path else None,
                    pptx_path,
                    slide_count,
                    utcnow().isoformat(),
                ],
            )
            self.db.commit()
        except Exception as exc:
            logger.warning("slide_agent.session_record_failed error=%s", exc)

    def _resolve_template_from_input(self, task: TaskEnvelope, *, persist_uploaded: bool) -> Path | None:
        for key in ("template_path", "template", "source_template_path"):
            raw_value = self._text(task.input.get(key))
            path = self._resolve_path(raw_value)
            if path is None or not path.exists():
                path = self._resolve_template_name(raw_value)
            if path is not None and path.exists() and path.suffix.lower() == ".pptx":
                return path

        pptx_artifact = self._find_pptx_artifact(task)
        if pptx_artifact is None:
            return None
        path = self._resolve_artifact_path(pptx_artifact)
        if path is None or not path.exists():
            return None
        if persist_uploaded:
            return self._persist_uploaded_template(path)
        return path

    def _resolve_template_name(self, raw: str) -> Path | None:
        value = self._text(raw)
        if not value:
            return None
        name = Path(value).name
        candidates = [name]
        if not name.lower().endswith(".pptx"):
            candidates.append(f"{name}.pptx")
        for candidate in candidates:
            path = (self.config.templates_dir / candidate).resolve()
            if path.exists() and path.is_file() and path.suffix.lower() == ".pptx":
                return path
        return None

    def _resolve_source_deck(self, task: TaskEnvelope) -> Path | None:
        for key in ("source_pptx_path", "pptx_path", "source_deck_path", "deck_path"):
            path = self._resolve_path(self._text(task.input.get(key)))
            if path is not None and path.exists() and path.suffix.lower() == ".pptx":
                return path
        artifact = self._find_pptx_artifact(task)
        if artifact is None:
            return None
        return self._resolve_artifact_path(artifact)

    def _find_pptx_artifact(self, task: TaskEnvelope) -> dict[str, Any] | None:
        for artifact in self._input_artifacts(task):
            if not isinstance(artifact, dict):
                continue
            mime = self._text(artifact.get("mime")) or self._text(artifact.get("mime_type"))
            filename = self._text(artifact.get("filename")) or Path(self._text(artifact.get("path"))).name
            if mime == _PPTX_MIME or filename.lower().endswith(".pptx"):
                return artifact
        return None

    def _input_artifacts(self, task: TaskEnvelope) -> list[dict[str, Any]]:
        artifacts = [item for item in (task.input_artifacts or []) if isinstance(item, dict)]
        extra = task.input.get("source_artifacts")
        if isinstance(extra, list):
            artifacts.extend(item for item in extra if isinstance(item, dict))
        return artifacts

    def _merge_artifact_descriptors(self, *artifact_lists: Any) -> list[dict[str, Any]]:
        merged: list[dict[str, Any]] = []
        seen: set[tuple[str, str]] = set()
        for artifact_list in artifact_lists:
            if not isinstance(artifact_list, list):
                continue
            for raw in artifact_list:
                normalized = self._normalize_artifact_descriptor(raw)
                if normalized is None:
                    continue
                key = (
                    self._text(normalized.get("artifact_id")),
                    self._text(normalized.get("path")),
                )
                if key in seen:
                    continue
                seen.add(key)
                merged.append(normalized)
        return merged

    def _normalize_artifact_descriptor(self, raw: Any) -> dict[str, Any] | None:
        if not isinstance(raw, dict):
            return None
        path = self._text(raw.get("path"))
        if not path:
            return None
        filename = self._text(raw.get("filename")) or Path(path).name
        artifact_id = self._text(raw.get("artifact_id")) or Path(filename).stem
        mime = (self._text(raw.get("mime")) or self._text(raw.get("mime_type")) or self._mime_for_path(Path(filename))).lower()
        normalized = {
            "artifact_id": artifact_id,
            "path": path,
            "mime": mime,
            "filename": filename,
        }
        for key in ("sha256", "created_by_agent", "kind", "audience", "source_artifact_id", "bundle_id", "doc_id"):
            if key in raw and raw.get(key) is not None:
                normalized[key] = raw.get(key)
        return normalized

    def _is_document_artifact(self, artifact: dict[str, Any]) -> bool:
        mime = self._text(artifact.get("mime")).lower()
        filename = self._text(artifact.get("filename")) or self._text(artifact.get("path"))
        suffix = Path(filename).suffix.lower()
        return mime in _DOCUMENT_MIMES or suffix in _DOCUMENT_EXTENSIONS

    def _is_pptx_artifact(self, artifact: dict[str, Any]) -> bool:
        mime = self._text(artifact.get("mime")).lower()
        filename = self._text(artifact.get("filename")) or self._text(artifact.get("path"))
        return mime == _PPTX_MIME or filename.lower().endswith(".pptx")

    def _is_docs_parser_artifact(self, artifact: dict[str, Any]) -> bool:
        path = self._text(artifact.get("path")).replace("\\", "/").lower()
        created_by_agent = self._text(artifact.get("created_by_agent")).lower()
        filename = (self._text(artifact.get("filename")) or Path(path).name).lower()
        return (
            "/docs_parser/" in path
            or "docs-parser" in created_by_agent
            or filename in {"manifest.json", "chunk_index.json", "document.md", "document.json"}
        )

    def _document_summaries_from_artifacts(self, artifacts: list[dict[str, Any]]) -> list[dict[str, Any]]:
        summaries: list[dict[str, Any]] = []
        for artifact in artifacts:
            resolved = self._resolve_artifact_path(artifact)
            if resolved is None or not resolved.exists() or resolved.name != "manifest.json":
                continue
            payload = self._safe_json_load(resolved)
            if not payload:
                continue
            counts = payload.get("counts") if isinstance(payload.get("counts"), dict) else {}
            outputs = payload.get("outputs") if isinstance(payload.get("outputs"), dict) else {}
            bundle_root = resolved.parent
            document_md_path = self._resolve_bundle_path(bundle_root, outputs.get("document_md") or "document.md")
            chunk_index_path = self._resolve_bundle_path(bundle_root, outputs.get("chunk_index") or "chunk_index.json")
            chunk_index = self._safe_json_load(chunk_index_path) if chunk_index_path is not None else {}
            summary = {
                "doc_id": self._text(payload.get("doc_id")),
                "bundle_id": self._text(payload.get("bundle_id")) or None,
                "artifact_id": self._text(payload.get("source_artifact_id")) or self._text(artifact.get("artifact_id")),
                "filename": self._text(payload.get("filename")) or self._text(artifact.get("filename")) or None,
                "mime": (self._text(payload.get("mime")) or self._text(artifact.get("mime"))).lower(),
                "title": self._text(payload.get("title")) or None,
                "section_count": self._int(counts.get("section_count"), 0),
                "chunk_count": self._int(counts.get("chunk_count"), 0),
                "table_count": self._int(counts.get("table_count"), 0),
                "figure_count": self._int(counts.get("figure_count"), 0),
                "page_count": counts.get("page_count"),
                "slide_count": counts.get("slide_count"),
                "asset_count": self._int(counts.get("asset_count"), 0),
                "preview_excerpt": self._read_compact_markdown(document_md_path, limit=1800),
                "top_sections": self._top_sections_from_chunk_index(chunk_index),
                "paths": {
                    "manifest": resolved.as_posix(),
                    "document_md": document_md_path.as_posix() if document_md_path is not None else None,
                    "chunk_index": chunk_index_path.as_posix() if chunk_index_path is not None else None,
                },
                "artifact_refs": [],
            }
            summaries.append(summary)
        return summaries

    def _enrich_document_summaries(self, documents: Any) -> list[dict[str, Any]]:
        if not isinstance(documents, list):
            return []
        return [
            self._enrich_document_summary(item)
            for item in documents
            if isinstance(item, dict)
        ]

    def _enrich_document_summary(self, document: dict[str, Any]) -> dict[str, Any]:
        summary = json.loads(json.dumps(document, ensure_ascii=False))
        paths = summary.get("paths") if isinstance(summary.get("paths"), dict) else {}
        document_md_path = self._resolve_path(self._text(paths.get("document_md")))
        chunk_index_path = self._resolve_path(self._text(paths.get("chunk_index")))
        if not self._text(summary.get("preview_excerpt")):
            summary["preview_excerpt"] = self._read_compact_markdown(document_md_path, limit=1800)
        if not isinstance(summary.get("top_sections"), list):
            chunk_index = self._safe_json_load(chunk_index_path) if chunk_index_path is not None else {}
            summary["top_sections"] = self._top_sections_from_chunk_index(chunk_index)
        return summary

    def _merge_document_summaries(self, *document_lists: Any) -> list[dict[str, Any]]:
        merged: list[dict[str, Any]] = []
        index_by_key: dict[tuple[str, str, str], int] = {}
        for document_list in document_lists:
            if not isinstance(document_list, list):
                continue
            for raw in document_list:
                if not isinstance(raw, dict):
                    continue
                document = self._enrich_document_summary(raw)
                key = (
                    self._text(document.get("artifact_id")),
                    self._text(document.get("doc_id")),
                    self._text(document.get("filename")),
                )
                if key in index_by_key:
                    merged[index_by_key[key]] = self._merge_document_fields(merged[index_by_key[key]], document)
                    continue
                index_by_key[key] = len(merged)
                merged.append(document)
        return merged

    def _merge_document_fields(self, current: dict[str, Any], incoming: dict[str, Any]) -> dict[str, Any]:
        merged = dict(current)
        for key, value in incoming.items():
            if key in {"artifact_refs", "top_sections"} and isinstance(value, list):
                existing = merged.get(key) if isinstance(merged.get(key), list) else []
                merged[key] = self._dedupe_json_list([*existing, *value])
                continue
            if key == "paths" and isinstance(value, dict):
                paths = dict(merged.get("paths") or {})
                for path_key, path_value in value.items():
                    if not paths.get(path_key) and path_value:
                        paths[path_key] = path_value
                merged["paths"] = paths
                continue
            if self._emptyish(merged.get(key)) and not self._emptyish(value):
                merged[key] = value
        return merged

    def _visual_asset_summaries_from_artifacts(self, artifacts: list[dict[str, Any]]) -> list[dict[str, Any]]:
        assets: list[dict[str, Any]] = []
        seen: set[str] = set()
        for artifact in artifacts:
            resolved = self._resolve_artifact_path(artifact)
            if resolved is None or not resolved.exists() or resolved.name != "chunk_index.json":
                continue
            chunk_index = self._safe_json_load(resolved)
            for collection_name in ("figures", "tables", "assets"):
                entries = chunk_index.get(collection_name) if isinstance(chunk_index.get(collection_name), list) else []
                for entry in entries[:12]:
                    if not isinstance(entry, dict):
                        continue
                    asset_ref = self._text(entry.get("asset_id") or entry.get("asset_ref") or entry.get("id") or entry.get("path"))
                    if not asset_ref or asset_ref in seen:
                        continue
                    seen.add(asset_ref)
                    assets.append(
                        {
                            "asset_ref": asset_ref,
                            "label": self._text(entry.get("label") or entry.get("title") or collection_name[:-1]) or None,
                            "caption": self._text(entry.get("caption") or entry.get("text") or entry.get("summary")) or None,
                            "description": self._text(entry.get("description")) or None,
                            "page_number": entry.get("page_number"),
                            "slide_number": entry.get("slide_number"),
                            "path": entry.get("path"),
                        }
                    )
                    if len(assets) >= 12:
                        return assets
        return assets

    def _resolve_bundle_path(self, bundle_root: Path, raw_path: Any) -> Path | None:
        resolved = self._resolve_path(self._text(raw_path))
        if resolved is not None and resolved.exists() and resolved.is_file():
            return resolved
        relative = self._text(raw_path)
        if not relative:
            return None
        candidate = (bundle_root / relative).resolve()
        return candidate if candidate.exists() and candidate.is_file() else None

    def _read_compact_markdown(self, path: Path | None, *, limit: int) -> str | None:
        if path is None or not path.exists() or not path.is_file():
            return None
        try:
            return self._compact_text(path.read_text(encoding="utf-8"), limit=limit)
        except Exception:
            return None

    def _top_sections_from_chunk_index(self, chunk_index: dict[str, Any], *, limit: int = 8) -> list[dict[str, Any]]:
        sections = chunk_index.get("sections") if isinstance(chunk_index.get("sections"), list) else []
        chunks = chunk_index.get("chunks") if isinstance(chunk_index.get("chunks"), list) else []
        source = sections or chunks
        normalized: list[dict[str, Any]] = []
        for entry in source:
            if not isinstance(entry, dict):
                continue
            section_id = self._text(entry.get("section_id") or entry.get("id") or entry.get("chunk_id"))
            title = self._text(entry.get("title") or entry.get("heading") or section_id)
            preview = self._compact_text(
                entry.get("summary") or entry.get("excerpt") or entry.get("text") or entry.get("search_text") or "",
                limit=320,
            )
            if not section_id and not title and not preview:
                continue
            normalized.append(
                {
                    "section_id": section_id or None,
                    "title": title or None,
                    "page_number": entry.get("page_number"),
                    "slide_number": entry.get("slide_number"),
                    "preview": preview,
                }
            )
            if len(normalized) >= limit:
                break
        return normalized

    def _safe_json_load(self, path: Path | None) -> dict[str, Any]:
        if path is None or not path.exists() or not path.is_file():
            return {}
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
        return payload if isinstance(payload, dict) else {}

    @staticmethod
    def _dedupe_json_list(items: list[Any]) -> list[Any]:
        deduped: list[Any] = []
        seen: set[str] = set()
        for item in items:
            fingerprint = json.dumps(item, ensure_ascii=False, sort_keys=True, default=str)
            if fingerprint in seen:
                continue
            seen.add(fingerprint)
            deduped.append(item)
        return deduped

    @staticmethod
    def _compact_text(value: Any, *, limit: int) -> str:
        text = re.sub(r"\s+", " ", str(value or "")).strip()
        if limit <= 0 or len(text) <= limit:
            return text
        return text[: max(0, limit - 1)].rstrip() + "..."

    @staticmethod
    def _emptyish(value: Any) -> bool:
        return value is None or value == "" or value == [] or value == {}

    def _resolve_artifact_path(self, artifact: dict[str, Any]) -> Path | None:
        return self._resolve_path(self._text(artifact.get("path")))

    def _resolve_path(self, raw: str) -> Path | None:
        value = self._text(raw)
        if not value:
            return None
        candidate = Path(value).expanduser()
        if candidate.is_absolute():
            return candidate.resolve()
        normalized = value.replace("\\", "/")
        raw_path = Path(normalized)
        if len(raw_path.parts) >= 2 and raw_path.parts[0] == "runs" and raw_path.parts[1] == "artifacts":
            return (BACKEND_ROOT / raw_path).resolve()
        local_template = (self.config.templates_dir / candidate).resolve()
        if local_template.exists():
            return local_template
        agent_path = (AGENT_ROOT / candidate).resolve()
        if agent_path.exists():
            return agent_path
        backend_path = (BACKEND_ROOT / candidate).resolve()
        if backend_path.exists():
            return backend_path
        return (self.artifacts_root / candidate).resolve()

    def _persist_uploaded_template(self, source_path: Path) -> Path:
        source_path = source_path.resolve()
        self.config.templates_dir.mkdir(parents=True, exist_ok=True)
        target = self.config.templates_dir / source_path.name
        source_hash = self._sha256(source_path)
        if target.exists():
            if self._sha256(target) == source_hash:
                return target.resolve()
            target = self.config.templates_dir / f"{source_path.stem}-{uuid4().hex[:8]}{source_path.suffix}"
        shutil.copy2(source_path, target)
        return target.resolve()

    def _artifact_manifest(
        self,
        *,
        task_id: str,
        path: Path,
        mime: str,
        kind: str = "output",
        audience: str = "deliverable",
    ) -> ArtifactManifest:
        digest = self._sha256(path)
        resolved = path.resolve()
        try:
            relative = resolved.relative_to(self.artifacts_root)
            logical_path = (Path("runs") / "artifacts" / relative).as_posix()
        except ValueError:
            logical_path = resolved.as_posix()
        return ArtifactManifest(
            artifact_id=f"art_{uuid4().hex[:12]}",
            task_id=task_id,
            mime=mime,
            sha256=digest,
            path=logical_path,
            created_by_agent=self.agent_id,
            kind=kind,
            audience=audience,
        )

    def _task_artifact_dir(self, task_id: str) -> Path:
        return self.artifacts_root / task_id / "slide_agent"

    def _catalog_dir(self, template_path: Path) -> Path:
        return self.config.catalogs_dir / template_path.stem

    def _catalog_path(self, template_path: Path) -> Path:
        return self._catalog_dir(template_path) / "catalog.json"

    def _slide_count_from_result(self, result: dict[str, Any]) -> int:
        build_spec = result.get("build_spec") if isinstance(result.get("build_spec"), dict) else {}
        if isinstance(build_spec.get("slides"), list):
            return len(build_spec["slides"])
        content_plan = result.get("content_plan") if isinstance(result.get("content_plan"), dict) else {}
        if isinstance(content_plan.get("slides"), list):
            return len(content_plan["slides"])
        slide_pngs = result.get("slide_pngs")
        if isinstance(slide_pngs, list):
            return len(slide_pngs)
        return 0

    def _count_pptx_slides(self, pptx_path: Path) -> int | None:
        try:
            from pptx import Presentation

            return len(Presentation(str(pptx_path)).slides)
        except Exception:
            return None

    def _result_title(self, result: dict[str, Any], *, fallback: str) -> str:
        build_spec = result.get("build_spec") if isinstance(result.get("build_spec"), dict) else {}
        title = self._text(build_spec.get("deck_title"))
        if title:
            return title
        content_plan = result.get("content_plan") if isinstance(result.get("content_plan"), dict) else {}
        title = self._text(content_plan.get("deck_title"))
        if title:
            return title
        plan_path = self._maybe_path(result.get("plan_path"))
        if plan_path is not None:
            try:
                plan = json.loads(plan_path.read_text(encoding="utf-8"))
                title = self._text(plan.get("deck_title")) if isinstance(plan, dict) else ""
                if title:
                    return title
            except Exception:
                pass
        return fallback

    @staticmethod
    def _write_json(path: Path, payload: Any) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")

    @staticmethod
    def _sha256(path: Path) -> str:
        h = hashlib.sha256()
        with path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(1024 * 1024), b""):
                h.update(chunk)
        return h.hexdigest()

    @staticmethod
    def _mime_for_path(path: Path) -> str:
        if path.suffix.lower() == ".pptx":
            return _PPTX_MIME
        if path.suffix.lower() == ".pdf":
            return "application/pdf"
        if path.suffix.lower() == ".json":
            return "application/json"
        if path.suffix.lower() == ".png":
            return "image/png"
        return mimetypes.guess_type(str(path))[0] or "application/octet-stream"

    @staticmethod
    def _normalize_workflow(value: Any) -> str:
        normalized = str(value or "").strip().lower()
        return normalized if normalized in {"html", "template", "auto"} else ""

    def _max_slides(self, payload: dict[str, Any]) -> int | None:
        value = self._requested_slide_count(payload)
        if value is None:
            return None
        return max(1, min(value, self.config.max_slides_per_deck))

    def _requested_slide_count(self, payload: dict[str, Any]) -> int | None:
        raw = payload.get("max_slides")
        if raw is None:
            raw = payload.get("slides")
        value = self._int(raw, 0)
        if value <= 0:
            return None
        return max(1, value)

    @staticmethod
    def _bool(value: Any, default: bool) -> bool:
        if value is None:
            return default
        if isinstance(value, bool):
            return value
        return str(value).strip().lower() not in {"0", "false", "no", "off"}

    @staticmethod
    def _int(value: Any, default: int) -> int:
        try:
            return int(value)
        except (TypeError, ValueError):
            return default

    @staticmethod
    def _text(value: Any) -> str:
        return str(value or "").strip()

    @staticmethod
    def _maybe_path(value: Any) -> Path | None:
        text = str(value or "").strip()
        if not text:
            return None
        path = Path(text).expanduser()
        if path.exists() and path.is_file():
            return path.resolve()
        return None

    @staticmethod
    def _failed(
        code: str,
        message: str,
        *,
        retryable: bool = False,
        next_action: str = "revise_input",
    ) -> AgentResult:
        return AgentResult(
            status="failed",
            output={},
            artifacts=[],
            error=AgentError(
                code=code,
                retryable=retryable,
                message=message,
                next_action=next_action,
            ),
        )
